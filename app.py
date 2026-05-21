#!/usr/bin/env python3.8
# -*- coding: utf-8 -*-
"""SmartCS v3.0 - 含后台管理和工单管理"""

import os, json, glob, hashlib, uuid, re, sqlite3, time, csv, io, shutil, threading, secrets, tempfile, subprocess, urllib.parse
from datetime import datetime, timedelta
from functools import wraps
import requests
from flask import Flask, request, jsonify, render_template, send_from_directory, session, g, redirect, Response
from werkzeug.security import generate_password_hash, check_password_hash

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
DB_PATH = os.path.join(BASE_DIR, 'data', 'smartcs.db')
KNOWLEDGE_DIR = os.path.join(BASE_DIR, 'knowledge')
UPLOAD_DIR = os.path.join(BASE_DIR, 'uploads')

DASHSCOPE_API_KEY = os.environ.get('DASHSCOPE_API_KEY', '')
API_BASE_URL = os.environ.get('API_BASE_URL', 'https://api.deepseek.com/v1/chat/completions')
MODEL_NAME = os.environ.get('MODEL_NAME', 'deepseek-chat')
ADMIN_PASSWORD = os.environ.get('ADMIN_PASSWORD', 'admin123')
SECRET_KEY = os.environ.get('SECRET_KEY', 'smart-cs-secret-2026')

app = Flask(__name__)
app.secret_key = SECRET_KEY
def _sync_app_config_from_db():
    """Sync Flask app config from system_config table"""
    try:
        cfg = get_system_config()
        sl = int(cfg.get('session_lifetime', '28800'))
        app.config['PERMANENT_SESSION_LIFETIME'] = sl
        mu = int(cfg.get('max_upload_size_mb', '50'))
        app.config['MAX_CONTENT_LENGTH'] = mu * 1024 * 1024
    except:
        pass

app.config.update(JSON_AS_ASCII=False, TEMPLATES_AUTO_RELOAD=True)
app.config.update(PERMANENT_SESSION_LIFETIME=28800)  # 启动默认值
app.config.update(MAX_CONTENT_LENGTH=50*1024*1024)   # 启动默认值
app.config.update(SESSION_COOKIE_HTTPONLY=True, SESSION_COOKIE_SAMESITE='Lax')

@app.before_request
def check_session_timeout():
    if session.get('agent_id') and request.endpoint and not request.endpoint.startswith('static'):
        now = datetime.now().timestamp()
        last_activity = session.get('last_activity', now)
        cfg = get_system_config()
        idle_timeout = int(cfg.get('session_idle_timeout', '1800'))
        if now - last_activity > idle_timeout:
            session.clear()
            if request.is_json:
                return jsonify({'error':'会话已超时，请重新登录','login_required':True}), 401
            return redirect('/agent/login')
        session['last_activity'] = now

# ====== 品牌/系统配置缓存 ======
_config_cache = {}
_config_cache_time = 0

def get_system_config():
    """Get all system_config as dict with caching (5 min expiry)"""
    global _config_cache, _config_cache_time
    now = time.time()
    if _config_cache and (now - _config_cache_time) < 300:
        return _config_cache
    try:
        rows = get_db().execute("SELECT key, value FROM system_config").fetchall()
        _config_cache = {r['key']: r['value'] for r in rows}
        _config_cache_time = now
    except:
        pass
    return _config_cache

def get_brand_config():
    """Get brand-related config values for templates"""
    cfg = get_system_config()
    return {
        'brand_name': cfg.get('brand_name', 'SmartCS 智能客服'),
        'brand_short': cfg.get('brand_short', 'SmartCS'),
        'brand_primary_color': cfg.get('brand_primary_color', '#1a73e8'),
        'brand_logo_path': cfg.get('brand_logo_path', '/static/icon-192.png'),
        'brand_favicon_path': cfg.get('brand_favicon_path', '/static/favicon.ico'),
    }

def invalidate_config_cache():
    """Invalidate config cache after updates"""
    global _config_cache, _config_cache_time
    _config_cache = {}
    _config_cache_time = 0

@app.context_processor
def inject_brand_config():
    """Inject brand/system config into all templates"""
    _sync_app_config_from_db()
    return get_brand_config()


def get_db():
    if 'db' not in g:
        g.db = sqlite3.connect(DB_PATH, timeout=10)
        g.db.row_factory = sqlite3.Row
        g.db.execute("PRAGMA journal_mode=WAL")
        g.db.execute("PRAGMA synchronous=NORMAL")
        g.db.execute("PRAGMA cache_size=-8000")
        g.db.execute("PRAGMA busy_timeout=5000")
        g.db.execute("PRAGMA foreign_keys=ON")
    return g.db

# ====== 统一角色装饰器工厂 ======
def role_required(*allowed_roles):
    """Unified role decorator factory - 三员管理核心"""
    def decorator(f):
        @wraps(f)
        def decorated(*args, **kwargs):
            if not session.get('agent_id'):
                if request.is_json:
                    return jsonify({'error':'未登录','login_required':True}), 401
                return redirect('/agent/login')
            role = session.get('agent_role', '')
            if role == 'superadmin':  # 兼容角色
                return f(*args, **kwargs)
            if role not in allowed_roles:
                return jsonify({'error':'权限不足'}), 403
            return f(*args, **kwargs)
        return decorated
    return decorator

sysadmin_required = role_required('sysadmin', 'superadmin')
secadmin_required = role_required('secadmin', 'superadmin')
audadmin_required = role_required('audadmin', 'superadmin')
logged_in_required = role_required('agent', 'sysadmin', 'secadmin', 'audadmin', 'superadmin')

agent_required = logged_in_required  # backward compatibility
admin_or_agent_required = logged_in_required  # backward compatibility

admin_required = role_required('sysadmin', 'superadmin')  # backward compat - now same as sysadmin

# ====== 密码策略辅助函数 ======
def _check_password_policy(password):
    """Check password against system_config policy. Returns (ok, error_msg)"""
    if not password:
        return False, '密码不能为空'
    cfg = get_db().execute("SELECT key, value FROM system_config").fetchall()
    config_dict = {r['key']: r['value'] for r in cfg}
    min_len = int(config_dict.get('password_min_length', '8'))
    req_upper = config_dict.get('password_require_upper', 'true') == 'true'
    if len(password) < min_len:
        return False, f'密码长度不能少于{min_len}位'
    if req_upper and not re.search(r'[A-Z]', password):
        return False, '密码必须包含大写字母'
    return True, ''

def _check_delete_agent_protection(aid):
    """Check if agent can be deleted. Returns (ok, error_msg)"""
    role = get_db().execute("SELECT role FROM agents WHERE id=?", (aid,)).fetchone()
    if not role:
        return True, ''
    admin_roles = {'sysadmin', 'secadmin', 'audadmin'}
    if role['role'] not in admin_roles:
        return True, ''
    # Check if this is the last of this role
    count = get_db().execute("SELECT COUNT(*) FROM agents WHERE role=?", (role['role'],)).fetchone()[0]
    if count <= 1:
        role_names = {'sysadmin': '系统管理员', 'secadmin': '安全管理员', 'audadmin': '审计管理员'}
        return False, f'不能删除最后一位{role_names.get(role["role"], role["role"])}'
    return True, ''

def _check_demote_agent_protection(aid):
    """Check if agent can be demoted from admin role"""
    role = get_db().execute("SELECT role FROM agents WHERE id=?", (aid,)).fetchone()
    if not role:
        return True, ''
    admin_roles = {'sysadmin', 'secadmin', 'audadmin'}
    if role['role'] not in admin_roles:
        return True, ''
    count = get_db().execute("SELECT COUNT(*) FROM agents WHERE role=?", (role['role'],)).fetchone()[0]
    if count <= 1:
        role_names = {'sysadmin': '系统管理员', 'secadmin': '安全管理员', 'audadmin': '审计管理员'}
        return False, f'不能降级最后一位{role_names.get(role["role"], role["role"])}'
    return True, ''

def generate_csrf_token():
    if 'csrf_token' not in session:
        session['csrf_token'] = secrets.token_hex(16)
    return session['csrf_token']

def csrf_protect(f):
    @wraps(f)
    def decorated(*args, **kwargs):
        if request.method == 'POST':
            token = request.json.get('csrf_token') if request.is_json else request.form.get('csrf_token')
            if not token or token != session.get('csrf_token'):
                return jsonify({'ok': False, 'error': 'CSRF token invalid or missing'}), 403
        return f(*args, **kwargs)
    return decorated

# ====== 安全头部中间件 ======
# ====== Phase 5: IM Adapter (企微/钉钉消息通道) ======
from abc import ABC, abstractmethod

class IMAdapter(ABC):
    """IM 平台消息适配器基类"""
    
    @abstractmethod
    def send_message(self, user_id: str, content: str, msg_type: str = 'text') -> dict:
        pass
    
    @abstractmethod
    def get_user_info(self, user_id: str) -> dict:
        pass
    
    def send_card(self, user_id: str, title: str, desc: str, url: str = '') -> dict:
        return {'error': 'not_supported'}
    
    def handle_callback(self, body: dict) -> dict:
        return {'error': 'not_supported'}


class WeComAdapter(IMAdapter):
    """企业微信适配器"""
    def __init__(self, config: dict):
        self.corpid = config.get('corpid', '')
        self.agentid = config.get('agentid', '')
        self.corpsecret = config.get('corpsecret', '')
        self._token = None
    
    def _get_access_token(self) -> str:
        try:
            url = f'https://qyapi.weixin.qq.com/cgi-bin/gettoken?corpid={self.corpid}&corpsecret={self.corpsecret}'
            resp = requests.get(url, timeout=10)
            data = resp.json()
            if data.get('errcode') == 0:
                self._token = data['access_token']
                return self._token
            return ''
        except: return ''
    
    def send_message(self, user_id: str, content: str, msg_type: str = 'text') -> dict:
        token = self._get_access_token()
        if not token:
            return {'error': 'get_token_failed'}
        payload = {
            'touser': user_id,
            'msgtype': 'text',
            'agentid': int(self.agentid) if self.agentid else 0,
            'text': {'content': content},
            'safe': 0
        }
        try:
            resp = requests.post(
                f'https://qyapi.weixin.qq.com/cgi-bin/message/send?access_token={token}',
                json=payload, timeout=10
            )
            return resp.json()
        except Exception as e:
            return {'error': str(e)}
    
    def get_user_info(self, user_id: str) -> dict:
        token = self._get_access_token()
        if not token:
            return {'error': 'get_token_failed'}
        try:
            resp = requests.get(
                f'https://qyapi.weixin.qq.com/cgi-bin/user/get?access_token={token}&userid={user_id}',
                timeout=10
            )
            return resp.json()
        except Exception as e:
            return {'error': str(e)}
    
    def handle_callback(self, body: dict) -> dict:
        """处理企微回调消息"""
        msg_type = body.get('MsgType', 'text')
        if msg_type == 'text':
            return {
                'platform_user_id': body.get('FromUserName', ''),
                'content': body.get('Content', ''),
                'msg_type': 'text'
            }
        elif msg_type == 'image':
            return {
                'platform_user_id': body.get('FromUserName', ''),
                'media_id': body.get('MediaId', ''),
                'msg_type': 'image'
            }
        return {'error': 'unsupported_msg_type'}


class DingTalkAdapter(IMAdapter):
    """钉钉适配器"""
    def __init__(self, config: dict):
        self.appkey = config.get('appkey', '')
        self.appsecret = config.get('appsecret', '')
        self._token = None
    
    def _get_access_token(self) -> str:
        try:
            resp = requests.post(
                'https://oapi.dingtalk.com/gettoken',
                json={'appkey': self.appkey, 'appsecret': self.appsecret},
                timeout=10
            )
            data = resp.json()
            if data.get('errcode') == 0:
                self._token = data['access_token']
                return self._token
            return ''
        except: return ''
    
    def send_message(self, user_id: str, content: str, msg_type: str = 'text') -> dict:
        token = self._get_access_token()
        if not token:
            return {'error': 'get_token_failed'}
        payload = {
            'agent_id': self.appkey,
            'userid_list': user_id,
            'msg': {
                'msgtype': 'text',
                'text': {'content': content}
            }
        }
        try:
            resp = requests.post(
                f'https://oapi.dingtalk.com/topapi/message/corpconversation/asyncsend_v2?access_token={token}',
                json=payload, timeout=10
            )
            return resp.json()
        except Exception as e:
            return {'error': str(e)}
    
    def get_user_info(self, user_id: str) -> dict:
        token = self._get_access_token()
        if not token:
            return {'error': 'get_token_failed'}
        try:
            resp = requests.post(
                f'https://oapi.dingtalk.com/topapi/v2/user/get?access_token={token}',
                json={'userid': user_id},
                timeout=10
            )
            return resp.json()
        except Exception as e:
            return {'error': str(e)}


_ADAPTER_CACHE = {}

def get_im_adapter(adapter_id: str) -> IMAdapter:
    """工厂方法：根据 adapter_id 获取对应的适配器实例"""
    if adapter_id in _ADAPTER_CACHE:
        return _ADAPTER_CACHE[adapter_id]
    row = get_db().execute("SELECT platform, config FROM im_adapters WHERE id=?", (adapter_id,)).fetchone()
    if not row:
        return None
    config = json.loads(row['config'] or '{}')
    if row['platform'] == 'wecom':
        inst = WeComAdapter(config)
    elif row['platform'] == 'dingtalk':
        inst = DingTalkAdapter(config)
    else:
        inst = None
    if inst:
        _ADAPTER_CACHE[adapter_id] = inst
    return inst


def _clear_adapter_cache():
    _ADAPTER_CACHE.clear()


def send_im_notification(adapter_id: str, user_id: str, content: str, msg_type: str = 'text') -> dict:
    """统一接口：通过指定适配器发送消息到 IM 用户"""
    adapter = get_im_adapter(adapter_id)
    if not adapter:
        return {'error': 'adapter_not_found'}
    result = adapter.send_message(user_id, content, msg_type)
    db = get_db()
    db.execute(
        "INSERT INTO im_messages(id,adapter_id,direction,platform_user_id,msg_type,content,status,platform_msg_id) "
        "VALUES(?,?,?,?,?,?,?,?)",
        (gen_id('im-'), adapter_id, 'outbound', user_id, msg_type, content,
         'sent' if result.get('errcode') == 0 else 'failed',
         json.dumps(result, ensure_ascii=False))
    )
    db.commit()
    return result


def send_im_notification_by_platform(platform: str, user_id: str, content: str) -> dict:
    """对外统一接口：找到平台上第一个启用的适配器发消息"""
    row = get_db().execute("SELECT id FROM im_adapters WHERE platform=? AND enabled=1 LIMIT 1", (platform,)).fetchone()
    if not row:
        return {'error': f'no_enabled_adapter_for_{platform}'}
    return send_im_notification(row['id'], user_id, content)

# ====== IM 适配器管理 API ======

@app.route('/api/admin/im-adapters', methods=['GET'])
@admin_required
def admin_list_im_adapters():
    adapters = get_db().execute("SELECT * FROM im_adapters ORDER BY platform, created_at").fetchall()
    return jsonify([dict(a) for a in adapters])


@app.route('/api/admin/im-adapters', methods=['POST'])
@admin_required
def admin_create_im_adapter():
    data = request.get_json()
    aid = 'im-' + uuid.uuid4().hex[:12]
    platform = data.get('platform', '')
    if platform not in ('wecom', 'dingtalk', 'feishu'):
        return jsonify({'error': '不支持的平台类型'}), 400
    get_db().execute(
        "INSERT INTO im_adapters(id,name,platform,enabled,config,callback_url,token,encoding_aes_key) "
        "VALUES(?,?,?,?,?,?,?,?)",
        (aid, data.get('name', platform), platform,
         data.get('enabled', 0), json.dumps(data.get('config', {}), ensure_ascii=False),
         data.get('callback_url', ''), data.get('token', ''),
         data.get('encoding_aes_key', ''))
    )
    get_db().commit()
    log_audit('', 'admin.im_adapter.create', session.get('agent_id',''), session.get('agent_name',''),
              {'adapter_id': aid, 'platform': platform, 'name': data.get('name', platform)})
    return jsonify({'ok': True, 'id': aid})


@app.route('/api/admin/im-adapters/<aid>', methods=['PUT'])
@admin_required
def admin_update_im_adapter(aid):
    data = request.get_json()
    existing = get_db().execute("SELECT id FROM im_adapters WHERE id=?", (aid,)).fetchone()
    if not existing:
        return jsonify({'error': '适配器不存在'}), 404
    fields = []
    vals = []
    for k in ('name', 'platform', 'enabled', 'callback_url', 'token', 'encoding_aes_key'):
        if k in data:
            fields.append(k + '=?')
            vals.append(data[k])
    if 'config' in data:
        fields.append('config=?')
        vals.append(json.dumps(data['config'], ensure_ascii=False))
    if fields:
        vals.append(aid)
        get_db().execute(
            'UPDATE im_adapters SET ' + ','.join(fields) + ', updated_at=datetime("now","localtime") WHERE id=?',
            vals
        )
        get_db().commit()
        _clear_adapter_cache()
    log_audit('', 'admin.im_adapter.update', session.get('agent_id',''), session.get('agent_name',''),
              {'adapter_id': aid, 'fields_updated': list(data.keys()) if 'data' in dir() else []})
    return jsonify({'ok': True})


@app.route('/api/admin/im-adapters/<aid>', methods=['DELETE'])
@admin_required
def admin_delete_im_adapter(aid):
    db = get_db()
    db.execute("DELETE FROM im_messages WHERE adapter_id=?", (aid,))
    db.execute("DELETE FROM im_user_mappings WHERE adapter_id=?", (aid,))
    db.execute("DELETE FROM im_adapters WHERE id=?", (aid,))
    db.commit()
    _clear_adapter_cache()
    log_audit('', 'admin.im_adapter.delete', session.get('agent_id',''), session.get('agent_name',''),
              {'adapter_id': aid})
    return jsonify({'ok': True})


@app.route('/api/admin/im-adapters/<aid>/test', methods=['POST'])
@admin_required
def admin_test_im_adapter(aid):
    row = get_db().execute("SELECT platform, config FROM im_adapters WHERE id=?", (aid,)).fetchone()
    if not row:
        return jsonify({'error': '适配器不存在'}), 404
    config = json.loads(row['config'] or '{}')
    if row['platform'] == 'wecom':
        imp = WeComAdapter(config)
        token = imp._get_access_token()
        if token:
            return jsonify({'ok': True, 'message': '连接成功, access_token 已获取'})
        return jsonify({'ok': False, 'message': '获取 access_token 失败，请检查 corpid/corpsecret 配置'})
    elif row['platform'] == 'dingtalk':
        imp = DingTalkAdapter(config)
        token = imp._get_access_token()
        if token:
            return jsonify({'ok': True, 'message': '连接成功, access_token 已获取'})
        return jsonify({'ok': False, 'message': '获取 access_token 失败，请检查 appkey/appsecret 配置'})
    return jsonify({'ok': False, 'message': '不支持的平台类型'})


@app.route('/api/im/callback/<adapter_id>', methods=['POST'])
def im_callback(adapter_id):
    """IM 平台回调入口"""
    row = get_db().execute("SELECT platform, config FROM im_adapters WHERE id=?", (adapter_id,)).fetchone()
    if not row:
        return jsonify({'error': 'adapter_not_found'}), 404
    config = json.loads(row['config'] or '{}')
    body = request.get_json(silent=True) or {}
    if row['platform'] == 'wecom':
        # Handle URL verification from wecom
        echostr = request.args.get('echostr', '')
        if echostr:
            return echostr, 200, {'Content-Type': 'text/plain'}
        adapter = WeComAdapter(config)
        parsed = adapter.handle_callback(body)
        if 'platform_user_id' in parsed:
            # Log inbound message
            db = get_db()
            db.execute(
                "INSERT INTO im_messages(id,adapter_id,direction,platform_user_id,msg_type,content,media_id) "
                "VALUES(?,?,?,?,?,?,?)",
                (gen_id('im-'), adapter_id, 'inbound', parsed['platform_user_id'],
                 parsed.get('msg_type', 'text'), parsed.get('content', ''),
                 parsed.get('media_id', ''))
            )
            db.commit()
            # Check if mapped to a customer, route to existing conversation
            mapping = db.execute(
                "SELECT smartcs_customer_id FROM im_user_mappings "
                "WHERE adapter_id=? AND platform_user_id=?",
                (adapter_id, parsed['platform_user_id'])
            ).fetchone()
            if mapping:
                return jsonify({'ok': True, 'customer_id': mapping['smartcs_customer_id']})
            return jsonify({'ok': True, 'message': 'inbound_logged'})
        return jsonify({'ok': False, 'error': 'parse_failed'})
    elif row['platform'] == 'dingtalk':
        adapter = DingTalkAdapter(config)
        parsed = adapter.handle_callback(body)
        return jsonify({'ok': True, 'data': parsed})
    return jsonify({'ok': False, 'error': 'unsupported_platform'}), 400


@app.route('/api/im/notify', methods=['POST'])
def im_notify_api():
    """通过配置的 IM 发送通知（外部调用接口）"""
    data = request.get_json()
    adapter_id = data.get('adapter_id', '')
    platform = data.get('platform', 'wecom')
    user_id = data.get('user_id', '')
    content = data.get('content', '')
    if not user_id or not content:
        return jsonify({'error': '缺少 user_id 或 content'}), 400
    if adapter_id:
        result = send_im_notification(adapter_id, user_id, content)
    else:
        result = send_im_notification_by_platform(platform, user_id, content)
    return jsonify(result)


@app.route('/api/admin/im-adapters/<aid>/mappings', methods=['GET'])
@admin_required
def admin_list_im_mappings(aid):
    mappings = get_db().execute("""
        SELECT m.*, c.name as customer_name, ag.name as agent_name
        FROM im_user_mappings m
        LEFT JOIN customers c ON m.smartcs_customer_id = c.id
        LEFT JOIN agents ag ON m.smartcs_agent_id = ag.id
        WHERE m.adapter_id=? ORDER BY m.created_at DESC""", (aid,)).fetchall()
    return jsonify([dict(m) for m in mappings])


@app.route('/api/admin/im-adapters/<aid>/mappings', methods=['POST'])
@admin_required
def admin_create_im_mapping(aid):
    data = request.get_json()
    pid = data.get('platform_user_id', '')
    if not pid:
        return jsonify({'error': '缺少 platform_user_id'}), 400
    mid = 'im-' + uuid.uuid4().hex[:12]
    try:
        get_db().execute(
            "INSERT INTO im_user_mappings(id,adapter_id,platform_user_id,smartcs_customer_id,smartcs_agent_id,platform_username,avatar_url) "
            "VALUES(?,?,?,?,?,?,?)",
            (mid, aid, pid,
             data.get('smartcs_customer_id', ''),
             data.get('smartcs_agent_id', ''),
             data.get('platform_username', ''),
             data.get('avatar_url', ''))
        )
        get_db().commit()
        return jsonify({'ok': True, 'id': mid})
    except Exception as e:
        return jsonify({'error': str(e)}), 400


@app.route('/api/admin/im-mappings/<mid>', methods=['DELETE'])
@admin_required
def admin_delete_im_mapping(mid):
    get_db().execute("DELETE FROM im_user_mappings WHERE id=?", (mid,))
    get_db().commit()
    return jsonify({'ok': True})


@app.route('/api/admin/im-messages', methods=['GET'])
@admin_required
def admin_list_im_messages():
    aid = request.args.get('adapter_id', '')
    page = int(request.args.get('page', 1))
    limit = int(get_system_config().get('pagination_per_page', '50'))
    offset = (page - 1) * limit
    where = ['1=1']
    params = []
    if aid:
        where.append('m.adapter_id=?')
        params.append(aid)
    wsql = ' AND '.join(where)
    total = get_db().execute(
        f"SELECT COUNT(*) FROM im_messages m WHERE {wsql}", params
    ).fetchone()[0]
    msgs = get_db().execute(
        f"SELECT m.*, a.name as adapter_name FROM im_messages m LEFT JOIN im_adapters a ON m.adapter_id=a.id "
        f"WHERE {wsql} ORDER BY m.created_at DESC LIMIT ? OFFSET ?",
        params + [limit, offset]
    ).fetchall()
    return jsonify({'messages': [dict(m) for m in msgs], 'total': total, 'page': page, 'pages': max(1, (total + limit - 1) // limit)})


# ====== IM 事件驱动集成 ======

def _send_im_event_notification(event_type: str, ticket_data: dict):
    """根据事件类型发送 IM 通知到对应人员"""
    db = get_db()
    ticket_id = ticket_data.get('ticket_id', '')
    if not ticket_id:
        return
    t = db.execute("SELECT agent_id, customer_id FROM service_tickets WHERE id=?", (ticket_id,)).fetchone()
    if not t:
        return
    
    # Build notification content based on event_type
    messages = {
        'ticket.processing': '🔔 新工单进入处理队列，等待IT工程师处理。',
        'ticket.assigned': '📋 您有一个新分配的工单需要处理。',
        'ticket.resolved': '🔧 工单已解决，等待客户确认评价。',
        'ticket.rated': '⭐ 客户已完成评价，等待IT工程师最终关闭。',
        'ticket.closed': '✅ 工单已关闭。',
    }
    content = messages.get(event_type, f'工单更新: {event_type}')
    
    # Notify assigned agent via configured IM
    if t['agent_id']:
        mappings = db.execute(
            "SELECT m.*, a.id as adapter_id, a.platform FROM im_user_mappings m "
            "JOIN im_adapters a ON m.adapter_id = a.id AND a.enabled=1 "
            "WHERE m.smartcs_agent_id=? LIMIT 1",
            (t['agent_id'],)
        ).fetchall()
        for mapping in mappings:
            send_im_notification(mapping['adapter_id'], mapping['platform_user_id'], content)


# ====== End Phase 5 ======

# ====== Phase 6: External Service Adapters (Jira/禅道/GitHub缺陷系统集成) ======

class ExternalServiceAdapter:
    """外部服务适配器基类"""
    name: str = "base"
    description: str = ""

    def __init__(self, config: dict):
        self.config = config

    def validate_config(self) -> tuple:
        return True, ""

    def get_actions(self) -> list:
        return []

    def execute_action(self, action: str, params: dict) -> dict:
        return {"error": "not_implemented"}


class DefectSystemAdapter(ExternalServiceAdapter):
    """缺陷系统适配器基类"""
    name = "缺陷系统"

    def create_ticket(self, title: str, description: str, priority: str = 'normal', reporter: str = '') -> dict:
        return {"error": "not_implemented"}

    def get_status(self, ticket_id: str) -> dict:
        return {"error": "not_implemented"}

    def get_actions(self):
        return [
            {"name": "create_defect", "params": ["title", "description", "priority"], "description": "创建缺陷工单"},
            {"name": "query_defect", "params": ["ticket_id"], "description": "查询缺陷状态"},
        ]


class JiraAdapter(DefectSystemAdapter):
    """Jira 适配器"""
    name = "Jira / JSM"
    description = "Atlassian Jira 和 JSM 缺陷跟踪"

    def __init__(self, config):
        super().__init__(config)
        self.base_url = config.get('url', '').rstrip('/')
        self.username = config.get('username', '')
        self.api_token = config.get('api_token', '')
        self.project_key = config.get('project_key', '')

    def validate_config(self) -> tuple:
        if not self.base_url or not self.api_token:
            return False, "请填写 URL 和 API Token"
        if not self.project_key:
            return False, "请填写项目 Key"
        return True, ""

    def create_ticket(self, title, description, priority='normal', reporter=''):
        # 模拟实现 - 实际可调用 Jira REST API: POST /rest/api/2/issue
        issue_id = f"{self.project_key}-{uuid.uuid4().hex[:6].upper()}"
        return {
            "success": True,
            "external_id": issue_id,
            "external_url": f"{self.base_url}/browse/{issue_id}",
            "error": ""
        }

    def get_status(self, ticket_id):
        return {"success": True, "status": "Open", "updated": "", "error": ""}


class ZentaoAdapter(DefectSystemAdapter):
    """禅道适配器"""
    name = "禅道"
    description = "禅道项目管理缺陷跟踪"

    def __init__(self, config):
        super().__init__(config)
        self.base_url = config.get('url', '').rstrip('/')
        self.api_key = config.get('api_key', '')
        self.project_id = config.get('project_id', '')

    def validate_config(self) -> tuple:
        if not self.base_url or not self.api_key:
            return False, "请填写 URL 和 API Key"
        if not self.project_id:
            return False, "请填写项目 ID"
        return True, ""

    def create_ticket(self, title, description, priority='normal', reporter=''):
        bug_id = f"BUG-{uuid.uuid4().hex[:8].upper()}"
        return {
            "success": True,
            "external_id": bug_id,
            "external_url": f"{self.base_url}/bug-view-{bug_id}.html",
            "error": ""
        }

    def get_status(self, ticket_id):
        return {"success": True, "status": "Active", "updated": "", "error": ""}


class GitHubIssuesAdapter(DefectSystemAdapter):
    """GitHub Issues 适配器"""
    name = "GitHub Issues"
    description = "GitHub Issues 缺陷跟踪"

    def __init__(self, config):
        super().__init__(config)
        self.repo = config.get('repo', '')
        self.token = config.get('token', '')

    def validate_config(self) -> tuple:
        if not self.repo:
            return False, "请填写仓库路径（如 owner/repo）"
        if not self.token:
            return False, "请填写 GitHub Token"
        return True, ""

    def create_ticket(self, title, description, priority='normal', reporter=''):
        issue_num = f"#{uuid.uuid4().hex[:4].upper()}"
        return {
            "success": True,
            "external_id": issue_num,
            "external_url": f"https://github.com/{self.repo}/issues/{issue_num[1:]}",
            "error": ""
        }

    def get_status(self, ticket_id):
        return {"success": True, "status": "Open", "updated": "", "error": ""}


# 适配器工厂
_ADAPTER_REGISTRY = {}

def register_adapter(platform: str, adapter_class):
    _ADAPTER_REGISTRY[platform] = adapter_class

def get_external_adapter(adapter_id: str):
    """根据 adapter_id 获取适配器实例"""
    row = get_db().execute("SELECT * FROM external_adapters WHERE id=?", (adapter_id,)).fetchone()
    if not row:
        return None
    cfg = row['config']
    config = json.loads(cfg) if isinstance(cfg, str) else cfg
    cls = _ADAPTER_REGISTRY.get(row['platform'])
    if cls:
        return cls(config)
    return None

# 注册内置适配器
register_adapter('jira', JiraAdapter)
register_adapter('zentao', ZentaoAdapter)
register_adapter('github', GitHubIssuesAdapter)
register_adapter('gitlab', GitHubIssuesAdapter)

# ====== End Phase 6 ======

_old_full_dispatch = app.full_dispatch_request
@app.after_request
def _secure_headers(resp):
    resp.headers['X-Content-Type-Options'] = 'nosniff'
    resp.headers['X-Frame-Options'] = 'DENY'
    resp.headers['X-XSS-Protection'] = '1; mode=block'
    resp.headers['Referrer-Policy'] = 'same-origin'
    return resp

@app.teardown_appcontext
def close_db(exception):
    db = g.pop('db', None)
    if db: db.close()

def init_db():
    with app.app_context():
        db = get_db()
        db.executescript('''
        CREATE TABLE IF NOT EXISTS customers (
            id TEXT PRIMARY KEY, name TEXT DEFAULT '游客', contact TEXT DEFAULT '',
            source TEXT DEFAULT 'web', created_at TEXT DEFAULT (datetime('now','localtime')));
        CREATE TABLE IF NOT EXISTS conversations (
            id TEXT PRIMARY KEY, customer_id TEXT REFERENCES customers(id),
            status TEXT DEFAULT 'active', source TEXT DEFAULT 'web',
            created_at TEXT DEFAULT (datetime('now','localtime')),
            updated_at TEXT DEFAULT (datetime('now','localtime')));
        CREATE TABLE IF NOT EXISTS messages (
            id TEXT PRIMARY KEY, conversation_id TEXT REFERENCES conversations(id),
            role TEXT NOT NULL CHECK(role IN ('user','bot','agent','system')),
            content TEXT NOT NULL, msg_type TEXT DEFAULT 'text',
            image_url TEXT DEFAULT '',
            created_at TEXT DEFAULT (datetime('now','localtime')));
        CREATE TABLE IF NOT EXISTS knowledge_files (
            id TEXT PRIMARY KEY, filename TEXT UNIQUE, word_count INTEGER DEFAULT 0,
            uploaded_by TEXT DEFAULT '', created_at TEXT DEFAULT (datetime('now','localtime')),
            updated_at TEXT DEFAULT (datetime('now','localtime')),
            title TEXT DEFAULT '', status TEXT DEFAULT 'approved', submitted_by TEXT DEFAULT '',
            review_notes TEXT DEFAULT '',
            created_by TEXT DEFAULT '', updated_by TEXT DEFAULT '',
            tags TEXT DEFAULT '{}');
        CREATE TABLE IF NOT EXISTS knowledge_history (
            id TEXT PRIMARY KEY, knowledge_id TEXT REFERENCES knowledge_files(id),
            updated_by TEXT DEFAULT '', updated_by_name TEXT DEFAULT '',
            updated_at TEXT DEFAULT (datetime('now','localtime')),
            change_summary TEXT DEFAULT '', old_content TEXT DEFAULT '', new_content TEXT DEFAULT '');
        CREATE TABLE IF NOT EXISTS escalations (
            id TEXT PRIMARY KEY, conversation_id TEXT REFERENCES conversations(id),
            agent_id TEXT REFERENCES agents(id), reason TEXT DEFAULT '',
            status TEXT DEFAULT 'pending', created_at TEXT DEFAULT (datetime('now','localtime')),
            assigned_at TEXT, resolved_at TEXT);
        CREATE TABLE IF NOT EXISTS agents (
            id TEXT PRIMARY KEY, name TEXT NOT NULL, email TEXT UNIQUE,
            password_hash TEXT NOT NULL, status TEXT DEFAULT 'offline', role TEXT DEFAULT 'agent',
            created_at TEXT DEFAULT (datetime('now','localtime')));
        CREATE TABLE IF NOT EXISTS customer_profiles (
            id TEXT PRIMARY KEY, customer_id TEXT UNIQUE REFERENCES customers(id),
            name TEXT DEFAULT '', phone TEXT DEFAULT '', company TEXT DEFAULT '',
            notes TEXT DEFAULT '', created_at TEXT DEFAULT (datetime('now','localtime')),
            updated_at TEXT DEFAULT (datetime('now','localtime')));
        CREATE TABLE IF NOT EXISTS agent_profiles (
            id TEXT PRIMARY KEY, agent_id TEXT UNIQUE REFERENCES agents(id),
            display_name TEXT DEFAULT '', department TEXT DEFAULT '',
            title TEXT DEFAULT '', phone TEXT DEFAULT '',
            created_at TEXT DEFAULT (datetime('now','localtime')));
        CREATE TABLE IF NOT EXISTS customer_agent_bindings (
            id TEXT PRIMARY KEY, customer_id TEXT NOT NULL REFERENCES customers(id),
            agent_id TEXT NOT NULL REFERENCES agents(id),
            binding_type TEXT DEFAULT 'temporary', status TEXT DEFAULT 'active',
            created_at TEXT DEFAULT (datetime('now','localtime')),
            UNIQUE(customer_id, agent_id));
        CREATE TABLE IF NOT EXISTS service_tickets (
            id TEXT PRIMARY KEY, escalation_id TEXT REFERENCES escalations(id),
            conversation_id TEXT NOT NULL REFERENCES conversations(id),
            customer_id TEXT NOT NULL REFERENCES customers(id),
            agent_id TEXT REFERENCES agents(id), status TEXT DEFAULT 'open',
            issue_description TEXT DEFAULT '', priority TEXT DEFAULT 'normal',
            image_url TEXT DEFAULT '', ticket_number TEXT DEFAULT '',
            close_reason TEXT DEFAULT '', transferred_from TEXT DEFAULT '',
            first_read_at TEXT DEFAULT '', assigned_at TEXT, resolved_at TEXT, 
            rated_at TEXT, confirmed_at TEXT, closed_at TEXT,
            resolution_notes TEXT DEFAULT '', customer_rating INTEGER DEFAULT 0,
            customer_feedback TEXT DEFAULT '',
            created_at TEXT DEFAULT (datetime('now','localtime')),
            updated_at TEXT DEFAULT (datetime('now','localtime')));

        CREATE INDEX IF NOT EXISTS idx_conv_customer ON conversations(customer_id);
        CREATE INDEX IF NOT EXISTS idx_msg_conv ON messages(conversation_id);
        CREATE INDEX IF NOT EXISTS idx_esc_status ON escalations(status);
        CREATE INDEX IF NOT EXISTS idx_ticket_conv ON service_tickets(conversation_id);
        CREATE INDEX IF NOT EXISTS idx_ticket_cust ON service_tickets(customer_id);
        CREATE INDEX IF NOT EXISTS idx_ticket_status ON service_tickets(status);
        CREATE TABLE IF NOT EXISTS audit_log (
            id TEXT PRIMARY KEY, ticket_id TEXT REFERENCES service_tickets(id),
            actor_id TEXT, actor_name TEXT DEFAULT '',
            action TEXT NOT NULL, detail TEXT DEFAULT '{}',
            ip_address TEXT DEFAULT '',
            created_at TEXT DEFAULT (datetime('now','localtime')));
        CREATE TABLE IF NOT EXISTS login_log (
            id TEXT PRIMARY KEY,
            user_id TEXT,
            user_name TEXT DEFAULT '',
            user_type TEXT CHECK(user_type IN ('customer','agent')),
            ip_address TEXT DEFAULT '',
            user_agent TEXT DEFAULT '',
            success INTEGER DEFAULT 0,
            fail_reason TEXT DEFAULT '',
            login_at TEXT DEFAULT (datetime('now','localtime'))
        );
        CREATE INDEX IF NOT EXISTS idx_audit_ticket ON audit_log(ticket_id);
        CREATE INDEX IF NOT EXISTS idx_audit_time ON audit_log(created_at);

        CREATE TABLE IF NOT EXISTS tickets_archive (
            id TEXT PRIMARY KEY, escalation_id TEXT DEFAULT '',
            conversation_id TEXT NOT NULL, customer_id TEXT NOT NULL,
            agent_id TEXT, status TEXT DEFAULT 'archived',
            issue_description TEXT DEFAULT '', priority TEXT DEFAULT 'normal',
            image_url TEXT DEFAULT '',
            assigned_at TEXT, resolved_at TEXT, confirmed_at TEXT, closed_at TEXT,
            resolution_notes TEXT DEFAULT '', customer_rating INTEGER DEFAULT 0,
            customer_feedback TEXT DEFAULT '',
            created_at TEXT DEFAULT (datetime('now','localtime')),
            updated_at TEXT DEFAULT (datetime('now','localtime')),
            ticket_number TEXT DEFAULT '', close_reason TEXT DEFAULT '',
            admin_remarks TEXT DEFAULT '', transferred_from TEXT DEFAULT '',
            level INTEGER DEFAULT 1, reopened_at TEXT DEFAULT '',
            reopened_count INTEGER DEFAULT 0,
            archived_at TEXT DEFAULT (datetime('now','localtime')));
        CREATE INDEX IF NOT EXISTS idx_archive_cust ON tickets_archive(customer_id);
        CREATE INDEX IF NOT EXISTS idx_archive_agent ON tickets_archive(agent_id);

        CREATE TABLE IF NOT EXISTS system_upgrades (id TEXT PRIMARY KEY, version TEXT NOT NULL, release_date TEXT NOT NULL DEFAULT '', content TEXT DEFAULT '', created_at TEXT DEFAULT (datetime('now','localtime')));

        CREATE TABLE IF NOT EXISTS ticket_seq (date_str TEXT PRIMARY KEY, seq INTEGER DEFAULT 0);
        CREATE TABLE IF NOT EXISTS close_reasons (id TEXT PRIMARY KEY, name TEXT UNIQUE, sort_order INTEGER DEFAULT 0);
        CREATE TABLE IF NOT EXISTS systems (id TEXT PRIMARY KEY, name TEXT NOT NULL UNIQUE, description TEXT DEFAULT '', created_at TEXT DEFAULT (datetime('now','localtime')));
        CREATE TABLE IF NOT EXISTS agent_systems (agent_id TEXT REFERENCES users(id), system_id TEXT REFERENCES systems(id), PRIMARY KEY(agent_id, system_id));
        CREATE TABLE IF NOT EXISTS system_config (key TEXT PRIMARY KEY, value TEXT NOT NULL, updated_at TEXT DEFAULT (datetime('now','localtime')));
        CREATE TABLE IF NOT EXISTS webhooks (
            id TEXT PRIMARY KEY, name TEXT NOT NULL, url TEXT NOT NULL,
            secret TEXT DEFAULT '', events TEXT NOT NULL DEFAULT '[]',
            retry INTEGER DEFAULT 3, timeout INTEGER DEFAULT 10,
            enabled INTEGER DEFAULT 1,
            created_at TEXT DEFAULT (datetime('now','localtime')),
            updated_at TEXT DEFAULT (datetime('now','localtime'))
        );
        CREATE TABLE IF NOT EXISTS webhook_logs (
            id TEXT PRIMARY KEY, webhook_id TEXT REFERENCES webhooks(id),
            event_type TEXT NOT NULL, payload TEXT NOT NULL,
            status TEXT NOT NULL DEFAULT 'pending',
            status_code INTEGER DEFAULT 0, response_body TEXT DEFAULT '',
            attempt INTEGER DEFAULT 1, error TEXT DEFAULT '',
            created_at TEXT DEFAULT (datetime('now','localtime'))
        );
        CREATE INDEX IF NOT EXISTS idx_whlog_webhook ON webhook_logs(webhook_id);
        CREATE INDEX IF NOT EXISTS idx_whlog_status ON webhook_logs(status);

        -- IM 适配器配置
        CREATE TABLE IF NOT EXISTS im_adapters (
            id TEXT PRIMARY KEY,
            name TEXT NOT NULL,
            platform TEXT NOT NULL,
            enabled INTEGER DEFAULT 0,
            config TEXT NOT NULL DEFAULT '{}',
            callback_url TEXT DEFAULT '',
            token TEXT DEFAULT '',
            encoding_aes_key TEXT DEFAULT '',
            created_at TEXT DEFAULT (datetime('now','localtime')),
            updated_at TEXT DEFAULT (datetime('now','localtime'))
        );
        CREATE TABLE IF NOT EXISTS im_user_mappings (
            id TEXT PRIMARY KEY,
            adapter_id TEXT NOT NULL REFERENCES im_adapters(id),
            platform_user_id TEXT NOT NULL,
            smartcs_customer_id TEXT REFERENCES customers(id),
            smartcs_agent_id TEXT REFERENCES agents(id),
            platform_username TEXT DEFAULT '',
            avatar_url TEXT DEFAULT '',
            created_at TEXT DEFAULT (datetime('now','localtime')),
            UNIQUE(adapter_id, platform_user_id)
        );
        CREATE TABLE IF NOT EXISTS im_messages (
            id TEXT PRIMARY KEY,
            adapter_id TEXT NOT NULL REFERENCES im_adapters(id),
            direction TEXT NOT NULL,
            platform_user_id TEXT NOT NULL,
            msg_type TEXT DEFAULT 'text',
            content TEXT DEFAULT '',
            media_id TEXT DEFAULT '',
            status TEXT DEFAULT 'sent',
            platform_msg_id TEXT DEFAULT '',
            created_at TEXT DEFAULT (datetime('now','localtime'))
        );
        CREATE INDEX IF NOT EXISTS idx_im_msgs_adapter ON im_messages(adapter_id, created_at);
        CREATE INDEX IF NOT EXISTS idx_im_maps_adapter ON im_user_mappings(adapter_id);

        -- 外部服务适配器配置
        CREATE TABLE IF NOT EXISTS external_adapters (
            id TEXT PRIMARY KEY,
            name TEXT NOT NULL,
            adapter_type TEXT NOT NULL,
            platform TEXT NOT NULL,
            enabled INTEGER DEFAULT 1,
            config TEXT NOT NULL DEFAULT '{}',
            created_at TEXT DEFAULT (datetime('now','localtime')),
            updated_at TEXT DEFAULT (datetime('now','localtime'))
        );
        -- 工单-外部系统关联
        CREATE TABLE IF NOT EXISTS ticket_external_links (
            id TEXT PRIMARY KEY,
            ticket_id TEXT NOT NULL REFERENCES service_tickets(id),
            external_system TEXT NOT NULL,
            external_id TEXT NOT NULL,
            external_url TEXT DEFAULT '',
            link_type TEXT DEFAULT 'defect',
            sync_status TEXT DEFAULT 'active',
            created_at TEXT DEFAULT (datetime('now','localtime')),
            updated_at TEXT DEFAULT (datetime('now','localtime'))
        );
        CREATE INDEX IF NOT EXISTS idx_ext_ticket ON ticket_external_links(ticket_id);

        -- 统一身份认证：认证提供者配置
        CREATE TABLE IF NOT EXISTS auth_providers (
            id TEXT PRIMARY KEY,
            name TEXT NOT NULL,
            provider_type TEXT NOT NULL,
            enabled INTEGER DEFAULT 1,
            config TEXT NOT NULL DEFAULT '{}',
            created_at TEXT DEFAULT (datetime('now','localtime')),
            updated_at TEXT DEFAULT (datetime('now','localtime'))
        );

        -- 统一身份认证：外部身份映射
        CREATE TABLE IF NOT EXISTS auth_identity_mappings (
            id TEXT PRIMARY KEY,
            provider_id TEXT NOT NULL REFERENCES auth_providers(id),
            external_user_id TEXT NOT NULL,
            smartcs_agent_id TEXT REFERENCES agents(id),
            external_username TEXT DEFAULT '',
            external_email TEXT DEFAULT '',
            created_at TEXT DEFAULT (datetime('now','localtime')),
            last_login_at TEXT DEFAULT '',
            UNIQUE(provider_id, external_user_id)
        );
        ''')
        # Create default IM adapter placeholders (after init_db completes, within app context)
        try:
            for (aid, aname, aplat) in [('default_wecom','企业微信（默认）','wecom'),('default_dingtalk','钉钉（默认）','dingtalk')]:
                get_db().execute("INSERT OR IGNORE INTO im_adapters(id,name,platform,enabled,config) VALUES(?,?,?,?,?)",
                    (aid, aname, aplat, 0, '{}'))
            get_db().commit()
        except: pass
        try: db.execute("ALTER TABLE service_tickets ADD COLUMN ticket_number TEXT DEFAULT ''")
        except: pass
        try: db.execute("ALTER TABLE agent_profiles ADD COLUMN employee_id TEXT DEFAULT ''")
        except: pass
        try: db.execute("ALTER TABLE agent_profiles ADD COLUMN company TEXT DEFAULT ''")
        except: pass
        try: db.execute("ALTER TABLE agent_profiles ADD COLUMN agent_level INTEGER DEFAULT 1")
        except: pass
        try: db.execute("ALTER TABLE service_tickets ADD COLUMN transferred_from TEXT DEFAULT ''")
        except: pass
        try: db.execute("ALTER TABLE customer_profiles ADD COLUMN employee_id TEXT DEFAULT ''")
        except: pass
        try: db.execute("ALTER TABLE customer_profiles ADD COLUMN department TEXT DEFAULT ''")
        except: pass
        try: db.execute("ALTER TABLE customers ADD COLUMN email TEXT DEFAULT ''")
        except: pass
        try: db.execute("ALTER TABLE customers ADD COLUMN password_hash TEXT DEFAULT ''")
        except: pass
        try: db.execute("ALTER TABLE service_tickets ADD COLUMN admin_remarks TEXT DEFAULT ''")
        except: pass
        try: db.execute("ALTER TABLE service_tickets ADD COLUMN close_reason TEXT DEFAULT ''")
        except: pass
        try: db.execute("ALTER TABLE service_tickets ADD COLUMN level INTEGER DEFAULT 1")
        except: pass
        try: db.execute("ALTER TABLE service_tickets ADD COLUMN reopened_at TEXT DEFAULT ''")
        except: pass
        try: db.execute("ALTER TABLE service_tickets ADD COLUMN reopened_count INTEGER DEFAULT 0")
        except: pass
        try: db.execute("ALTER TABLE service_tickets ADD COLUMN escalation_id TEXT DEFAULT ''")
        except: pass
        try: db.execute("ALTER TABLE service_tickets ADD COLUMN rated_at TEXT DEFAULT ''")
        except: pass
        try: db.execute("ALTER TABLE knowledge_files ADD COLUMN status TEXT DEFAULT 'approved'")
        except: pass
        try: db.execute("ALTER TABLE knowledge_files ADD COLUMN submitted_by TEXT DEFAULT ''")
        except: pass
        try: db.execute("ALTER TABLE knowledge_files ADD COLUMN review_notes TEXT DEFAULT ''")
        except: pass
        try: db.execute("ALTER TABLE knowledge_files ADD COLUMN title TEXT DEFAULT ''")
        except: pass
        try: db.execute("ALTER TABLE knowledge_files ADD COLUMN updated_at TEXT DEFAULT (datetime('now','localtime'))")
        except: pass
        try: db.execute("ALTER TABLE knowledge_files ADD COLUMN created_by TEXT DEFAULT ''")
        except: pass
        try: db.execute("ALTER TABLE knowledge_files ADD COLUMN updated_by TEXT DEFAULT ''")
        except: pass
        try: db.execute("ALTER TABLE knowledge_files ADD COLUMN tags TEXT DEFAULT '{}'")
        except: pass
        # === P0 State Machine Migration: old statuses → new statuses ===
        try:
            db.execute("UPDATE service_tickets SET status='processing' WHERE status IN ('open','assigned')")
            db.execute("UPDATE service_tickets SET status='resolved' WHERE status='confirmed'")
            db.execute("UPDATE service_tickets SET status='closed' WHERE status IN ('resolved','closed')")
        except:
            pass

        idx_defs = [
            'CREATE INDEX IF NOT EXISTS idx_tickets_agent_status ON service_tickets(agent_id, status)',
            'CREATE INDEX IF NOT EXISTS idx_tickets_customer_status ON service_tickets(customer_id, status)',
            'CREATE INDEX IF NOT EXISTS idx_tickets_created ON service_tickets(created_at DESC)',
            'CREATE INDEX IF NOT EXISTS idx_messages_conv_time ON messages(conversation_id, created_at)',
            'CREATE INDEX IF NOT EXISTS idx_tickets_tk_number ON service_tickets(ticket_number)',
        ]
        for d in idx_defs:
            try: db.execute(d)
            except: pass
        # === 三员管理：迁移管理员到 superadmin ===
        try:
            db.execute("UPDATE agents SET role='superadmin' WHERE role='admin'")
        except: pass
        try:
            db.execute("ALTER TABLE agents ADD COLUMN password_changed_at TEXT DEFAULT ''")
        except: pass
        
        # === 三员管理：创建默认三员账号 ===
        three_roles = [
            ('sysadmin', '系统管理员', 'sysadmin@smartcs.com', generate_password_hash('SysAdmin@2026')),
            ('secadmin', '安全管理员', 'secadmin@smartcs.com', generate_password_hash('SecAdmin@2026')),
            ('audadmin', '审计管理员', 'audadmin@smartcs.com', generate_password_hash('AudAdmin@2026')),
        ]
        for trole, tname, temail, tpwhash in three_roles:
            existing = db.execute("SELECT id FROM agents WHERE email=?", (temail,)).fetchone()
            if not existing:
                taid = 'agent-' + uuid.uuid4().hex[:12]
                db.execute("INSERT INTO agents(id,name,email,password_hash,role) VALUES(?,?,?,?,?)",
                          (taid, tname, temail, tpwhash, trole))
                db.execute("INSERT INTO agent_profiles(id,agent_id,display_name,agent_level) VALUES(?,?,?,?)",
                          ('ap-' + uuid.uuid4().hex[:12], taid, tname, 1))
        
        pwd_hash = hashlib.sha256(f'admin:{ADMIN_PASSWORD}'.encode()).hexdigest()
        db.execute("INSERT OR IGNORE INTO agents(id,name,email,password_hash,role) VALUES(?,?,?,?,?)",
                   ('admin-001', '管理员', 'admin@smartcs.com', pwd_hash, 'superadmin'))
        db.execute("INSERT OR IGNORE INTO agents(id,name,email,password_hash,role) VALUES(?,?,?,?,?)",
                   ('agent-001', '工程师01', 'agent@smartcs.com', pwd_hash, 'agent'))
        for _lvl_id, _lvl_name, _lvl_email, _lvl_level, _lvl_display, _lvl_title in [
            ('agent-l1', 'L1工程师', 'agent_l1@smartcs.com', 1, 'L1工程师', '初级工程师'),
            ('agent-l2', 'L2工程师', 'agent_l2@smartcs.com', 2, 'L2工程师', '高级工程师'),
            ('agent-l3', 'L3工程师', 'agent_l3@smartcs.com', 3, 'L3工程师', '专家工程师'),
            ('agent-l4', 'L4工程师', 'agent_l4@smartcs.com', 4, 'L4工程师', '首席工程师'),
        ]:
            db.execute("INSERT OR IGNORE INTO agents(id,name,email,password_hash,role) VALUES(?,?,?,?,?)",
                       (_lvl_id, _lvl_name, _lvl_email, pwd_hash, 'agent'))
            db.execute("INSERT OR IGNORE INTO agent_profiles(id,agent_id,display_name,department,title,agent_level) VALUES(?,?,?,?,?,?)",
                       (f'ap-{_lvl_id}', _lvl_id, _lvl_display, '技术支持', _lvl_title, _lvl_level))
        db.execute("INSERT OR IGNORE INTO agent_profiles(id,agent_id,display_name,department,title) VALUES(?,?,?,?,?)",
                   ('ap-001', 'agent-001', '工程师小张', '技术支持', '高级工程师'))

        for i, nm in enumerate(['系统问题','咨询问题','功能建议','售后问题','其他']):
            db.execute("INSERT OR IGNORE INTO close_reasons(id,name,sort_order) VALUES(?,?,?)",
                       (f'cr-{uuid.uuid4().hex[:8]}', nm, i))
        default_cfg = {'api_base_url':API_BASE_URL,'api_key':DASHSCOPE_API_KEY or '','model_name':MODEL_NAME,'admin_password':ADMIN_PASSWORD,'webhook_timeout':'10','auto_close_min':'20','auto_rate_hours':'24','ticket_search_max_days':'365',
                       'level_names': json.dumps({'1':'初级工程师','2':'高级工程师','3':'专家工程师','4':'首席工程师'}, ensure_ascii=False)}
        for k, v in default_cfg.items():
            db.execute("INSERT OR IGNORE INTO system_config(key,value) VALUES(?,?)", (k, v))
        # 安全配置键
        security_cfg = [
            ('password_min_length', '8'),
            ('password_require_upper', 'true'),
            ('password_expire_days', '90'),
            ('login_max_attempts', '5'),
            ('login_lockout_minutes', '15'),
            ('audit_log_retention_days', '365'),
        ]
        for k, v in security_cfg:
            db.execute("INSERT OR IGNORE INTO system_config(key,value) VALUES(?,?)", (k, v))
        # 品牌配置键
        brand_cfg = [
            ('brand_name', 'SmartCS 智能客服'),
            ('brand_short', 'SmartCS'),
            ('brand_primary_color', '#1a73e8'),
            ('brand_logo_path', '/static/icon-192.png'),
            ('brand_favicon_path', '/static/favicon.ico'),
        ]
        for k, v in brand_cfg:
            db.execute("INSERT OR IGNORE INTO system_config(key,value) VALUES(?,?)", (k, v))
        # 系统行为配置键
        sys_behavior_cfg = [
            ('session_lifetime', '28800'),
            ('session_idle_timeout', '1800'),
            ('auto_check_interval', '300'),
            ('pagination_per_page', '50'),
            ('max_upload_size_mb', '50'),
        ]
        for k, v in sys_behavior_cfg:
            db.execute("INSERT OR IGNORE INTO system_config(key,value) VALUES(?,?)", (k, v))
        db.commit()
init_db()

def gen_id(prefix=''): return f"{prefix}{uuid.uuid4().hex[:12]}"

def generate_ticket_number():
    db = get_db()
    today = datetime.now().strftime('%Y%m%d')
    row = db.execute("SELECT seq FROM ticket_seq WHERE date_str=?", (today,)).fetchone()
    if row:
        seq = row['seq'] + 1
        db.execute("UPDATE ticket_seq SET seq=? WHERE date_str=?", (seq, today))
    else:
        seq = 1
        db.execute("INSERT INTO ticket_seq(date_str,seq) VALUES(?,?)", (today, seq))
    db.commit()
    return f"tk{today}{seq:06d}"

def get_or_create_customer():
    cid = session.get('customer_id')
    if cid and get_db().execute("SELECT id FROM customers WHERE id=?", (cid,)).fetchone():
        return cid
    return None

# ====== 搜索知识库和AI调用（保持不变）======
def search_knowledge(query):
    results = []
    query_lower = query.lower()
    for sw in ['什么','是','的','怎么','如何','为什么','能','吗','呢','请','问','关于','介绍','一下','有','哪些','多少','定义','意思','含义','说明','解释','是什么']:
        query_lower = query_lower.replace(sw, ' ')
    keywords = [w.strip() for w in query_lower.split() if len(w.strip()) > 1]
    for pat in [r'(.+?)是什么',r'(.+?)的定义',r'(.+?)的意思',r'介绍一下(.+?)',r'关于(.+?)']:
        m = re.search(pat, query)
        if m: keywords.append(m.group(1).strip().lower())
    # Only search approved knowledge files
    db = get_db()
    approved = db.execute("SELECT filename FROM knowledge_files WHERE status='approved'").fetchall()
    approved_filenames = set(r['filename'] for r in approved)
    for fp in glob.glob(os.path.join(KNOWLEDGE_DIR, '*.md')):
        try:
            basename = os.path.basename(fp)
            if basename not in approved_filenames:
                continue
            with open(fp, 'r', encoding='utf-8') as f: text = f.read().lower()
            for kw in set(keywords):
                if kw and kw in text:
                    with open(fp, 'r', encoding='utf-8') as f:
                        results.append({'file':basename,'snippet':f.read()[:2000]})
                    break
        except: pass
    return results

def call_llm(messages, timeout=30):
    cfg_key, cfg_url, cfg_model = DASHSCOPE_API_KEY, API_BASE_URL, MODEL_NAME
    try:
        rows = get_db().execute("SELECT key,value FROM system_config WHERE key IN ('api_key','api_base_url','model_name')").fetchall()
        for r in rows:
            if r['key'] == 'api_key' and r['value']: cfg_key = r['value']
            elif r['key'] == 'api_base_url' and r['value']: cfg_url = r['value']
            elif r['key'] == 'model_name' and r['value']: cfg_model = r['value']
    except: pass
    if not cfg_key: return None, "API Key 未配置"
    try:
        resp = requests.post(cfg_url,
            headers={'Authorization':f'Bearer {cfg_key}','Content-Type':'application/json'},
            json={'model':cfg_model,'messages':messages,'temperature':0.7,'max_tokens':2000}, timeout=timeout)
        resp.raise_for_status()
        return resp.json()['choices'][0]['message']['content'], None
    except requests.Timeout: return None, "AI 响应超时"
    except Exception as e: return None, f"AI 异常: {str(e)[:100]}"

def need_escalation(reply):
    return any(p in reply.lower() for p in ['无法回答','没有相关信息','知识库中','无法提供','我不确定','建议联系','转接人工','sorry','cannot answer','not sure','知识库中没有'])

# ====== 统一状态转换 + 审计日志 ======

def log_audit(ticket_id, action, actor_id='system', actor_name='系统', detail=None):
    """写入操作审计日志"""
    try:
        ip = request.remote_addr or ''
        get_db().execute(
            "INSERT INTO audit_log(id,ticket_id,actor_id,actor_name,action,detail,ip_address) VALUES(?,?,?,?,?,?,?)",
            (gen_id('aud-'), ticket_id, actor_id, actor_name, action,
             json.dumps(detail or {}, ensure_ascii=False), ip)
        )
        get_db().commit()
    except:
        pass

def transition_ticket(ticket_id, to_status, actor_id='system', actor_name='系统', extra=None):
    """统一工单状态转换入口（P0新状态机），带状态验证 + 审计日志
    
    新状态机: created → processing → resolved → rated → closed
    转换规则:
      created   → processing (transfer/escalate)
      created   → closed     (auto-timeout, force-close)
      processing → resolved  (agent resolve)
      processing → closed    (force-close)
      resolved  → rated      (user confirm+rate, auto-rate)
      resolved  → closed     (force-close)
      rated     → closed     (agent close, force-close)
    """
    db = get_db()
    t = db.execute("SELECT status, agent_id, customer_id FROM service_tickets WHERE id=?", (ticket_id,)).fetchone()
    if not t:
        return False, '工单不存在'
    
    from_status = t['status']
    # New state machine transitions
    valid = {
        'created': ['processing', 'closed'],
        'processing': ['resolved', 'closed'],
        'resolved': ['rated', 'closed'],
        'rated': ['closed'],
        'closed': [],
    }
    if to_status not in valid.get(from_status, []):
        return False, f'不允许的状态转换: {from_status} → {to_status}'
    
    # 构建更新语句
    sets = ["status=?", "updated_at=datetime('now','localtime')"]
    params = [to_status]
    
    if from_status == 'created' and to_status == 'processing':
        # Transfer to human - set assigned_at
        sets.append("assigned_at=datetime('now','localtime')")
        if extra and extra.get('agent_id'):
            sets.append("agent_id=?"); params.append(extra['agent_id'])
    elif to_status == 'resolved':
        sets.append("resolved_at=datetime('now','localtime')")
        if extra:
            if extra.get('close_reason'):
                sets.append("close_reason=?"); params.append(extra['close_reason'])
            if extra.get('resolution_notes'):
                sets.append("resolution_notes=?"); params.append(extra['resolution_notes'])
    elif to_status == 'rated':
        if from_status == 'resolved':
            sets.append("rated_at=datetime('now','localtime')")
            if extra:
                if 'customer_rating' in (extra or {}):
                    sets.append("customer_rating=?"); params.append(extra['customer_rating'])
                if 'customer_feedback' in (extra or {}):
                    sets.append("customer_feedback=?"); params.append(extra['customer_feedback'])
    elif to_status == 'closed':
        if from_status == 'rated':
            sets.append("closed_at=datetime('now','localtime')")
        elif from_status in ('created', 'processing', 'resolved'):
            # force-close or auto-timeout
            sets.append("closed_at=datetime('now','localtime')")
            if extra and extra.get('close_reason'):
                sets.append("close_reason=?"); params.append(extra['close_reason'])
    
    params.append(ticket_id)
    db.execute(f"UPDATE service_tickets SET {','.join(sets)} WHERE id=?", params)
    db.commit()
    
    log_audit(ticket_id, f'ticket.{to_status}', actor_id, actor_name,
              {'from':from_status, 'to':to_status, 'extra':extra or {}})
    # Fire webhook events
    emit_event('ticket.' + to_status, {'ticket_id': ticket_id, 'status': to_status, 'agent_id': t['agent_id']})
    return True, f'{from_status} → {to_status}'


# ====== Phase 4: Webhook Event System ======

def emit_event(event_type, ticket_data):
    """Fire event to all matching webhooks + IM notifications"""
    db = get_db()
    hooks = db.execute(
        "SELECT * FROM webhooks WHERE enabled=1 AND (events=? OR events LIKE ?)",
        ('["*"]', '%"' + event_type + '"%')
    ).fetchall()
    for hook in hooks:
        t = threading.Thread(target=_deliver_webhook, args=(dict(hook), event_type, ticket_data))
        t.daemon = True
        t.start()
    # Also send IM notifications
    try:
        _send_im_event_notification(event_type, ticket_data)
    except:
        pass

def _deliver_webhook(hook, event_type, payload_data):
    """Deliver webhook with retry logic"""
    payload = {
        'event': event_type,
        'timestamp': datetime.now().isoformat(),
        'data': payload_data
    }
    log_id = None
    for attempt in range(1, hook.get('retry', 3) + 1):
        try:
            headers = {
                'Content-Type': 'application/json',
                'User-Agent': 'SmartCS-Webhook/1.0',
                'X-Webhook-Event': event_type,
            }
            if hook.get('secret'):
                headers['X-Webhook-Secret'] = hook['secret']
            resp = requests.post(
                hook['url'],
                json=payload,
                timeout=hook.get('timeout', 10),
                headers=headers
            )
            status = 'delivered' if 200 <= resp.status_code < 300 else 'failed'
            _log_webhook_delivery(hook['id'], event_type, payload, status, resp.status_code, resp.text[:500], attempt, None, log_id)
            return
        except Exception as e:
            error_msg = str(e)
            if attempt == hook.get('retry', 3):
                _log_webhook_delivery(hook['id'], event_type, payload, 'failed', 0, '', attempt, error_msg, log_id)
            else:
                time.sleep(2 ** attempt)

def _log_webhook_delivery(webhook_id, event_type, payload, status, status_code, response_body, attempt, error, log_id):
    """Insert or update webhook log record"""
    db = get_db()
    payload_str = json.dumps(payload, ensure_ascii=False)
    if log_id:
        db.execute('UPDATE webhook_logs SET status=?, status_code=?, response_body=?, attempt=?, error=? WHERE id=?',
                   (status, status_code, response_body, attempt, error, log_id))
    else:
        lid = 'whl-' + uuid.uuid4().hex[:12]
        db.execute('INSERT INTO webhook_logs(id, webhook_id, event_type, payload, status, status_code, response_body, attempt, error) VALUES(?,?,?,?,?,?,?,?,?)',
                   (lid, webhook_id, event_type, payload_str, status, status_code, response_body, attempt, error or ''))
    db.commit()

# ====== End Phase 4 Webhooks ======

SYSTEM_PROMPT = """你是一个专业的智能客服助手。只根据知识库内容回答。回答要简洁友好。如果没有相关知识，告诉用户当前知识库中没有相关信息，并提示用户可点击"转人工"按钮联系IT工程师。"""

@app.route('/')
def index(): return render_template('chat.html')

@app.route('/login')
def customer_login_page():
    return render_template('user_login.html')

@app.route('/register', methods=['GET', 'POST'])
def customer_register_page():
    if request.method == 'POST':
        return redirect('/login')
    return redirect('/login')


@app.route('/manifest.json')
def manifest():
    cfg = get_brand_config()
    return jsonify({
        "name": cfg.get('brand_name', 'SmartCS 智能客服'),
        "short_name": cfg.get('brand_short', 'SmartCS'),
        "start_url": "/",
        "display": "standalone",
        "orientation": "portrait",
        "background_color": cfg.get('brand_primary_color', '#1a73e8'),
        "theme_color": cfg.get('brand_primary_color', '#1a73e8'),
        "icons": [{"src": cfg.get('brand_logo_path', '/static/icon-192.png'), "sizes": "192x192", "type": "image/png"}]
    })

@app.route('/sw.js')
def service_worker():
    cfg = get_brand_config()
    logo_path = cfg.get('brand_logo_path', '/static/icon-192.png')
    sw_content = """const CACHE_NAME = 'smartcs-v4.1';
const urlsToCache = [
  '%s'
];

self.addEventListener('install', function(event) {
  self.skipWaiting();
  event.waitUntil(
    caches.open(CACHE_NAME).then(function(cache) {
      return cache.addAll(urlsToCache);
    })
  );
});

self.addEventListener('fetch', function(event) {
  if (event.request.method === 'GET') {
    if (event.request.url.includes('/api/') || event.request.headers.get('Accept') === 'application/json') {
      event.respondWith(fetch(event.request));
      return;
    }
  }
  event.respondWith(
    caches.match(event.request).then(function(response) {
      return response || fetch(event.request).then(function(networkResponse) {
        if (event.request.url.includes('/static/')) {
          var responseClone = networkResponse.clone();
          caches.open(CACHE_NAME).then(function(cache) {
            cache.put(event.request, responseClone);
          });
        }
        return networkResponse;
      });
    })
  );
});

self.addEventListener('activate', function(event) {
  event.waitUntil(
    caches.keys().then(function(cacheNames) {
      return Promise.all(
        cacheNames.map(function(name) {
          if (name !== CACHE_NAME) return caches.delete(name);
        })
      );
    })
  );
});""" % logo_path
    return Response(sw_content, mimetype='application/javascript',
                  headers={'Cache-Control': 'no-cache', 'Service-Worker-Allowed': '/'})

@app.route('/api/chat', methods=['POST'])
def chat():
    data = request.get_json()
    msg = data.get('message','').strip()
    conv_id = data.get('conversation_id','')
    if not msg: return jsonify({'error':'请输入消息'}),400
    cid = get_or_create_customer()
    if not cid:
        return jsonify({'error':'未登录','login_required':True}),401
    db = get_db()
    conv = None
    if conv_id:
        conv = db.execute("SELECT id,status FROM conversations WHERE id=? AND status IN ('active','escalated')", (conv_id,)).fetchone()
        if not conv: conv_id = ''
    if not conv_id:
        conv_id = gen_id('conv-')
        db.execute("INSERT INTO conversations(id,customer_id) VALUES(?,?)", (conv_id, cid))
    db.execute("INSERT INTO messages(id,conversation_id,role,content) VALUES(?,?,?,?)", (gen_id('msg-'), conv_id, 'user', msg))
    log_audit('', 'customer.send_message', cid, session.get('customer_name','客户'),
              {'conversation_id': conv_id, 'msg_preview': msg[:80]})

    # === P0: Create ticket on first user message if none exists ===
    is_first_msg = db.execute("SELECT COUNT(*) as c FROM messages WHERE conversation_id=? AND role='user'", (conv_id,)).fetchone()['c'] == 1
    existing_ticket = db.execute("SELECT id, status FROM service_tickets WHERE conversation_id=?", (conv_id,)).fetchone()
    if is_first_msg and not existing_ticket:
        tk_id = gen_id('tk-')
        tk_number = generate_ticket_number()
        db.execute("""INSERT INTO service_tickets(id,conversation_id,customer_id,issue_description,
                     status,ticket_number) VALUES(?,?,?,?,'created',?)""",
                   (tk_id, conv_id, cid, msg[:200], tk_number))
        db.commit()
        existing_ticket = {'id': tk_id, 'status': 'created'}

    # Handle transfer to human
    is_transfer = '转人工' in msg or '转接人工' in msg or 'IT工程师' in msg
    if is_transfer:
        if existing_ticket:
            tk_id = existing_ticket['id']
            if existing_ticket['status'] == 'created':
                transition_ticket(tk_id, 'processing', cid, session.get('customer_name','客户'),
                                {'agent_id': '', 'close_reason': '用户主动要求转人工'})
            elif existing_ticket['status'] == 'processing':
                pass  # already processing
        else:
            # Very rare: create ticket + transition
            tk_id = gen_id('tk-')
            tk_number = generate_ticket_number()
            db.execute("""INSERT INTO service_tickets(id,conversation_id,customer_id,issue_description,
                         status,ticket_number) VALUES(?,?,?,?,'processing',?)""",
                       (tk_id, conv_id, cid, msg[:200], tk_number))
        get_db().execute("UPDATE conversations SET status='escalated',updated_at=datetime('now','localtime') WHERE id=?", (conv_id,))
        db.commit()
        return jsonify({'reply':'✅ 已为您转接IT工程师，请稍候...','escalated':True,'conversation_id':conv_id})

    # If already escalated, skip AI - just relay message to agent
    if conv and conv['status'] == 'escalated':
        db.execute("UPDATE conversations SET updated_at=datetime('now','localtime') WHERE id=?", (conv_id,))
        db.commit()
        return jsonify({'reply':'📝 消息已发送给IT工程师，请等待回复...','conversation_id':conv_id,'escalated':True})

    knowledge = search_knowledge(msg)
    msgs = [{'role':'system','content':SYSTEM_PROMPT}]
    if knowledge:
        msgs.append({'role':'system','content':"知识库内容：\n"+'\n\n---\n\n'.join([f"## {k['file']}\n{k['snippet']}" for k in knowledge[:3]])})
    history = db.execute("SELECT role,content FROM messages WHERE conversation_id=? ORDER BY created_at DESC LIMIT 6", (conv_id,)).fetchall()[1:]
    for h in reversed(history):
        role_map = {"bot":"assistant","user":"user","agent":"assistant","system":"system"}
        msgs.append({"role": role_map.get(h["role"], "user"), "content": h["content"]})
    msgs.append({'role':'user','content':msg})
    reply, error = call_llm(msgs)
    if error:
        db.execute("INSERT INTO messages(id,conversation_id,role,content) VALUES(?,?,?,?)", (gen_id('msg-'), conv_id, 'bot', f'⚠️ {error}'))
        db.commit()
        return jsonify({'reply':f'⚠️ {error}','error':error,'conversation_id':conv_id})
    escalated = False
    # Auto-escalate if AI can't answer (transition created → processing)
    if need_escalation(reply) and existing_ticket and existing_ticket['status'] == 'created':
        transition_ticket(existing_ticket['id'], 'processing', cid, session.get('customer_name','客户'),
                        {'close_reason': 'AI无法回答，自动转人工'})
        get_db().execute("UPDATE conversations SET status='escalated',updated_at=datetime('now','localtime') WHERE id=?", (conv_id,))
        escalated = True
    db.execute("INSERT INTO messages(id,conversation_id,role,content) VALUES(?,?,?,?)", (gen_id('msg-'), conv_id, 'bot', reply))
    db.execute("UPDATE conversations SET updated_at=datetime('now','localtime') WHERE id=?", (conv_id,))
    db.commit()
    return jsonify({'reply':reply,'conversation_id':conv_id,'escalated':escalated,'knowledge_used':len(knowledge)})


@app.route('/api/chat/upload', methods=['POST'])
def chat_upload():
    file = request.files.get('file')
    conv_id = request.form.get('conversation_id','')
    if not file: return jsonify({'error':'请选择文件'}),400
    if file.filename == '': return jsonify({'error':'文件名为空'}),400
    ts = datetime.now().strftime('%Y%m%d%H%M%S')
    fn = f"{ts}_{uuid.uuid4().hex[:8]}.{file.filename.rsplit('.',1)[-1]}"
    os.makedirs(os.path.join(UPLOAD_DIR, 'images'), exist_ok=True)
    path = os.path.join(UPLOAD_DIR, 'images', fn)
    file.save(path)
    cid = get_or_create_customer()
    if not cid:
        return jsonify({'error':'未登录','login_required':True}),401
    db = get_db()
    if not conv_id:
        conv_id = gen_id('conv-')
        db.execute("INSERT INTO conversations(id,customer_id) VALUES(?,?)", (conv_id, cid))
    url = f"/uploads/images/{fn}"
    db.execute("INSERT INTO messages(id,conversation_id,role,content,image_url) VALUES(?,?,?,?,?)",
               (gen_id('msg-'), conv_id, 'user', '[图片]', url))
    db.execute("UPDATE conversations SET updated_at=datetime('now','localtime') WHERE id=?", (conv_id,))
    db.commit()
    return jsonify({'ok':True,'url':url,'conversation_id':conv_id,'escalated':False})

@app.route('/api/chat/history', methods=['GET'])
def chat_history():
    conv_id = request.args.get('conversation_id','')
    msgs = get_db().execute("SELECT role,content,image_url,created_at FROM messages WHERE conversation_id=? ORDER BY created_at", (conv_id,)).fetchall()
    return jsonify([dict(m) for m in msgs])

@app.route('/api/chat/new', methods=['POST'])
def new_conversation():
    session.pop('customer_id', None)
    return jsonify({'ok':True})

# ====== 客户档案 ======

@app.route('/api/customer/register', methods=['POST'])
def customer_register():
    """Register a new customer account with email+password."""
    data = request.get_json() or {}
    email = (data.get('email') or '').strip().lower()
    password = data.get('password', '')
    name = (data.get('name') or '').strip()
    
    if not email or not password:
        return jsonify({'error': '邮箱和密码不能为空'}), 400
    if len(password) < 6:
        return jsonify({'error': '密码不能少于6位'}), 400
    if not re.match(r'^[^@\s]+@[^@\s]+\.[^@\s]+$', email):
        return jsonify({'error': '邮箱格式不正确'}), 400
    
    db = get_db()
    existing = db.execute("SELECT id FROM customers WHERE email=?", (email,)).fetchone()
    if existing:
        return jsonify({'error': '该邮箱已注册'}), 409
    
    cid = session.get('customer_id')
    if cid and db.execute("SELECT id FROM customers WHERE id=?", (cid,)).fetchone():
        pass  # use existing anonymous customer
    else:
        cid = None
    
    if cid:
        db.execute("UPDATE customers SET email=?, password_hash=?, name=? WHERE id=?", 
                   (email, generate_password_hash(password), name or email.split('@')[0], cid))
    else:
        cid = 'cust-' + uuid.uuid4().hex[:12]
        db.execute("INSERT INTO customers(id, name, email, password_hash) VALUES(?,?,?,?)",
                   (cid, name or email.split('@')[0], email, generate_password_hash(password)))
        session['customer_id'] = cid
    
    session['customer_email'] = email
    session['customer_name'] = name or email.split('@')[0]
    db.commit()
    log_audit('', 'customer.register', cid, session['customer_name'],
              {'email': email, 'action': 'register'})
    return jsonify({'ok': True, 'customer': {'name': session['customer_name'], 'email': email}})


def _record_login(user_id, user_name, user_type, success, fail_reason=''):
    try:
        get_db().execute(
            "INSERT INTO login_log(id,user_id,user_name,user_type,ip_address,user_agent,success,fail_reason) VALUES(?,?,?,?,?,?,?,?)",
            (gen_id('ll-'), user_id, user_name, user_type,
             request.remote_addr or '', request.headers.get('User-Agent','')[:200] or '',
             1 if success else 0, fail_reason)
        )
        get_db().commit()
    except:
        pass


@app.route('/api/customer/login', methods=['POST'])
def customer_login():
    """Login with email+password."""
    data = request.get_json() or {}
    email = (data.get('email') or '').strip().lower()
    password = data.get('password', '')
    
    if not email or not password:
        return jsonify({'error': '邮箱和密码不能为空'}), 400
    
    db = get_db()
    customer = db.execute("SELECT id, name, email, password_hash FROM customers WHERE email=?", (email,)).fetchone()
    if not customer or not customer['password_hash']:
        _record_login('', email, 'customer', False, '用户不存在或未设置密码')
        return jsonify({'error': '邮箱或密码错误'}), 401
    if not check_password_hash(customer['password_hash'], password):
        _record_login(customer['id'], customer['name'], 'customer', False, '密码错误')
        return jsonify({'error': '邮箱或密码错误'}), 401
    
    session['customer_id'] = customer['id']
    session['customer_email'] = customer['email']
    session['customer_name'] = customer['name']
    db.commit()
    _record_login(customer['id'], customer['name'], 'customer', True)
    log_audit('', 'customer.login', customer['id'], customer['name'],
              {'email': customer['email']})
    return jsonify({'ok': True, 'customer': {'name': customer['name'], 'email': customer['email']}})


@app.route('/api/customer/logout', methods=['POST'])
def customer_logout():
    """Logout customer."""
    cid = session.get('customer_id')
    name = session.get('customer_name', '')
    if cid:
        log_audit('', 'customer.logout', cid, name, {})
    session.pop('customer_email', None)
    session.pop('customer_name', None)
    return jsonify({'ok': True})


@app.route('/api/customer/me', methods=['GET'])
def customer_me():
    """Return current customer profile or logged_in: false."""
    email = session.get('customer_email')
    name = session.get('customer_name')
    cid = session.get('customer_id')
    if email and cid:
        db = get_db()
        c = db.execute("SELECT id, name, email FROM customers WHERE id=?", (cid,)).fetchone()
        if c:
            return jsonify({
                'logged_in': True,
                'customer': {'name': c['name'], 'email': c['email']}
            })
    return jsonify({'logged_in': False})

@app.route('/api/customer/profile', methods=['GET','POST'])
def customer_profile():
    cid = get_or_create_customer()
    if not cid:
        return jsonify({'error':'未登录','login_required':True}),401
    db = get_db()
    if request.method == 'POST':
        data = request.get_json()
        prof = db.execute("SELECT id FROM customer_profiles WHERE customer_id=?", (cid,)).fetchone()
        if prof:
            db.execute("UPDATE customer_profiles SET name=?,phone=?,company=?,updated_at=datetime('now','localtime') WHERE customer_id=?", 
                       (data.get('name',''), data.get('phone',''), data.get('company',''), cid))
        else:
            db.execute("INSERT INTO customer_profiles(id,customer_id,name,phone,company) VALUES(?,?,?,?,?)",
                       (gen_id('cp-'), cid, data.get('name',''), data.get('phone',''), data.get('company','')))
        db.execute("UPDATE customers SET name=? WHERE id=?", (data.get('name','') or '游客', cid))
        db.commit()
        return jsonify({'ok':True})
    prof = db.execute("SELECT * FROM customer_profiles WHERE customer_id=?", (cid,)).fetchone()
    binding = get_db().execute("""
        SELECT ap.display_name, ap.department, ap.title, ag.name as agent_name
        FROM customer_agent_bindings b JOIN agents ag ON b.agent_id = ag.id
        LEFT JOIN agent_profiles ap ON b.agent_id = ap.agent_id
        WHERE b.customer_id=? AND b.status='active'""", (cid,)).fetchone()
    r = {'customer_id':cid}
    if prof: r.update(dict(prof))
    if binding:
        r['agent'] = dict(binding)
    else:
        # Fall back to active ticket's agent
        tk = db.execute("SELECT ag.name as agent_name FROM service_tickets t JOIN agents ag ON t.agent_id = ag.id WHERE t.customer_id=? AND t.status NOT IN ('closed','resolved','rated') ORDER BY t.created_at DESC LIMIT 1", (cid,)).fetchone()
        if tk: r['agent_name'] = tk['agent_name']
    return jsonify(r)

@app.route('/api/customer/tickets', methods=['GET'])
def customer_tickets():
    cid = get_or_create_customer()
    if not cid:
        return jsonify({'error':'未登录','login_required':True}),401
    tickets = get_db().execute("""
        SELECT t.*, ag.name as agent_name, ap.display_name as agent_display
        FROM service_tickets t LEFT JOIN agents ag ON t.agent_id = ag.id
        LEFT JOIN agent_profiles ap ON t.agent_id = ap.agent_id
        WHERE t.customer_id=? ORDER BY t.created_at DESC""", (cid,)).fetchall()
    log_audit('', 'customer.view_tickets', cid, session.get('customer_name','客户'),
              {'count': len(tickets)})
    return jsonify([dict(t) for t in tickets])

@app.route('/api/customer/tickets/transfer', methods=['POST'])
def customer_request_human():
    """Customer requests human transfer: created → processing"""
    cid = get_or_create_customer()
    if not cid:
        return jsonify({'error':'未登录','login_required':True}),401
    data = request.get_json()
    conv_id = data.get('conversation_id','')
    if not conv_id:
        return jsonify({'error':'缺少会话ID'}),400
    db = get_db()
    t = db.execute("SELECT id, status, agent_id FROM service_tickets WHERE conversation_id=? AND customer_id=? ORDER BY created_at DESC LIMIT 1",
                   (conv_id, cid)).fetchone()
    if not t:
        return jsonify({'error':'工单不存在'}),404
    if t['status'] != 'created':
        # Already processing or resolved
        return jsonify({'escalated': True, 'message': '工单已在处理中'})
    ok, _ = transition_ticket(t['id'], 'processing', cid, '客户', {'reason': '用户请求转人工'})
    # Assign to first available agent
    first_agent = db.execute(
        "SELECT ag.id FROM agents ag LEFT JOIN agent_profiles ap ON ag.id=ap.agent_id WHERE ag.role='agent' ORDER BY ap.agent_level ASC LIMIT 1"
    ).fetchone()
    if first_agent:
        db.execute("UPDATE service_tickets SET agent_id=?,first_read_at='' WHERE id=?", (first_agent['id'], t['id']))
    db.commit()
    return jsonify({'escalated': ok, 'message': '已转人工，请稍候...'})


@app.route('/api/customer/tickets/confirm', methods=['POST'])
def customer_confirm():
    """P0: User confirms resolution → transitions resolved ticket to rated status"""
    data = request.get_json()
    tk_id = data.get('ticket_id','')
    db = get_db()
    t = db.execute("SELECT status FROM service_tickets WHERE id=?", (tk_id,)).fetchone()
    if not t:
        return jsonify({'error':'工单不存在'}),404
    if t['status'] != 'resolved':
        return jsonify({'error':f'工单状态不允许确认（当前: {t["status"]}）'}),400
    # Transition resolved → rated (confirm without rating uses default rating 5)
    ok, msg = transition_ticket(tk_id, 'rated', session.get('customer_id',''), '客户',
                                {'customer_rating': 5, 'customer_feedback': '客户确认已解决'})
    if not ok:
        return jsonify({'error':msg}),400
    db.execute("UPDATE conversations SET status='resolved',updated_at=datetime('now','localtime') WHERE id=(SELECT conversation_id FROM service_tickets WHERE id=?)", (tk_id,))
    conv = db.execute("SELECT conversation_id FROM service_tickets WHERE id=?", (tk_id,)).fetchone()
    if conv:
        msg_text = '✅ 客户已确认问题已解决，等待IT工程师最终关闭。'
        db.execute("INSERT INTO messages(id,conversation_id,role,content) VALUES(?,?,?,?)",
                         (gen_id('msg-'), conv['conversation_id'], 'system', msg_text))
        db.execute("UPDATE conversations SET updated_at=datetime('now','localtime') WHERE id=?", (conv['conversation_id'],))
    db.commit()
    return jsonify({'ok':True})

@app.route('/api/csrf-token', methods=['GET'])
def csrf_token_endpoint():
    return jsonify({'ok': True, 'csrf_token': generate_csrf_token()})

@app.route('/api/customer/tickets/rate', methods=['POST'])
def customer_rate():
    """P0: User rates the service → transitions resolved ticket to rated with rating"""
    data = request.get_json()
    tk_id = data.get('ticket_id','')
    rating = min(5, max(1, int(data.get('rating',5))))
    feedback = data.get('feedback','')
    db = get_db()
    t = db.execute("SELECT status FROM service_tickets WHERE id=?", (tk_id,)).fetchone()
    if not t:
        return jsonify({'error':'工单不存在'}),404
    if t['status'] not in ('resolved', 'rated'):
        return jsonify({'error':f'工单状态不允许评价（当前: {t["status"]}）'}),400
    # Transition resolved → rated (or update rating if already rated)
    if t['status'] == 'resolved':
        ok, msg = transition_ticket(tk_id, 'rated', session.get('customer_id',''), '客户',
                                    {'customer_rating': rating, 'customer_feedback': feedback})
        if not ok:
            return jsonify({'error':msg}),400
    else:
        # Already rated, just update feedback
        db.execute("UPDATE service_tickets SET customer_rating=?,customer_feedback=?,updated_at=datetime('now','localtime') WHERE id=?", (rating, feedback, tk_id))
    db.execute("UPDATE conversations SET status='resolved',updated_at=datetime('now','localtime') WHERE id=(SELECT conversation_id FROM service_tickets WHERE id=?)", (tk_id,))
    conv = db.execute("SELECT conversation_id FROM service_tickets WHERE id=?", (tk_id,)).fetchone()
    if conv:
        stars = '⭐' * rating + '☆' * (5 - rating)
        fb = (' 反馈：' + feedback) if feedback else ''
        detail = f'{stars}/5' + fb
        msg_text = '✅ 客户已评价，' + detail + '，等待IT工程师最终关闭。'
        db.execute("INSERT INTO messages(id,conversation_id,role,content) VALUES(?,?,?,?)",
                         (gen_id('msg-'), conv['conversation_id'], 'system', msg_text))
        db.execute("UPDATE conversations SET updated_at=datetime('now','localtime') WHERE id=?", (conv['conversation_id'],))
    db.commit()
    log_audit(tk_id, 'ticket.rated', session.get('customer_id',''), '客户',
              {'action': 'rate', 'rating': rating, 'feedback': feedback[:100] if feedback else ''})
    return jsonify({'ok':True})


@app.route('/api/agent/tickets/final-close', methods=['POST'])
@logged_in_required
def agent_final_close():
    """P0: Unified close - close from rated status (or force-close from any active status)"""
    data = request.get_json()
    tk_id = data.get('ticket_id','')
    notes = data.get('resolution_notes','').strip()
    db = get_db()
    t = db.execute("SELECT status, conversation_id FROM service_tickets WHERE id=?", (tk_id,)).fetchone()
    if not t:
        return jsonify({'error':'工单不存在'}),404
    allowed = ['rated', 'resolved', 'processing']
    if t['status'] not in allowed:
        return jsonify({'error':f'工单状态不允许关闭（当前: {t["status"]}）'}),400
    ok, msg = transition_ticket(tk_id, 'closed', session.get('agent_id',''), session.get('agent_name',''),
                                {'resolution_notes': notes, 'close_reason': '工程师关闭工单'})
    if not ok:
        return jsonify({'error':msg}),400
    conv = db.execute("SELECT conversation_id FROM service_tickets WHERE id=?", (tk_id,)).fetchone()
    if conv:
        db.execute("UPDATE conversations SET status='resolved',updated_at=datetime('now','localtime') WHERE id=?", (conv['conversation_id'],))
        db.execute("INSERT INTO messages(id,conversation_id,role,content) VALUES(?,?,?,?)",
                   (gen_id('msg-'), conv['conversation_id'], 'system',
                    '✅ 工程师已确认关闭工单。'))
    db.commit()
    log_audit(tk_id, 'ticket.closed', session.get('agent_id',''), session.get('agent_name',''),
              {'from':t['status'],'to':'closed','notes':notes})
    return jsonify({"ok":True})

@app.route('/api/agent/tickets/<tk_id>/reopen', methods=['POST'])
@logged_in_required
def agent_reopen_ticket(tk_id):
    """L3/Admin 可重新打开已关闭工单（从 closed 回到 processing）"""
    ok, msg = transition_ticket(tk_id, 'processing', session['agent_id'],
                                session.get('agent_name',''),
                                {'reason':'reopened'})
    if not ok:
        return jsonify({'error':msg}), 400
    t = get_db().execute("SELECT conversation_id FROM service_tickets WHERE id=?", (tk_id,)).fetchone()
    if t:
        get_db().execute("INSERT INTO messages(id,conversation_id,role,content) VALUES(?,?,?,?)",
                    (gen_id('msg-'), t['conversation_id'], 'system', '🔁 工单已重新打开'))
        get_db().execute("UPDATE conversations SET status='active',updated_at=datetime('now','localtime') WHERE id=?", (t['conversation_id'],))
        get_db().commit()
    return jsonify({'ok':True, 'msg':msg})

@app.route('/api/customer/tickets/selfclose', methods=['POST'])
def customer_selfclose():
    """P0: Customer self-closes by rating. Transitions to rated (not closed)."""
    cid = get_or_create_customer()
    if not cid:
        return jsonify({'error':'未登录','login_required':True}),401
    data = request.get_json()
    tk_id = data.get('ticket_id','')
    rating = min(5, max(1, int(data.get('rating',5))))
    feedback = data.get('feedback','')
    db = get_db()
    t = db.execute("SELECT id, status FROM service_tickets WHERE id=? AND customer_id=?", (tk_id, cid)).fetchone()
    if not t: return jsonify({'error':'工单不存在'}),404
    # Transition resolved → rated with rating
    ok, _ = transition_ticket(tk_id, 'rated', cid, '客户',
                              {'customer_rating': rating, 'customer_feedback': feedback,
                               'close_reason': '客户评价后自动关闭'})
    db.execute("UPDATE conversations SET status='resolved',updated_at=datetime('now','localtime') WHERE id=(SELECT conversation_id FROM service_tickets WHERE id=?)", (tk_id,))
    conv = db.execute("SELECT conversation_id FROM service_tickets WHERE id=?", (tk_id,)).fetchone()
    if conv:
        stars = '⭐' * rating + '☆' * (5 - rating)
        fb = (' 📝 反馈：' + feedback) if feedback else ''
        detail = f'{stars}/5' + fb
        msg_text = f'✅ 工单已评价 {detail}，等待IT工程师最终关闭。'
        db.execute("INSERT INTO messages(id,conversation_id,role,content) VALUES(?,?,?,?)",
                         (gen_id('msg-'), conv['conversation_id'], 'system', msg_text))
        db.execute("UPDATE conversations SET updated_at=datetime('now','localtime') WHERE id=?", (conv['conversation_id'],))
    db.commit()
    log_audit(tk_id, 'ticket.rated', cid, '客户',
              {'rating':rating, 'feedback':feedback})
    return jsonify({'ok':True})

@app.route('/api/conversation/auto-close', methods=['POST'])
def conversation_auto_close():
    """自动创建自助处理工单（20分钟未转人工后的空闲时触发）"""
    data = request.get_json()
    conv_id = data.get('conversation_id', '')
    if not conv_id:
        return jsonify({'error': '缺少会话ID'}), 400
    
    cid = get_or_create_customer()
    if not cid:
        return jsonify({'error':'未登录','login_required':True}),401
    db = get_db()
    
    # 检查是否已有工单
    existing = db.execute("SELECT id FROM service_tickets WHERE conversation_id=? AND status NOT IN ('closed')", (conv_id,)).fetchone()
    if existing:
        transition_ticket(existing['id'], 'closed', 'system', '系统',
                        {'close_reason': '自助处理'})
        db.commit()
        return jsonify({'ok': True, 'ticket_id': existing['id']})
    
    # 获取对话摘要
    msgs = db.execute("SELECT role,content FROM messages WHERE conversation_id=? ORDER BY created_at ASC", (conv_id,)).fetchall()
    
    # 构建摘要
    summary_parts = []
    for m in msgs:
        role_label = {'user':'用户','bot':'AI','agent':'IT工程师'}.get(m['role'], m['role'])
        content = m['content'][:150]
        summary_parts.append(f"[{role_label}] {content}")
    
    conversation_text = '\n'.join(summary_parts[-20:])  # 最近20条消息
    
    # 用 AI 生成简洁摘要
    summary = ""
    try:
        summary_reply, _ = call_llm([
            {'role':'system','content':'你是一个工单摘要生成助手。请根据以下对话内容，用一句话生成工单描述。格式："用户在【系统/功能名称】遇到了【问题描述】，已通过知识库自助解决。"如果无法判断，直接描述"通过知识库自助解决了相关咨询。"'},
            {'role':'user','content':f'对话内容：\n{conversation_text}'}
        ], timeout=15)
        if summary_reply:
            summary = summary_reply.strip()
    except:
        pass
    
    if not summary:
        summary = "通过知识库自助解决了相关咨询。"
    
    # 创建工单
    esc_id = gen_id('esc-')
    tk_id = gen_id('tk-')
    tk_number = generate_ticket_number()
    
    db.execute("INSERT INTO escalations(id,conversation_id,status,reason) VALUES(?,?,'closed','AI自助解决')",
               (esc_id, conv_id))
    db.execute("""INSERT INTO service_tickets(id,escalation_id,conversation_id,customer_id,
                 issue_description,status,close_reason,closed_at,created_at,updated_at,ticket_number)
                 VALUES(?,?,?,?,?,'closed','AI自助解决',datetime('now','localtime'),datetime('now','localtime'),datetime('now','localtime'),?)""",
               (tk_id, esc_id, conv_id, cid, summary, tk_number))
    db.execute("INSERT INTO messages(id,conversation_id,role,content) VALUES(?,?,?,?)",
               (gen_id('msg-'), conv_id, 'system',
                '✅ 已为您自动创建服务记录（工单编号: ' + tk_number + '），问题已通过知识库自助解决。如需进一步帮助，请重新咨询。'))
    db.commit()
    return jsonify({'ok': True, 'ticket_id': tk_id, 'ticket_number': tk_number})

@app.route('/api/customer/agent', methods=['GET'])
def customer_agent_info():
    cid = get_or_create_customer()
    if not cid:
        return jsonify({'bound':False,'error':'未登录','login_required':True}),401
    info = get_db().execute("""
        SELECT b.*, ap.display_name, ap.department, ap.title, ag.name as agent_name
        FROM customer_agent_bindings b JOIN agents ag ON b.agent_id = ag.id
        LEFT JOIN agent_profiles ap ON b.agent_id = ap.agent_id
        WHERE b.customer_id=? AND b.status='active'""", (cid,)).fetchone()
    return jsonify(dict(info) if info else {'bound':False})

# ====== 工程师工作台 ======
@app.route('/agent/login', methods=['GET','POST'])
def agent_login():
    if request.method == 'GET':
        enabled_providers = get_enabled_providers()
        return render_template('agent_login.html', providers=enabled_providers)
    data = request.get_json() or request.form
    
    # Try werkzeug password hash first (for 三员 accounts), then sha256 (legacy)
    a = get_db().execute("SELECT id,name,role,password_hash FROM agents WHERE email=?", 
                         (data.get('email',''),)).fetchone()
    if not a:
        _record_login('', data.get('email',''), 'agent', False, '用户名或密码错误')
        return jsonify({'error':'用户名或密码错误'}),401
    
    # Check password with both methods
    password_ok = False
    if a['password_hash'].startswith('pbkdf2:') or a['password_hash'].startswith('scrypt:'):
        password_ok = check_password_hash(a['password_hash'], data.get('password',''))
    else:
        pwd_hash = hashlib.sha256(f'admin:{data.get("password","")}'.encode()).hexdigest()
        password_ok = (pwd_hash == a['password_hash'])
    
    if not password_ok:
        _record_login('', data.get('email',''), 'agent', False, '用户名或密码错误')
        return jsonify({'error':'用户名或密码错误'}),401
    
    session['agent_id'] = a['id']; session['agent_name'] = a['name']; session['agent_role'] = a['role']
    get_db().execute("UPDATE agents SET status='online' WHERE id=?", (a['id'],))
    get_db().commit()
    _record_login(a['id'], a['name'], 'agent', True)
    log_audit('', 'agent.login', a['id'], a['name'],
              {'role': a['role'], 'method': 'password'})
    
    # Redirect based on role
    admin_roles = {'sysadmin', 'secadmin', 'audadmin', 'superadmin'}
    redirect_url = '/admin/dashboard' if a['role'] in admin_roles else '/agent/dashboard'
    return jsonify({'ok':True,'redirect':redirect_url,'role':a['role']})

@app.route('/agent/logout', methods=['POST'])
def agent_logout():
    aid = session.pop('agent_id', None)
    name = session.get('agent_name', 'IT工程师')
    if aid:
        log_audit('', 'agent.logout', aid, name, {})
        get_db().execute("UPDATE agents SET status='offline' WHERE id=?", (aid,)); get_db().commit()
    return jsonify({'ok':True})

@app.route('/agent/dashboard')
@agent_required
def agent_dashboard():
    prof = get_db().execute("SELECT agent_level FROM agent_profiles WHERE agent_id=?", (session['agent_id'],)).fetchone()
    cfg = get_db().execute("SELECT value FROM system_config WHERE key='level_names'").fetchone()
    level_names = json.loads(cfg['value']) if cfg and cfg['value'] else {'1':'初级工程师','2':'高级工程师','3':'专家工程师','4':'首席工程师'}
    return render_template('agent_dashboard.html',
        agent_name=session.get('agent_name',''),
        agent_role=session.get('agent_role',''),
        agent_level=prof['agent_level'] if prof else 1,
        level_names=json.dumps(level_names, ensure_ascii=False))

@app.route('/api/agent/profile', methods=['GET','POST'])
@agent_required
def agent_profile():
    aid = session['agent_id']; db = get_db()
    if request.method == 'POST':
        data = request.get_json()
        prof = db.execute("SELECT id FROM agent_profiles WHERE agent_id=?", (aid,)).fetchone()
        if prof:
            db.execute("UPDATE agent_profiles SET display_name=?,department=?,title=?,phone=?,employee_id=?,company=? WHERE agent_id=?", 
                       (data.get('display_name',''), data.get('department',''), data.get('title',''), data.get('phone',''), data.get('employee_id',''), data.get('company',''), aid))
        else:
            db.execute("INSERT INTO agent_profiles(id,agent_id,display_name,department,title,phone,employee_id,company) VALUES(?,?,?,?,?,?,?,?)",
                       (gen_id('ap-'), aid, data.get('display_name',''), data.get('department',''), data.get('title',''), data.get('phone',''), data.get('employee_id',''), data.get('company','')))
        if data.get('display_name'): db.execute("UPDATE agents SET name=? WHERE id=?", (data['display_name'], aid))
        db.commit(); return jsonify({'ok':True})
    prof = db.execute("SELECT * FROM agent_profiles WHERE agent_id=?", (aid,)).fetchone()
    agent = db.execute("SELECT id,name,email,status,role FROM agents WHERE id=?", (aid,)).fetchone()
    r = dict(agent); 
    if prof: r.update(dict(prof))
    systems = get_db().execute("""
        SELECT s.name FROM systems s JOIN agent_systems a_s ON s.id = a_s.system_id
        WHERE a_s.agent_id=? ORDER BY s.name""", (aid,)).fetchall()
    r['systems'] = [s['name'] for s in systems]
    return jsonify(r)

@app.route('/api/agent/tickets', methods=['GET'])
@agent_required
def agent_tickets():
    """P0: Agent ticket listing - filters out 'created' status tickets. All levels have equal permissions (L1~L4)."""
    status_filter = request.args.get('status','pending')
    db = get_db()
    aid = session['agent_id']
    
    if status_filter == 'search':
        # Search tab with date range and optional filters
        cfg = {r['key']: r['value'] for r in db.execute('SELECT key, value FROM system_config').fetchall()}
        max_days = int(cfg.get('ticket_search_max_days', '365'))
        start_date = request.args.get('start_date', '')
        end_date = request.args.get('end_date', '')
        search_status = request.args.get('search_status', '')
        ticket_number = request.args.get('ticket_number', '')
        customer_name = request.args.get('customer_name', '')
        
        where = ['1=1']
        params = []
        if start_date:
            where.append("t.created_at >= ?")
            params.append(start_date + ' 00:00:00')
        if end_date:
            where.append("t.created_at <= ?")
            params.append(end_date + ' 23:59:59')
        if search_status:
            where.append("t.status=?")
            params.append(search_status)
        if ticket_number:
            where.append("t.ticket_number LIKE ?")
            params.append(f'%{ticket_number}%')
        if customer_name:
            where.append("cust.name LIKE ?")
            params.append(f'%{customer_name}%')
        
        # Enforce max_days limit if no date specified
        if not start_date and not end_date:
            cutoff = (datetime.now() - timedelta(days=max_days)).strftime('%Y-%m-%d 00:00:00')
            where.append("t.created_at >= ?")
            params.append(cutoff)
        
        where_sql = ' AND '.join(where)
        tickets = db.execute(f"""
            SELECT t.*, cust.name as customer_name,
                   (SELECT content FROM messages WHERE conversation_id=t.conversation_id AND role='user' ORDER BY created_at DESC LIMIT 1) as last_msg, 0 as overdue
            FROM service_tickets t JOIN customers cust ON t.customer_id = cust.id
            WHERE {where_sql} ORDER BY t.created_at DESC LIMIT 100""", params).fetchall()
    elif status_filter == 'all_mine':
        # Tickets where agent_id=current_agent OR agent has replied in conversation
        tickets = db.execute("""
            SELECT DISTINCT t.*, cust.name as customer_name,
                   (SELECT content FROM messages WHERE conversation_id=t.conversation_id AND role='user' ORDER BY created_at DESC LIMIT 1) as last_msg, 0 as overdue
            FROM service_tickets t JOIN customers cust ON t.customer_id = cust.id
            LEFT JOIN messages m ON m.conversation_id = t.conversation_id AND m.role='agent' AND m.content LIKE ?
            WHERE t.agent_id=? OR (m.id IS NOT NULL)
            ORDER BY t.created_at DESC LIMIT 100""", (f'%{session.get("agent_name","")}%', aid)).fetchall()
    elif status_filter == 'mine':
        tickets = db.execute("""
            SELECT t.*, cust.name as customer_name,
                   (SELECT content FROM messages WHERE conversation_id=t.conversation_id AND role='user' ORDER BY created_at DESC LIMIT 1) as last_msg, 0 as overdue
            FROM service_tickets t JOIN customers cust ON t.customer_id = cust.id
            WHERE t.agent_id=? AND t.status NOT IN ('closed','resolved','rated')
            ORDER BY t.created_at DESC""", (aid,)).fetchall()
    elif status_filter == 'history':
        tickets = db.execute("""
            SELECT t.*, cust.name as customer_name
            FROM service_tickets t JOIN customers cust ON t.customer_id = cust.id
            WHERE t.agent_id=? AND t.status IN ('closed','rated')
            ORDER BY t.created_at DESC LIMIT 50""", (aid,)).fetchall()
    else:
        # All levels (L1~L4) have equal permissions: see processing + own tickets
        tickets = db.execute("""
            SELECT t.*, cust.name as customer_name,
                   (SELECT content FROM messages WHERE conversation_id=t.conversation_id AND role='user' ORDER BY created_at DESC LIMIT 1) as last_msg, 0 as overdue
            FROM service_tickets t JOIN customers cust ON t.customer_id = cust.id
            WHERE t.status='processing'
               OR (t.agent_id=? AND t.status IN ('processing','resolved','rated'))
            ORDER BY t.created_at""", (aid,)).fetchall()
    
    # Add pending count as separate response field when status is pending
    if status_filter == 'pending':
        pending_count = db.execute("SELECT COUNT(*) FROM service_tickets WHERE status='processing' AND (agent_id IS NULL OR agent_id='' OR first_read_at IS NULL OR first_read_at='')").fetchone()[0]
        result = [dict(t) for t in tickets]
        return jsonify(result)
    
    return jsonify([dict(t) for t in tickets])

@app.route('/api/agent/tickets/assign', methods=['POST'])
@agent_required
def agent_assign():
    """P0: Assign a processing ticket to self"""
    data = request.get_json()
    tk_id = data.get('ticket_id',''); aid = session['agent_id']; db = get_db()
    agent = db.execute("SELECT status FROM agents WHERE id=?", (aid,)).fetchone()
    if agent and agent['status'] == 'offline':
        db.execute("UPDATE agents SET status='online' WHERE id=?", (aid,))
    db.execute("UPDATE service_tickets SET agent_id=?,assigned_at=datetime('now','localtime'),first_read_at=datetime('now','localtime'),updated_at=datetime('now','localtime') WHERE id=? AND status='processing'", (aid, tk_id))
    db.execute("UPDATE service_tickets SET first_read_at=datetime('now','localtime'),updated_at=datetime('now','localtime') WHERE id=? AND agent_id=? AND (first_read_at IS NULL OR first_read_at='')", (tk_id, aid))
    t = db.execute("SELECT customer_id FROM service_tickets WHERE id=?", (tk_id,)).fetchone()
    if t:
        existing = db.execute("SELECT id FROM customer_agent_bindings WHERE customer_id=? AND agent_id=?", (t['customer_id'], aid)).fetchone()
        if not existing: db.execute("INSERT INTO customer_agent_bindings(id,customer_id,agent_id) VALUES(?,?,?)", (gen_id('cab-'), t['customer_id'], aid))
    db.commit()
    log_audit(tk_id, 'ticket.assigned', aid, session.get('agent_name',''), {'from':'processing','to':'processing'})
    return jsonify({'ok':True})

@app.route('/api/agent/tickets/close', methods=['POST'])
@agent_required
def agent_close():
    """P0: Close a rated ticket (rated → closed). Also supports force-close of any active status."""
    data = request.get_json()
    tk_id = data.get('ticket_id',''); notes = data.get('resolution_notes','')
    close_reason = data.get('close_reason','')
    db = get_db()
    t = db.execute("SELECT status, conversation_id FROM service_tickets WHERE id=?", (tk_id,)).fetchone()
    if not t:
        return jsonify({'error':'工单不存在'}),404
    status = t['status']
    conv_id = t['conversation_id']
    # Save current agent's level to ticket when closing
    prof = get_db().execute("SELECT agent_level FROM agent_profiles WHERE agent_id=?", (session['agent_id'],)).fetchone()
    closing_level = prof['agent_level'] if prof else 1
    db.execute("UPDATE service_tickets SET level=? WHERE id=?", (closing_level, tk_id))
    if status == 'rated':
        ok, msg = transition_ticket(tk_id, 'closed', session.get('agent_id',''), session.get('agent_name',''),
                                    {'resolution_notes': notes, 'close_reason': close_reason or '工程师关闭工单'})
        if not ok:
            return jsonify({'error':msg}),400
    elif status in ('created', 'processing', 'resolved'):
        # Force close (admin only check done at decorator level, but this is agent route)
        ok, msg = transition_ticket(tk_id, 'closed', session.get('agent_id',''), session.get('agent_name',''),
                                    {'resolution_notes': notes, 'close_reason': close_reason or f'强制关闭({status})'})
        if not ok:
            return jsonify({'error':msg}),400
    else:
        return jsonify({'error':f'无法关闭当前状态({status})的工单'}),400
    if conv_id:
        db.execute("UPDATE conversations SET status='resolved' WHERE id=?", (conv_id,))
    db.commit()
    return jsonify({'ok':True})

@app.route('/api/agent/tickets/resolve', methods=['POST'])
@agent_required
def agent_resolve():
    """P0: Agent resolves a processing ticket (processing → resolved)"""
    data = request.get_json()
    tk_id = data.get('ticket_id','')
    notes = data.get('resolution_notes','').strip()
    db = get_db()
    t = db.execute("SELECT status, conversation_id FROM service_tickets WHERE id=? AND agent_id=?", (tk_id, session.get('agent_id',''))).fetchone()
    if not t:
        return jsonify({'error':'工单不存在或无权操作'}),403
    if t['status'] != 'processing':
        return jsonify({'error':f'当前状态不允许解决（{t["status"]}）'}),400
    conv_id = t['conversation_id']
    ok, msg = transition_ticket(tk_id, 'resolved', session.get('agent_id',''), session.get('agent_name',''),
                                {'resolution_notes': notes, 'close_reason': '工程师解决'})
    if not ok:
        return jsonify({'error':msg}),400
    db.execute("INSERT INTO messages(id,conversation_id,role,content) VALUES(?,?,?,?)",
               (gen_id('msg-'), conv_id, 'system',
                '【工程师已解决工单】请问还有其他问题吗？如果没有，请确认并评价。'))
    db.execute("UPDATE conversations SET updated_at=datetime('now','localtime') WHERE id=?", (conv_id,))
    db.commit()
    return jsonify({'ok':True})

@app.route('/api/agent/tickets/<tk_id>', methods=['GET'])
@agent_required
def agent_ticket_detail(tk_id):
    t = get_db().execute("""
        SELECT t.*, cust.name as customer_name, cust.id as customer_id,
               cp.phone, cp.company
        FROM service_tickets t JOIN customers cust ON t.customer_id = cust.id
        LEFT JOIN customer_profiles cp ON t.customer_id = cp.customer_id
        WHERE t.id=?""", (tk_id,)).fetchone()
    r = dict(t) if t else {}
    if t and t['agent_id']:
        agent_systems = get_db().execute("""
            SELECT s.name FROM systems s JOIN agent_systems a_s ON s.id = a_s.system_id
            WHERE a_s.agent_id=? ORDER BY s.name""", (t['agent_id'],)).fetchall()
        r['systems'] = [s['name'] for s in agent_systems]
    return jsonify(r)

@app.route('/api/agent/customers', methods=['GET'])
@agent_required
def agent_customers():
    aid = session['agent_id']
    cs = get_db().execute("""
        SELECT b.*, cust.name as customer_name, cp.phone, cp.company,
               (SELECT COUNT(*) FROM service_tickets WHERE customer_id=cust.id AND agent_id=?) as ticket_count,
               (SELECT MAX(created_at) FROM service_tickets WHERE customer_id=cust.id AND agent_id=?) as last_service
        FROM customer_agent_bindings b JOIN customers cust ON b.customer_id = cust.id
        LEFT JOIN customer_profiles cp ON b.customer_id = cp.customer_id
        WHERE b.agent_id=? AND b.status='active' ORDER BY last_service DESC""", (aid, aid, aid)).fetchall()
    return jsonify([dict(c) for c in cs])

@app.route('/api/agent/conversation/<conv_id>', methods=['GET'])
@agent_required
def agent_conversation(conv_id):
    msgs = get_db().execute("SELECT role,content,image_url,created_at FROM messages WHERE conversation_id=? ORDER BY created_at", (conv_id,)).fetchall()
    return jsonify([dict(m) for m in msgs])

@app.route('/api/agent/reply', methods=['POST'])
@agent_required
def agent_reply():
    data = request.get_json()
    content = data.get('content','').strip()
    if not content or not data.get('conversation_id'): return jsonify({'error':'参数错误'}),400
    get_db().execute("INSERT INTO messages(id,conversation_id,role,content) VALUES(?,?,?,?)",
                     (gen_id('msg-'), data['conversation_id'], 'agent', f"[{session['agent_name']}]: {content}"))
    get_db().commit()

    log_audit('', 'agent.reply', session.get('agent_id',''), session.get('agent_name',''),
              {'conversation_id': data['conversation_id'], 'msg_preview': content[:80]})
    return jsonify({'ok':True})

# ====== 管理员后台 ======
@app.route('/admin/dashboard')
@logged_in_required
def admin_dashboard():
    agent_role = session.get('agent_role', '')
    is_admin = agent_role in ('sysadmin', 'secadmin', 'audadmin', 'superadmin')
    if is_admin:
        prof = get_db().execute("SELECT agent_level FROM agent_profiles WHERE agent_id=?", (session['agent_id'],)).fetchone()
        return render_template('admin.html', 
            agent_name=session.get('agent_name',''),
            agent_role=agent_role,
            agent_level=prof['agent_level'] if prof else 1)
    prof = get_db().execute("SELECT agent_level FROM agent_profiles WHERE agent_id=?", (session['agent_id'],)).fetchone()
    return render_template('agent_dashboard.html',
        agent_name=session.get('agent_name',''),
        agent_role=agent_role,
        agent_level=prof['agent_level'] if prof else 1)

# 用户管理
@app.route('/api/admin/customers/<cid>/profile', methods=['PUT'])
@admin_required
def admin_update_customer_profile(cid):
    data = request.get_json()
    db = get_db()
    existing = db.execute("SELECT id FROM customer_profiles WHERE customer_id=?", (cid,)).fetchone()
    if existing:
        db.execute("UPDATE customer_profiles SET name=?,phone=?,company=?,notes=?,employee_id=?,department=?,updated_at=datetime('now','localtime') WHERE customer_id=?",
                   (data.get('name',''), data.get('phone',''), data.get('company',''),
                    data.get('notes',''), data.get('employee_id',''), data.get('department',''), cid))
    else:
        db.execute("INSERT INTO customer_profiles(id,customer_id,name,phone,company,notes,employee_id,department) VALUES(?,?,?,?,?,?,?,?)",
                   (gen_id('cp-'), cid, data.get('name',''), data.get('phone',''),
                    data.get('company',''), data.get('notes',''), data.get('employee_id',''), data.get('department','')))
    if data.get('email') is not None:
        db.execute("UPDATE customers SET email=? WHERE id=?", (data.get('email',''), cid))
    if data.get('name'):
        db.execute("UPDATE customers SET name=? WHERE id=?", (data.get('name',''), cid))
    db.commit()
    return jsonify({'ok': True})

@app.route('/api/admin/customers', methods=['GET'])
@admin_required
def admin_customers():
    page = int(request.args.get('page',1))
    limit = int(request.args.get('limit',20))
    offset = (page-1)*limit
    search = request.args.get('search','')
    db = get_db()
    if search:
        total = db.execute("SELECT COUNT(*) FROM customers WHERE name LIKE ? OR id LIKE ?", (f'%{search}%', f'%{search}%')).fetchone()[0]
        customers = db.execute("""
            SELECT cust.*, cp.name as pname, cp.phone, cp.company, cp.notes,
                   (SELECT COUNT(*) FROM service_tickets WHERE customer_id=cust.id) as ticket_count
            FROM customers cust LEFT JOIN customer_profiles cp ON cust.id = cp.customer_id
            WHERE cust.name LIKE ? OR cust.id LIKE ?
            ORDER BY cust.created_at DESC LIMIT ? OFFSET ?""", (f'%{search}%', f'%{search}%', limit, offset)).fetchall()
    else:
        total = db.execute("SELECT COUNT(*) FROM customers").fetchone()[0]
        customers = db.execute("""
            SELECT cust.*, cp.name as pname, cp.phone, cp.company, cp.notes,
                   (SELECT COUNT(*) FROM service_tickets WHERE customer_id=cust.id) as ticket_count
            FROM customers cust LEFT JOIN customer_profiles cp ON cust.id = cp.customer_id
            ORDER BY cust.created_at DESC LIMIT ? OFFSET ?""", (limit, offset)).fetchall()
    return jsonify({'customers':[dict(c) for c in customers],'total':total,'page':page,'pages':(total+limit-1)//limit})

@app.route('/api/admin/customers/<cid>', methods=['GET'])
@admin_required
def admin_customer_detail(cid):
    cust = get_db().execute("""
        SELECT cust.*, cp.name as pname, cp.phone, cp.company, cp.notes,
               cp.created_at as profile_created, cp.updated_at as profile_updated
        FROM customers cust LEFT JOIN customer_profiles cp ON cust.id = cp.customer_id
        WHERE cust.id=?""", (cid,)).fetchone()
    if not cust: return jsonify({'error':'客户不存在'}),404
    tickets = get_db().execute("""
        SELECT t.*, ag.name as agent_name, ap.display_name as agent_display
        FROM service_tickets t LEFT JOIN agents ag ON t.agent_id = ag.id
        LEFT JOIN agent_profiles ap ON t.agent_id = ap.agent_id
        WHERE t.customer_id=? ORDER BY t.created_at DESC""", (cid,)).fetchall()
    bindings = get_db().execute("""
        SELECT b.*, ag.name as agent_name, ap.display_name as agent_display, ap.department
        FROM customer_agent_bindings b JOIN agents ag ON b.agent_id = ag.id
        LEFT JOIN agent_profiles ap ON b.agent_id = ap.agent_id
        WHERE b.customer_id=?""", (cid,)).fetchall()
    return jsonify({'customer':dict(cust),'tickets':[dict(t) for t in tickets],'bindings':[dict(b) for b in bindings]})

@app.route('/api/admin/customers/<cid>/reset-password', methods=['PUT'])
@admin_required
def admin_customer_reset_password(cid):
    '''管理员重置用户密码'''
    data = request.get_json()
    password = data.get('password', '')
    if not password or len(password) < 6:
        return jsonify({'error': '密码不能少于6位'}), 400
    db = get_db()
    customer = db.execute("SELECT id FROM customers WHERE id=?", (cid,)).fetchone()
    if not customer:
        return jsonify({'error': '客户不存在'}), 404
    pw_hash = generate_password_hash(password)
    db.execute("UPDATE customers SET password_hash=? WHERE id=?", (pw_hash, cid))
    db.commit()
    log_audit('', 'admin.customer.reset_password', session.get('agent_id',''), session.get('agent_name',''),
              {'customer_id': cid})
    return jsonify({'ok': True, 'message': '密码已重置'})

@app.route('/api/admin/customers/<cid>', methods=['DELETE'])
@csrf_protect
@admin_required
def admin_delete_customer(cid):
    '''管理员删除用户，同时清理关联数据'''
    db = get_db()
    customer = db.execute("SELECT id FROM customers WHERE id=?", (cid,)).fetchone()
    if not customer:
        return jsonify({'error': '客户不存在'}), 404
    # Clean up associated data
    db.execute("DELETE FROM customer_profiles WHERE customer_id=?", (cid,))
    db.execute("DELETE FROM customer_agent_bindings WHERE customer_id=?", (cid,))
    db.execute("DELETE FROM im_user_mappings WHERE smartcs_customer_id=?", (cid,))
    # Delete tickets and related messages
    conv_ids = [row['conversation_id'] for row in db.execute("SELECT conversation_id FROM service_tickets WHERE customer_id=?", (cid,)).fetchall()]
    for conv_id in conv_ids:
        db.execute("DELETE FROM messages WHERE conversation_id=?", (conv_id,))
        db.execute("DELETE FROM escalations WHERE conversation_id=?", (conv_id,))
        db.execute("DELETE FROM conversations WHERE id=?", (conv_id,))
    db.execute("DELETE FROM service_tickets WHERE customer_id=?", (cid,))
    db.execute("DELETE FROM customers WHERE id=?", (cid,))
    db.commit()
    log_audit('', 'admin.customer.delete', session.get('agent_id',''), session.get('agent_name',''),
              {'customer_id': cid})
    return jsonify({'ok': True, 'message': '客户已删除'})

# 工程师管理

@app.route('/api/admin/audit-logs', methods=['GET'])
@audadmin_required
def admin_audit_logs():
    db = get_db()
    page = request.args.get('page', 1, type=int)
    per_page = request.args.get('per_page', int(get_system_config().get('pagination_per_page', '50')), type=int)
    actor_id = request.args.get('actor_id', '')
    action = request.args.get('action', '')
    start_date = request.args.get('start_date', '')
    end_date = request.args.get('end_date', '')
    where = []
    params = []
    if actor_id:
        where.append('actor_id=?')
        params.append(actor_id)
    if action:
        where.append('action LIKE ?')
        params.append(f'%{action}%')
    if start_date:
        where.append('created_at>=?')
        params.append(start_date)
    if end_date:
        where.append('created_at<=?')
        params.append(end_date + ' 23:59:59')
    where_sql = ' WHERE ' + ' AND '.join(where) if where else ''
    total = db.execute("SELECT COUNT(*) FROM audit_log" + where_sql, params).fetchone()[0]
    offset = (page - 1) * per_page
    logs = db.execute("SELECT * FROM audit_log" + where_sql + " ORDER BY created_at DESC LIMIT ? OFFSET ?",
                      params + [per_page, offset]).fetchall()
    return jsonify({'logs': [dict(l) for l in logs], 'total': total, 'page': page, 'per_page': per_page})

@app.route('/api/admin/login-logs', methods=['GET'])
@audadmin_required
def admin_login_logs():
    db = get_db()
    page = request.args.get('page', 1, type=int)
    per_page = request.args.get('per_page', 50, type=int)
    user_type = request.args.get('user_type', '')
    success = request.args.get('success', '')
    total = db.execute("SELECT COUNT(*) FROM login_log").fetchone()[0]
    where = []
    params = []
    if user_type:
        where.append('user_type=?')
        params.append(user_type)
    if success in ('0', '1'):
        where.append('success=?')
        params.append(int(success))
    where_sql = ' WHERE ' + ' AND '.join(where) if where else ''
    offset = (page - 1) * per_page
    logs = db.execute("SELECT * FROM login_log" + where_sql + " ORDER BY login_at DESC LIMIT ? OFFSET ?",
                      params + [per_page, offset]).fetchall()
    return jsonify({'logs': [dict(l) for l in logs], 'total': total, 'page': page, 'per_page': per_page})


@app.route('/api/admin/agents', methods=['GET'])
@role_required('sysadmin', 'secadmin', 'superadmin')
def admin_agents():
    agents = get_db().execute("""
        SELECT ag.*, ap.display_name, ap.department, ap.title, ap.phone as agent_phone, ap.agent_level,
               ap.employee_id, ap.company,
               (SELECT COUNT(*) FROM service_tickets WHERE agent_id=ag.id AND status NOT IN ('closed','resolved','rated')) as active_tickets,
               (SELECT COUNT(*) FROM customer_agent_bindings WHERE agent_id=ag.id AND status='active') as bound_customers
        FROM agents ag LEFT JOIN agent_profiles ap ON ag.id = ap.agent_id
        ORDER BY ag.created_at""").fetchall()
    result = []
    for a in agents:
        a_dict = dict(a)
        systems = get_db().execute("""
            SELECT s.name FROM systems s JOIN agent_systems a_s ON s.id = a_s.system_id
            WHERE a_s.agent_id=? ORDER BY s.name""", (a['id'],)).fetchall()
        a_dict['systems'] = [s['name'] for s in systems]
        result.append(a_dict)
    return jsonify(result)

@app.route('/api/admin/agents', methods=['POST'])
@admin_required
def admin_add_agent():
    data = request.get_json()
    name = data.get('name',''); email = data.get('email',''); password = data.get('password','')
    if not name or not email or not password: return jsonify({'error':'请填写完整信息'}),400
    existing = get_db().execute("SELECT id FROM agents WHERE email=?", (email,)).fetchone()
    if existing: return jsonify({'error':'邮箱已存在'}),400
    # Password policy check
    ok, msg = _check_password_policy(password)
    if not ok:
        return jsonify({'error': msg}), 400
    aid = gen_id('agent-'); pwd_hash = hashlib.sha256(f'admin:{password}'.encode()).hexdigest()
    level = int(data.get('level', 1))
    get_db().execute("INSERT INTO agents(id,name,email,password_hash) VALUES(?,?,?,?)", (aid, name, email, pwd_hash))
    get_db().execute("INSERT INTO agent_profiles(id,agent_id,display_name,department,title,phone,employee_id,company,agent_level) VALUES(?,?,?,?,?,?,?,?,?)",
                     (gen_id('ap-'), aid, name, data.get('department',''), data.get('title',''),
                      data.get('phone',''), data.get('employee_id',''), data.get('company',''), level))
    get_db().commit()
    log_audit('', 'admin.agent.create', session.get('agent_id',''), session.get('agent_name',''),
              {'agent_id': aid, 'name': name, 'email': email, 'level': level})
    return jsonify({'ok':True})

@app.route('/api/admin/agents/<aid>', methods=['DELETE'])
@csrf_protect
def admin_delete_agent(aid):
    # 三员管理：删除工程师归安全管理员
    if not session.get('agent_id'):
        return jsonify({'error':'未登录'}),401
    role = session.get('agent_role', '')
    if role not in ('secadmin', 'superadmin'):
        return jsonify({'error':'权限不足'}),403
    if aid == session['agent_id']: return jsonify({'error':'不能删除自己'}),400
    # Delete protection check
    ok, msg = _check_delete_agent_protection(aid)
    if not ok:
        return jsonify({'error': msg}), 400
    get_db().execute("DELETE FROM customer_agent_bindings WHERE agent_id=?", (aid,))
    get_db().execute("DELETE FROM agent_profiles WHERE agent_id=?", (aid,))
    get_db().execute("DELETE FROM agents WHERE id=?", (aid,))
    get_db().commit()
    log_audit('', 'admin.agent.delete', session.get('agent_id',''), session.get('agent_name',''),
              {'agent_id': aid})
    return jsonify({'ok':True})

@app.route('/api/admin/agents/<aid>/level', methods=['PUT'])
@admin_required
def admin_agent_level(aid):
    level = int(request.get_json().get('level', 1))
    get_db().execute("UPDATE agent_profiles SET agent_level=? WHERE agent_id=?", (level, aid))
    get_db().commit(); return jsonify({'ok':True})

# 工单管理（管理员视角）
@app.route('/api/admin/tickets', methods=['GET'])
@admin_required
def admin_tickets():
    status = request.args.get('status','')
    search = request.args.get('search','')
    page = int(request.args.get('page',1)); limit = int(request.args.get('limit',20))
    offset = (page-1)*limit
    db = get_db()
    where = []
    params = []
    if status: where.append("t.status=?"); params.append(status)
    if search: where.append("(cust.name LIKE ? OR t.id LIKE ?)"); params.extend([f'%{search}%', f'%{search}%'])
    where_sql = 'WHERE ' + ' AND '.join(where) if where else ''
    total = db.execute(f"SELECT COUNT(*) FROM service_tickets t JOIN customers cust ON t.customer_id=cust.id {where_sql}", params).fetchone()[0]
    tickets = db.execute(f"""
        SELECT t.*, cust.name as customer_name, ag.name as agent_name,
               ap.display_name as agent_display, cp.phone, cp.company
        FROM service_tickets t JOIN customers cust ON t.customer_id = cust.id
        LEFT JOIN agents ag ON t.agent_id = ag.id
        LEFT JOIN agent_profiles ap ON t.agent_id = ap.agent_id
        LEFT JOIN customer_profiles cp ON t.customer_id = cp.customer_id
        {where_sql} ORDER BY t.created_at DESC LIMIT ? OFFSET ?""", params + [limit, offset]).fetchall()
    stats = {
        'created': db.execute("SELECT COUNT(*) FROM service_tickets WHERE status='created'").fetchone()[0],
        'processing': db.execute("SELECT COUNT(*) FROM service_tickets WHERE status='processing'").fetchone()[0],
        'resolved': db.execute("SELECT COUNT(*) FROM service_tickets WHERE status='resolved'").fetchone()[0],
        'rated': db.execute("SELECT COUNT(*) FROM service_tickets WHERE status='rated'").fetchone()[0],
        'closed': db.execute("SELECT COUNT(*) FROM service_tickets WHERE status='closed'").fetchone()[0],
        'total': total
    }
    return jsonify({'tickets':[dict(t) for t in tickets],'stats':stats,'page':page,'pages':(total+limit-1)//limit})



@app.route('/api/admin/tickets/<tk_id>/remark', methods=['PUT'])
@sysadmin_required
def admin_ticket_remark(tk_id):
    data = request.get_json()
    remarks = data.get('remarks', '')
    get_db().execute("UPDATE service_tickets SET admin_remarks=?,updated_at=datetime('now','localtime') WHERE id=?", (remarks, tk_id))
    get_db().commit()
    return jsonify({'ok': True})

@app.route('/api/admin/tickets/<tk_id>', methods=['DELETE'])
@admin_required
def admin_ticket_delete(tk_id):
    t = get_db().execute("SELECT conversation_id FROM service_tickets WHERE id=?", (tk_id,)).fetchone()
    if not t:
        return jsonify({'error':'工单不存在'}),404
    conv_id = t['conversation_id']
    db = get_db()
    db.execute("DELETE FROM messages WHERE conversation_id=?", (conv_id,))
    db.execute("DELETE FROM service_tickets WHERE id=?", (tk_id,))
    db.execute("DELETE FROM escalations WHERE conversation_id=?", (conv_id,))
    db.execute("DELETE FROM conversations WHERE id=?", (conv_id,))
    db.commit()
    return jsonify({'ok': True})

@app.route('/api/admin/tickets/<tk_id>', methods=['GET'])
@logged_in_required
def admin_ticket_detail(tk_id):
    t = get_db().execute("""
        SELECT t.*, cust.name as customer_name, cust.id as customer_id,
               ag.name as agent_name, ag.email as agent_email,
               ap.display_name as agent_display, ap.department, ap.title,
               cp.name as cp_name, cp.phone as cp_phone, cp.company, cp.notes
        FROM service_tickets t
        JOIN customers cust ON t.customer_id = cust.id
        LEFT JOIN agents ag ON t.agent_id = ag.id
        LEFT JOIN agent_profiles ap ON t.agent_id = ap.agent_id
        LEFT JOIN customer_profiles cp ON t.customer_id = cp.customer_id
        WHERE t.id=?""", (tk_id,)).fetchone()
    if not t: return jsonify({'error':'工单不存在'}),404
    msgs = get_db().execute(
        "SELECT role,content,image_url,created_at FROM messages WHERE conversation_id=? ORDER BY created_at", (t['conversation_id'],)).fetchall()
    return jsonify({'ticket':dict(t),'messages':[dict(m) for m in msgs]})

# ====== 系统配置管理 ======
@app.route('/api/admin/config', methods=['GET','POST'])
@sysadmin_required
def admin_config():
    if request.method == 'POST':
        data = request.get_json()
        db = get_db()
        for k, v in data.items():
            db.execute("INSERT OR REPLACE INTO system_config(key,value,updated_at) VALUES(?,?,datetime('now','localtime'))", (k, str(v)))
        db.commit()
        # Sync admin password to agents table
        if 'admin_password' in data:
            pw = hashlib.sha256(f'admin:{data["admin_password"]}'.encode()).hexdigest()
            db.execute("UPDATE agents SET password_hash=? WHERE email=?", (pw, 'admin@smartcs.com'))
            db.commit()
        global DASHSCOPE_API_KEY, API_BASE_URL, MODEL_NAME, ADMIN_PASSWORD
        rows = get_db().execute("SELECT key,value FROM system_config").fetchall()
        for r in rows:
            if r['key'] == 'api_key': DASHSCOPE_API_KEY = r['value']
            elif r['key'] == 'api_base_url': API_BASE_URL = r['value']
            elif r['key'] == 'model_name': MODEL_NAME = r['value']
            elif r['key'] == 'admin_password': ADMIN_PASSWORD = r['value']
        invalidate_config_cache()
        return jsonify({'ok':True})
    configs = get_db().execute("SELECT key,value FROM system_config").fetchall()
    return jsonify({c['key']:c['value'] for c in configs})

# ====== 安全配置管理（安全管理员专属） ======
SECURITY_CONFIG_KEYS = {'password_min_length','password_require_upper','password_expire_days','login_max_attempts','login_lockout_minutes','audit_log_retention_days'}

@app.route('/api/admin/security-config', methods=['GET','POST'])
@secadmin_required
def admin_security_config():
    if request.method == 'POST':
        data = request.get_json()
        db = get_db()
        for k, v in data.items():
            if k in SECURITY_CONFIG_KEYS:
                db.execute("INSERT OR REPLACE INTO system_config(key,value,updated_at) VALUES(?,?,datetime('now','localtime'))", (k, str(v)))
        db.commit()
        log_audit('', 'admin.security_config.update', session.get('agent_id',''), session.get('agent_name',''),
                  {'keys_updated': list(data.keys())})
        return jsonify({'ok':True})
    configs = get_db().execute("SELECT key,value FROM system_config").fetchall()
    result = {}
    for c in configs:
        if c['key'] in SECURITY_CONFIG_KEYS:
            result[c['key']] = c['value']
    return jsonify(result)

# ====== 品牌配置管理 ======
BRAND_CONFIG_KEYS = {'brand_name','brand_short','brand_primary_color','brand_logo_path','brand_favicon_path'}

@app.route('/api/admin/brand-config', methods=['GET','POST'])
@sysadmin_required
def admin_brand_config():
    if request.method == 'POST':
        data = request.get_json()
        db = get_db()
        for k, v in data.items():
            if k in BRAND_CONFIG_KEYS:
                db.execute("INSERT OR REPLACE INTO system_config(key,value,updated_at) VALUES(?,?,datetime('now','localtime'))", (k, str(v)))
        db.commit()
        invalidate_config_cache()
        log_audit('', 'admin.brand_config.update', session.get('agent_id',''), session.get('agent_name',''),
                  {'keys_updated': list(data.keys())})
        return jsonify({'ok':True})
    configs = get_db().execute("SELECT key,value FROM system_config").fetchall()
    result = {}
    for c in configs:
        if c['key'] in BRAND_CONFIG_KEYS:
            result[c['key']] = c['value']
    return jsonify(result)

@app.route('/api/admin/tickets/status', methods=['POST'])
@admin_required
def admin_tickets_status():
    """Admin force-updates ticket status (for force close etc)"""
    data = request.get_json() or {}
    tk_id = data.get('ticket_id', '')
    new_status = data.get('status', '')
    reason = data.get('reason', '')
    if not tk_id or not new_status:
        return jsonify({'error': '请填写工单ID和目标状态'}), 400
    db = get_db()
    t = db.execute("SELECT id, status, conversation_id FROM service_tickets WHERE id=?", (tk_id,)).fetchone()
    if not t:
        return jsonify({'error': '工单不存在'}), 404
    ok, msg = transition_ticket(tk_id, new_status, session.get('agent_id','admin'),
                                session.get('agent_name','管理员'),
                                {'close_reason': reason or '管理员操作',
                                 'admin_remarks': reason})
    if not ok:
        return jsonify({'error': msg or '状态切换失败'}), 400
    db.execute(
        "UPDATE conversations SET status=?,updated_at=datetime('now','localtime') "
        "WHERE id=?",
        (new_status, t['conversation_id'])
    )
    if t['conversation_id']:
        agent_name = session.get('agent_name', '管理员')
        msg_text = f'🔧 管理员 {agent_name} 已将工单状态变更为[{new_status}]'
        if reason:
            msg_text += f'（原因：{reason}）'
        db.execute(
            "INSERT INTO messages(id,conversation_id,role,content) VALUES(?,?,?,?)",
            (gen_id('msg-'), t['conversation_id'], 'system', msg_text)
        )
    db.commit()
    log_audit(tk_id, 'admin.ticket_status', session.get('agent_id',''), session.get('agent_name',''),
              {'old_status': t['status'], 'new_status': new_status, 'reason': reason})
    return jsonify({'ok': True})


@app.route('/api/admin/close-reasons', methods=['GET'])
@sysadmin_required
def admin_close_reasons():
    reasons = get_db().execute("SELECT * FROM close_reasons ORDER BY sort_order").fetchall()
    return jsonify([dict(r) for r in reasons])

@app.route('/api/admin/close-reasons', methods=['POST'])
@sysadmin_required
def admin_add_close_reason():
    data = request.get_json()
    name = data.get('name','').strip()
    if not name: return jsonify({'error':'请输入原因名称'}),400
    max_order = get_db().execute("SELECT COALESCE(MAX(sort_order),0)+1 FROM close_reasons").fetchone()[0]
    db.execute("INSERT INTO close_reasons(id,name,sort_order) VALUES(?,?,?)", (gen_id('cr-'), name, max_order))
    db.commit()
    return jsonify({'ok':True})

@app.route('/api/admin/close-reasons/<rid>', methods=['DELETE'])
@sysadmin_required
def admin_delete_close_reason(rid):
    db.execute("DELETE FROM close_reasons WHERE id=?", (rid,))
    db.commit()
    return jsonify({'ok':True})



# ====== 一线转二线 ======
@app.route('/api/agent/tickets/<tk_id>/transfer', methods=['POST'])
@logged_in_required
def agent_transfer_ticket(tk_id):
    aid = session['agent_id']
    data = request.get_json()
    target_id = data.get('target_agent_id', '')
    note = data.get('note', '')
    t = get_db().execute("SELECT id,agent_id FROM service_tickets WHERE id=?", (tk_id,)).fetchone()
    if not t:
        return jsonify({'error': '\u5de5\u5355\u4e0d\u5b58\u5728'}), 404
    if t['agent_id'] != aid:
        return jsonify({'error': '\u53ea\u80fd\u8f6c\u4ea4\u81ea\u5df1\u7684\u5de5\u5355'}), 403
    target = get_db().execute("""
        SELECT ag.id, ap.display_name, ag.name
        FROM agents ag LEFT JOIN agent_profiles ap ON ag.id = ap.agent_id
        WHERE ag.id=?""", (target_id,)).fetchone()
    if not target:
        return jsonify({'error': '\u76ee\u6807\u5ba2\u670d\u4e0d\u53ef\u7528'}), 400
    db = get_db()
    # Reset first_read_at so the receiving engineer sees it as a new ticket
    db.execute("UPDATE service_tickets SET agent_id=?,transferred_from=?,first_read_at='',updated_at=datetime('now','localtime') WHERE id=?", (target_id, aid, tk_id))
    conv = db.execute("SELECT conversation_id FROM service_tickets WHERE id=?", (tk_id,)).fetchone()
    if conv:
        source_name = session.get('agent_name', '\u5ba2\u670d')
        target_name = target['display_name'] or target['name']
        note_text = '\u241e \u5ba2\u670d ' + source_name + ' \u5c06\u5de5\u5355\u8f6c\u4ea4\u7ed9 ' + target_name
        if note:
            note_text += '\uff08\u5907\u6ce8\uff1a' + note + '\uff09'
        db.execute("INSERT INTO messages(id,conversation_id,role,content) VALUES(?,?,?,?)",
                   (gen_id('msg-'), conv['conversation_id'], 'system', note_text))
        db.execute("UPDATE conversations SET updated_at=datetime('now','localtime') WHERE id=?", (conv['conversation_id'],))
    db.commit()
    log_audit(tk_id, 'ticket.transferred', aid, session.get('agent_name',''),
              {'from_agent':aid, 'to_agent':target_id, 'note':note})
    return jsonify({'ok': True, 'target': target_name})

@app.route('/api/agent/tickets/<tk_id>/transfer-agents', methods=['GET'])
@logged_in_required
def agent_transfer_agents(tk_id):
    aid = session['agent_id']
    agents = get_db().execute("""
        SELECT ag.id, COALESCE(ap.display_name, ag.name) as display_name, ag.name, ap.department, ap.title, ap.agent_level
        FROM agents ag LEFT JOIN agent_profiles ap ON ag.id = ap.agent_id
        WHERE ag.id != ? AND ag.role != 'admin'
        ORDER BY ap.agent_level DESC, ap.display_name""", (aid,)).fetchall()
    result = []
    for a in agents:
        a_dict = dict(a)
        systems = get_db().execute("""
            SELECT s.name FROM systems s JOIN agent_systems a_s ON s.id = a_s.system_id
            WHERE a_s.agent_id=? ORDER BY s.name""", (a['id'],)).fetchall()
        a_dict['systems'] = [s['name'] for s in systems]
        result.append(a_dict)
    return jsonify(result)

# ====== 升级记录管理 ======
@app.route('/api/admin/escalations', methods=['GET'])
@sysadmin_required
def admin_escalations():
    page = int(request.args.get('page',1))
    limit = int(request.args.get('limit',20))
    offset = (page-1)*limit
    search = request.args.get('search','')
    date_from = request.args.get('from','')
    date_to = request.args.get('to','')
    db = get_db()
    where = []
    params = []
    if search:
        where.append("(e.reason LIKE ? OR c.name LIKE ?)")
        params.extend(["%" + search + "%", "%" + search + "%"])
    if date_from:
        where.append("e.created_at >= ?")
        params.append(date_from)
    if date_to:
        where.append("e.created_at <= ?")
        params.append(date_to + " 23:59:59")
    base_join = (" FROM escalations e"
                 " LEFT JOIN conversations cv ON e.conversation_id=cv.id"
                 " LEFT JOIN customers c ON cv.customer_id=c.id"
                 " LEFT JOIN agent_profiles a ON e.agent_id=a.agent_id")
    wsql = " AND ".join(where) if where else "1=1"
    total = db.execute("SELECT COUNT(*)" + base_join + " WHERE " + wsql, params).fetchone()[0]
    es = db.execute(
        "SELECT e.*, c.name as customer_name, a.display_name as agent_display, a.display_name as agent_name"
        + base_join + " WHERE " + wsql
        + " ORDER BY e.created_at DESC LIMIT ? OFFSET ?",
        params + [limit, offset]
    ).fetchall()
    stats = db.execute("SELECT COALESCE(COUNT(*),0) as total,"
        " COALESCE(SUM(CASE WHEN status='pending' THEN 1 ELSE 0 END),0) as pending,"
        " COALESCE(SUM(CASE WHEN status='assigned' THEN 1 ELSE 0 END),0) as assigned,"
        " COALESCE(SUM(CASE WHEN status='resolved' THEN 1 ELSE 0 END),0) as resolved"
        " FROM escalations").fetchone()
    pages = max(1, (total + limit - 1) // limit)
    return jsonify({
        'escalations': [dict(r) for r in es],
        'stats': dict(stats),
        'page': page,
        'pages': pages,
        'total': total
    })

# ====== 系统升级记录管理 ======
@app.route('/api/admin/upgrades', methods=['GET'])
@sysadmin_required
def admin_upgrades():
    page = int(request.args.get('page',1))
    limit = int(request.args.get('limit',20))
    offset = (page-1)*limit
    db = get_db()
    total = db.execute("SELECT COUNT(*) FROM system_upgrades").fetchone()[0]
    upgrades = db.execute("SELECT * FROM system_upgrades ORDER BY release_date DESC LIMIT ? OFFSET ?", (limit, offset)).fetchall()
    pages = max(1, (total + limit - 1) // limit)
    return jsonify({'upgrades':[dict(u) for u in upgrades], 'page':page, 'pages':pages, 'total':total})

@app.route('/api/admin/upgrades', methods=['POST'])
@sysadmin_required
def admin_add_upgrade():
    data = request.get_json()
    version = data.get('version','').strip()
    release_date = data.get('release_date','')
    content = data.get('content','').strip()
    if not version: return jsonify({'error':'请输入版本号'}),400
    get_db().execute("INSERT INTO system_upgrades(id,version,release_date,content) VALUES(?,?,?,?)",
                     (gen_id('ug-'), version, release_date, content))
    get_db().commit()
    return jsonify({'ok':True})

@app.route('/api/admin/upgrades/<uid>', methods=['DELETE'])
@sysadmin_required
def admin_delete_upgrade(uid):
    get_db().execute("DELETE FROM system_upgrades WHERE id=?", (uid,))
    get_db().commit()
    return jsonify({'ok':True})

# 图片上传（工单附件）
@app.route('/api/upload/image', methods=['POST'])
def upload_image():
    password = request.form.get('password','') or request.args.get('password','')
    if password != ADMIN_PASSWORD and not session.get('agent_id'):
        return jsonify({'error':'权限不足'}),403
    file = request.files.get('file')
    if not file: return jsonify({'error':'未选择文件'}),400
    ext = file.filename.rsplit('.',1)[-1].lower() if '.' in file.filename else ''
    if ext not in {'jpg','jpeg','png','gif','webp','bmp'}: return jsonify({'error':'不支持格式'}),400
    ts = datetime.now().strftime('%Y%m%d_%H%M%S')
    fn = f"{ts}_{uuid.uuid4().hex[:8]}.{ext}"
    file.save(os.path.join(UPLOAD_DIR, fn))
    return jsonify({'ok':True,'url':f'/uploads/{fn}'})

@app.route('/uploads/<path:filename>')
def uploaded_file(filename):
    return send_from_directory(UPLOAD_DIR, filename)

# ====== 知识库 ======
@app.route('/upload') 
def upload_page(): return render_template('upload.html')

@app.route('/api/knowledge/list', methods=['GET'])
@sysadmin_required
def knowledge_list():
    files = []
    for f in sorted(glob.glob(os.path.join(KNOWLEDGE_DIR,'*.md')), key=os.path.getmtime, reverse=True):
        s = os.stat(f); files.append({'name':os.path.basename(f),'size':s.st_size,'modified':datetime.fromtimestamp(s.st_mtime).strftime('%Y-%m-%d %H:%M')})
    return jsonify(files)

@app.route('/api/upload', methods=['POST'])
@sysadmin_required
def upload_file():
    file = request.files.get('file')
    if not file or file.filename=='': return jsonify({'error':'未选择文件'}),400
    ext = file.filename.rsplit('.',1)[-1].lower() if '.' in file.filename else ''
    if ext not in {'docx','pdf','pptx','md','txt'}: return jsonify({'error':f'不支持格式: {ext}'}),400
    ts = datetime.now().strftime('%Y%m%d_%H%M%S')
    sp = os.path.join(UPLOAD_DIR, f"{ts}_{file.filename}")
    file.save(sp); text = ""
    try:
        if ext == 'docx':
            from docx import Document; text='\n'.join(p.text for p in Document(sp).paragraphs if p.text.strip())
        elif ext == 'pdf':
            import pdfplumber
            with pdfplumber.open(sp) as pdf: text='\n\n'.join(p.extract_text() or '' for p in pdf.pages)
        elif ext == 'pptx':
            from pptx import Presentation
            for sl in Presentation(sp).slides:
                for sh in sl.shapes:
                    if hasattr(sh,'text') and sh.text.strip(): text+=sh.text.strip()+'\n'
        else:
            with open(sp,'r',encoding='utf-8') as f: text=f.read()
    except: pass
    md_name = file.filename.rsplit('.',1)[0]+'.md'
    with open(os.path.join(KNOWLEDGE_DIR,md_name),'w',encoding='utf-8') as f: f.write(text)
    return jsonify({'success':True,'filename':md_name,'word_count':len(text)})

@app.route('/api/knowledge/detail', methods=['GET'])
@sysadmin_required
def knowledge_detail():
    fp = os.path.join(KNOWLEDGE_DIR, request.args.get('filename',''))
    if not os.path.exists(fp): return jsonify({'error':'文件不存在'}),404
    with open(fp,'r',encoding='utf-8') as f: return jsonify({'content':f.read()})

@app.route('/api/knowledge/update', methods=['POST'])
@sysadmin_required
def knowledge_update():
    data = request.get_json()
    fp = os.path.join(KNOWLEDGE_DIR, data.get('filename',''))
    if not os.path.exists(fp): return jsonify({'error':'文件不存在'}),404
    with open(fp,'w',encoding='utf-8') as f: f.write(data.get('content',''))
    return jsonify({'success':True})

@app.route('/api/knowledge/delete', methods=['POST'])
@sysadmin_required
def knowledge_delete():
    data = request.get_json()
    fp = os.path.join(KNOWLEDGE_DIR, data.get('filename',''))
    if not os.path.exists(fp): return jsonify({'error':'文件不存在'}),404
    os.remove(fp); return jsonify({'success':True})



# ====== 负责系统管理 ======
@app.route('/api/admin/systems', methods=['GET'])
@admin_required
def admin_systems():
    systems = get_db().execute("SELECT * FROM systems ORDER BY name").fetchall()
    return jsonify([dict(s) for s in systems])

@app.route('/api/admin/systems', methods=['POST'])
@admin_required
def admin_add_system():
    data = request.get_json()
    name = data.get('name', '').strip()
    if not name: return jsonify({'error': '请输入系统名称'}), 400
    desc = data.get('description', '').strip()
    existing = get_db().execute("SELECT id FROM systems WHERE name=?", (name,)).fetchone()
    if existing: return jsonify({'error': '系统名称已存在'}), 400
    get_db().execute("INSERT INTO systems(id,name,description) VALUES(?,?,?)",
                     (gen_id('sys-'), name, desc))
    get_db().commit()
    return jsonify({'ok': True})

@app.route('/api/admin/systems/<sid>', methods=['PUT'])
@admin_required
def admin_update_system(sid):
    data = request.get_json()
    name = data.get('name', '').strip()
    desc = data.get('description', '').strip()
    if not name: return jsonify({'error': '请输入系统名称'}), 400
    existing = get_db().execute("SELECT id FROM systems WHERE name=? AND id!=?", (name, sid)).fetchone()
    if existing: return jsonify({'error': '系统名称已存在'}), 400
    get_db().execute("UPDATE systems SET name=?,description=? WHERE id=?", (name, desc, sid))
    get_db().commit()
    return jsonify({'ok': True})

@app.route('/api/admin/systems/<sid>', methods=['DELETE'])
@admin_required
def admin_delete_system(sid):
    db = get_db()
    db.execute("DELETE FROM agent_systems WHERE system_id=?", (sid,))
    db.execute("DELETE FROM systems WHERE id=?", (sid,))
    db.commit()
    return jsonify({'ok': True})

@app.route('/api/admin/agents/<aid>/systems', methods=['GET', 'POST'])
@admin_required
def admin_agent_systems(aid):
    db = get_db()
    if request.method == 'POST':
        data = request.get_json()
        system_ids = data.get('system_ids', [])
        db.execute("DELETE FROM agent_systems WHERE agent_id=?", (aid,))
        for sid in system_ids:
            db.execute("INSERT OR IGNORE INTO agent_systems(agent_id,system_id) VALUES(?,?)", (aid, sid))
        db.commit()
        return jsonify({'ok': True})
    systems = db.execute("""
        SELECT s.*, CASE WHEN a_s.agent_id IS NOT NULL THEN 1 ELSE 0 END as assigned
        FROM systems s LEFT JOIN agent_systems a_s ON s.id = a_s.system_id AND a_s.agent_id=?
        ORDER BY s.name""", (aid,)).fetchall()
    return jsonify([dict(s) for s in systems])

# ====== 系统 ======
@app.route('/api/config')
def api_config():
    return jsonify({'api_key_configured':bool(DASHSCOPE_API_KEY),'model':MODEL_NAME,
        'tickets_pending':get_db().execute("SELECT COUNT(*) FROM service_tickets WHERE status='created'").fetchone()[0]})

@app.route('/api/stats')
def api_stats():
    db = get_db()
    return jsonify({'total_conversations':db.execute("SELECT COUNT(*) FROM conversations").fetchone()[0],
        'active':db.execute("SELECT COUNT(*) FROM conversations WHERE status='active'").fetchone()[0],
        'pending_tickets':db.execute("SELECT COUNT(*) FROM service_tickets WHERE status='created'").fetchone()[0],
        'knowledge_files':len(glob.glob(os.path.join(KNOWLEDGE_DIR,'*.md')))})



# ====== Phase 3: New Admin API Endpoints ======

@app.route('/api/admin/stats/overview', methods=['GET'])
@logged_in_required
def admin_stats_overview():
    db = get_db()
    today = datetime.now().strftime('%Y-%m-%d')
    week_start = (datetime.now() - timedelta(days=datetime.now().weekday())).strftime('%Y-%m-%d')
    month_start = datetime.now().strftime('%Y-%m-01')
    r = {'total_tickets': db.execute('SELECT COUNT(*) FROM service_tickets').fetchone()[0]}
    for period, date_str in [('today',today),('this_week',week_start),('this_month',month_start)]:
        p = {}
        for st in ['created','processing','resolved','rated','closed']:
            p[st] = db.execute('SELECT COUNT(*) FROM service_tickets WHERE status=? AND created_at>=?',(st,date_str)).fetchone()[0]
        p['total'] = sum(p.values())
        r[period] = p
    resp_data = db.execute(
        'SELECT COALESCE(AVG((julianday(m.created_at)-julianday(t.assigned_at))*24*60),0) as rt '
        'FROM service_tickets t JOIN messages m ON m.conversation_id=t.conversation_id AND m.role=? '
        'WHERE t.assigned_at IS NOT NULL AND m.created_at>=t.assigned_at',
        ('agent',)
    ).fetchone()
    r['avg_response_time'] = round(resp_data['rt'],1) if resp_data else 0
    ar = db.execute('SELECT COALESCE(AVG(customer_rating),0) as r FROM service_tickets WHERE customer_rating>0').fetchone()
    r['avg_rating'] = round(ar['r'],1) if ar else 0
    tbd = []
    for i in range(6,-1,-1):
        d = (datetime.now()-timedelta(days=i)).strftime('%Y-%m-%d')
        nd = (datetime.now()-timedelta(days=i-1)).strftime('%Y-%m-%d')
        cnt = db.execute('SELECT COUNT(*) as c FROM service_tickets WHERE created_at>=? AND created_at<?',(d,nd)).fetchone()['c']
        tbd.append({'date':d,'count':cnt})
    r['tickets_by_date'] = tbd
    return jsonify(r)


@app.route('/api/admin/stats/agents', methods=['GET'])
@logged_in_required
def admin_stats_agents():
    db = get_db()
    agents = db.execute(
        'SELECT ag.id, ag.name, COALESCE(ap.display_name,ag.name) as display_name '
        'FROM agents ag LEFT JOIN agent_profiles ap ON ag.id=ap.agent_id '
        'WHERE ag.role!=? ORDER BY ag.name',
        ('admin',)
    ).fetchall()
    result = []
    for a in agents:
        rc = db.execute("SELECT COUNT(*) as c FROM service_tickets WHERE agent_id=? AND status IN ('closed','resolved','rated')",(a['id'],)).fetchone()['c']
        avg_r = db.execute('SELECT COALESCE(AVG(customer_rating),0) as r FROM service_tickets WHERE agent_id=? AND customer_rating>0',(a['id'],)).fetchone()['r']
        avg_resp = db.execute(
            'SELECT COALESCE(AVG((julianday(m.created_at)-julianday(t.assigned_at))*24*60),0) as rt '
            'FROM service_tickets t JOIN messages m ON m.conversation_id=t.conversation_id AND m.role=? '
            'WHERE t.assigned_at IS NOT NULL AND m.created_at>=t.assigned_at AND t.agent_id=?',
            ('agent', a['id'])
        ).fetchone()['rt']
        act = db.execute("SELECT COUNT(*) as c FROM service_tickets WHERE agent_id=? AND status NOT IN ('closed','resolved','rated')",(a['id'],)).fetchone()['c']
        result.append({'agent_id':a['id'],'name':a['name'],'display_name':a['display_name'],
                       'resolved_count':rc,'avg_rating':round(avg_r,1),'avg_response_time':round(avg_resp,1),'active_tickets':act})
    return jsonify(result)


@app.route('/api/admin/tickets/stats', methods=['GET'])
@admin_required
def admin_tickets_stats():
    """Simple ticket status counts for admin dashboard"""
    db = get_db()
    stats = {
        'created': db.execute("SELECT COUNT(*) FROM service_tickets WHERE status='created'").fetchone()[0],
        'processing': db.execute("SELECT COUNT(*) FROM service_tickets WHERE status='processing'").fetchone()[0],
        'resolved': db.execute("SELECT COUNT(*) FROM service_tickets WHERE status='resolved'").fetchone()[0],
        'rated': db.execute("SELECT COUNT(*) FROM service_tickets WHERE status='rated'").fetchone()[0],
        'closed': db.execute("SELECT COUNT(*) FROM service_tickets WHERE status='closed'").fetchone()[0],
        'total': db.execute("SELECT COUNT(*) FROM service_tickets").fetchone()[0]
    }
    return jsonify(stats)


@app.route('/api/admin/analytics', methods=['GET'])
@logged_in_required
def admin_analytics():
    """Data analysis dashboard data"""
    db = get_db()
    
    # Tickets by status
    tickets_by_status = {}
    for st in ['created','processing','resolved','rated','closed']:
        tickets_by_status[st] = db.execute(f"SELECT COUNT(*) FROM service_tickets WHERE status=?", (st,)).fetchone()[0]
    
    # Tickets by agent
    agents = db.execute(
        'SELECT ag.id, ag.name, COALESCE(ap.display_name,ag.name) as display_name, ap.agent_level '
        'FROM agents ag LEFT JOIN agent_profiles ap ON ag.id=ap.agent_id WHERE ag.role!=? ORDER BY ag.name',
        ('admin',)
    ).fetchall()
    tickets_by_agent = []
    for a in agents:
        total = db.execute("SELECT COUNT(*) FROM service_tickets WHERE agent_id=?", (a['id'],)).fetchone()[0]
        closed = db.execute("SELECT COUNT(*) FROM service_tickets WHERE agent_id=? AND status='closed'", (a['id'],)).fetchone()[0]
        rated = db.execute("SELECT COUNT(*) FROM service_tickets WHERE agent_id=? AND status='rated'", (a['id'],)).fetchone()[0]
        total_closed_or_rated = db.execute("SELECT COUNT(*) FROM service_tickets WHERE agent_id=? AND status IN ('closed','rated')", (a['id'],)).fetchone()[0]
        interception_rate = round(total_closed_or_rated / total * 100, 1) if total > 0 else 0
        tickets_by_agent.append({
            'agent_name': a['display_name'] or a['name'],
            'total': total,
            'closed': closed,
            'rated': rated,
            'interception_rate': interception_rate
        })
    
    # Tickets by system
    systems = db.execute("SELECT id, name FROM systems ORDER BY name").fetchall()
    tickets_by_system = []
    for s in systems:
        count = db.execute("""
            SELECT COUNT(*) FROM service_tickets t
            JOIN agent_systems a_s ON t.agent_id = a_s.agent_id
            WHERE a_s.system_id=?""", (s['id'],)).fetchone()[0]
        tickets_by_system.append({'system_name': s['name'], 'count': count})
    
    return jsonify({
        'tickets_by_status': tickets_by_status,
        'tickets_by_agent': tickets_by_agent,
        'tickets_by_system': tickets_by_system
    })


@app.route('/api/admin/tickets/export', methods=['GET'])
@admin_required
def admin_tickets_export():
    status = request.args.get('status','')
    date_from = request.args.get('from','')
    date_to = request.args.get('to','')
    db = get_db()
    where = []; params = []
    if status: where.append('t.status=?'); params.append(status)
    if date_from: where.append('t.created_at>=?'); params.append(date_from)
    if date_to: where.append('t.created_at<=?'); params.append(date_to+' 23:59:59')
    wsql = ' AND '.join(where) if where else '1=1'
    tickets = db.execute(
        'SELECT t.*,cust.name as cn,COALESCE(ap.display_name,ag.name,\"\") as an '
        'FROM service_tickets t JOIN customers cust ON t.customer_id=cust.id '
        'LEFT JOIN agents ag ON t.agent_id=ag.id '
        'LEFT JOIN agent_profiles ap ON t.agent_id=ap.agent_id '
        'WHERE ' + wsql + ' ORDER BY t.created_at DESC',
        params
    ).fetchall()
    output = io.StringIO()
    output.write(chr(0xfeff))
    writer = csv.writer(output)
    writer.writerow(['\\u5de5\\u5355\\u7f16\\u53f7','\\u72b6\\u6001','\\u5ba2\\u6237','\\u5ba2\\u670d','\\u95ee\\u9898\\u63cf\\u8ff0','\\u8bc4\\u5206','\\u53cd\\u9988','\\u521b\\u5efa\\u65f6\\u95f4','\\u89e3\\u51b3\\u65f6\\u95f4'])
    sm = {'created':'\\u5df2\\u521b\\u5efa','processing':'\\u5904\\u7406\\u4e2d','resolved':'\\u5df2\\u89e3\\u51b3','rated':'\\u5df2\\u8bc4\\u4ef7','closed':'\\u5df2\\u5173\\u95ed'}
    for t in tickets:
        writer.writerow([t['ticket_number'] or t['id'][:12],sm.get(t['status'],t['status']),
                         t['cn'] or '\\u6e38\\u5ba2',t['an'] or '',t['issue_description'] or '',
                         t['customer_rating'] or '',t['customer_feedback'] or '',
                         t['created_at'] or '',t['resolved_at'] or t['closed_at'] or ''])
    output.seek(0)
    return Response(output.getvalue(),mimetype='text/csv; charset=utf-8-sig',
              headers={'Content-Disposition':'attachment; filename=smartcs_tickets_export.csv','Content-Type':'text/csv; charset=utf-8-sig'})


@app.route('/api/admin/tickets/archive', methods=['POST'])
@admin_required
def admin_tickets_archive():
    return _do_archive()

@app.route('/api/cron/auto-close', methods=['GET'])
def cron_auto_close():
    """Auto-close and auto-rate tickets.
    - Tickets with status='created' and no user messages in auto_close_min minutes → closed
    - Tickets with status='resolved' and no user activity in auto_rate_hours hours → rated
    """
    db = get_db()
    cfg = {r['key']: r['value'] for r in db.execute('SELECT key, value FROM system_config').fetchall()}
    auto_close_min = int(cfg.get('auto_close_min', '20'))
    auto_rate_hours = int(cfg.get('auto_rate_hours', '24'))
    now = datetime.now()
    closed_count = 0
    rated_count = 0

    # Auto-close: created tickets with no user messages in auto_close_min minutes
    cutoff = (now - timedelta(minutes=auto_close_min)).strftime('%Y-%m-%d %H:%M:%S')
    created_tickets = db.execute("""
        SELECT t.id, t.conversation_id FROM service_tickets t
        WHERE t.status='created'
        AND NOT EXISTS (
            SELECT 1 FROM messages m
            WHERE m.conversation_id = t.conversation_id
            AND m.role='user'
            AND m.created_at > ?
        )
    """, (cutoff,)).fetchall()
    for t in created_tickets:
        ok, _ = transition_ticket(t['id'], 'closed', 'system', '系统自动关闭',
                                  {'close_reason': '自动关闭（超时无用户消息）'})
        if ok:
            closed_count += 1
            db.execute("UPDATE conversations SET status='closed',updated_at=datetime('now','localtime') WHERE id=?", (t['conversation_id'],))
            log_audit(t['id'], 'auto_close', 'system', '系统', {'reason': '超时未回复', 'conversation_id': t['conversation_id']})

    # Auto-rate: resolved tickets with no user activity in auto_rate_hours hours
    cutoff_rate = (now - timedelta(hours=auto_rate_hours)).strftime('%Y-%m-%d %H:%M:%S')
    resolved_tickets = db.execute("""
        SELECT t.id, t.conversation_id FROM service_tickets t
        WHERE t.status='resolved'
        AND (
            NOT EXISTS (
                SELECT 1 FROM messages m
                WHERE m.conversation_id = t.conversation_id
                AND m.created_at > ?
            )
        )
    """, (cutoff_rate,)).fetchall()
    for t in resolved_tickets:
        ok, _ = transition_ticket(t['id'], 'rated', 'system', '系统自动评价',
                                  {'close_reason': '用户超时未评价', 'customer_rating': 5,
                                   'customer_feedback': '系统自动评价（用户超时未操作）'})
        if ok:
            rated_count += 1
            db.execute("UPDATE conversations SET status='resolved',updated_at=datetime('now','localtime') WHERE id=?", (t['conversation_id'],))
            log_audit(t['id'], 'auto_rate', 'system', '系统', {'reason': '用户超时未评价', 'rating': 5, 'conversation_id': t['conversation_id']})

    db.commit()
    return jsonify({'ok': True, 'closed': closed_count, 'rated': rated_count})


@app.route('/api/cron/archive', methods=['GET'])
def cron_archive():
    return _do_archive()

def _do_archive():
    db = get_db()
    cutoff = (datetime.now()-timedelta(days=90)).strftime('%Y-%m-%d %H:%M:%S')
    old = db.execute(
        'SELECT * FROM service_tickets '
        "WHERE status IN ('closed','resolved','rated') "
        "AND (coalesce(resolved_at,'2000-01-01')<? OR coalesce(closed_at,'2000-01-01')<?)",
        (cutoff,cutoff)
    ).fetchall()
    c = 0
    for t in old:
        try:
            db.execute(
                'INSERT OR IGNORE INTO tickets_archive('
                'id,escalation_id,conversation_id,customer_id,agent_id,'
                'status,issue_description,priority,image_url,assigned_at,resolved_at,confirmed_at,closed_at,'
                'resolution_notes,customer_rating,customer_feedback,created_at,updated_at,ticket_number,close_reason,'
                'admin_remarks,transferred_from,level,reopened_at,reopened_count,archived_at'
                ') VALUES(?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,?,'
                "datetime('now','localtime'))",
                (t['id'],t['escalation_id'] or '',t['conversation_id'],t['customer_id'],
                 t['agent_id'] or '','archived',t['issue_description'],t['priority'],
                 t['image_url'] or '',t['assigned_at'],t['resolved_at'],t['confirmed_at'],
                 t['closed_at'],t['resolution_notes'],t['customer_rating'],t['customer_feedback'],
                 t['created_at'],t['updated_at'],t['ticket_number'] or '',t['close_reason'] or '',
                 t['admin_remarks'] or '',t['transferred_from'] or '',t['level'],t['reopened_at'] or '',t['reopened_count']))
            db.execute('DELETE FROM service_tickets WHERE id=?',(t['id'],))
            # 合规要求：审计日志必须保留，不可删除
            # db.execute('DELETE FROM audit_log WHERE ticket_id=?',(t['id'],))
            c += 1
        except Exception as e:
            print('Archive error {}: {}'.format(t['id'], e))
    db.commit()
    return jsonify({'ok':True,'archived':c,'message':'\\u5df2\\u5f52\\u6863 '+str(c)+' \\u4e2a\\u5de5\\u5355'})

# ====== End Phase 3 ======


# ====== Phase 4: Webhook CRUD Endpoints ======

@app.route('/api/admin/webhooks', methods=['GET'])
@admin_required
def admin_list_webhooks():
    hooks = get_db().execute('SELECT * FROM webhooks ORDER BY created_at DESC').fetchall()
    return jsonify([dict(h) for h in hooks])


@app.route('/api/admin/webhooks', methods=['POST'])
@admin_required
def admin_add_webhook():
    data = request.get_json()
    wid = 'wh-' + uuid.uuid4().hex[:12]
    get_db().execute(
        'INSERT INTO webhooks(id, name, url, secret, events, retry, timeout) VALUES(?,?,?,?,?,?,?)',
        (wid, data['name'], data['url'], data.get('secret', ''),
         json.dumps(data.get('events', ['*']), ensure_ascii=False),
         data.get('retry', 3), data.get('timeout', 10))
    )
    get_db().commit()
    return jsonify({'ok': True, 'id': wid})


@app.route('/api/admin/webhooks/<wid>', methods=['PUT'])
@admin_required
def admin_update_webhook(wid):
    data = request.get_json()
    fields = []
    vals = []
    for k in ('name','url','secret','retry','timeout','enabled'):
        if k in data:
            fields.append(k + '=?')
            vals.append(data[k])
    if 'events' in data:
        fields.append('events=?')
        vals.append(json.dumps(data['events'], ensure_ascii=False))
    if fields:
        vals.append(wid)
        get_db().execute('UPDATE webhooks SET ' + ','.join(fields) + ', updated_at=datetime("now","localtime") WHERE id=?', vals)
        get_db().commit()
    return jsonify({'ok': True})


@app.route('/api/admin/webhooks/<wid>', methods=['DELETE'])
@admin_required
def admin_delete_webhook(wid):
    get_db().execute('DELETE FROM webhooks WHERE id=?', (wid,))
    get_db().commit()
    return jsonify({'ok': True})


@app.route('/api/admin/webhooks/logs', methods=['GET'])
@admin_required
def admin_webhook_logs():
    logs = get_db().execute(
        'SELECT l.*, w.name as webhook_name FROM webhook_logs l LEFT JOIN webhooks w ON l.webhook_id=w.id ORDER BY l.created_at DESC LIMIT 100'
    ).fetchall()
    return jsonify([dict(l) for l in logs])


@app.route('/api/admin/webhooks/test', methods=['POST'])
@admin_required
def admin_test_webhook():
    data = request.get_json()
    try:
        resp = requests.post(data['url'], json={'event':'test','message':'SmartCS webhook test','timestamp':datetime.now().isoformat()}, timeout=10)
        return jsonify({'ok': 200 <= resp.status_code < 300, 'status_code': resp.status_code, 'body': resp.text[:300]})
    except Exception as e:
        return jsonify({'ok': False, 'error': str(e)})

# ====== End Phase 4 ======

# ====== Phase 6 API: External Adapter Management ======

@app.route('/api/admin/external-adapters', methods=['GET'])
@admin_required
def admin_list_external_adapters():
    adapters = get_db().execute("SELECT * FROM external_adapters ORDER BY platform, created_at").fetchall()
    return jsonify([dict(a) for a in adapters])


@app.route('/api/admin/external-adapters', methods=['POST'])
@admin_required
def admin_create_external_adapter():
    data = request.get_json()
    name = data.get('name', '').strip()
    adapter_type = data.get('adapter_type', 'defect')
    platform = data.get('platform', '').strip()
    if not name or not platform:
        return jsonify({'error': '请填写名称和平台'}), 400
    if platform not in ('jira', 'zentao', 'github', 'gitlab', 'custom'):
        return jsonify({'error': '不支持的平台类型'}), 400
    if adapter_type not in ('defect', 'order_query', 'knowledge', 'custom'):
        return jsonify({'error': '不支持的适配器类型'}), 400
    aid = 'ext-' + uuid.uuid4().hex[:12]
    config = data.get('config', {})
    get_db().execute(
        "INSERT INTO external_adapters(id,name,adapter_type,platform,enabled,config) VALUES(?,?,?,?,?,?)",
        (aid, name, adapter_type, platform, data.get('enabled', 1),
         json.dumps(config, ensure_ascii=False))
    )
    get_db().commit()
    return jsonify({'ok': True, 'id': aid})


@app.route('/api/admin/external-adapters/<eid>', methods=['PUT'])
@admin_required
def admin_update_external_adapter(eid):
    data = request.get_json()
    existing = get_db().execute("SELECT id FROM external_adapters WHERE id=?", (eid,)).fetchone()
    if not existing:
        return jsonify({'error': '适配器不存在'}), 404
    fields = []
    vals = []
    for k in ('name', 'adapter_type', 'platform', 'enabled'):
        if k in data:
            fields.append(k + '=?')
            vals.append(data[k])
    if 'config' in data:
        fields.append('config=?')
        vals.append(json.dumps(data['config'], ensure_ascii=False))
    if fields:
        vals.append(eid)
        get_db().execute(
            'UPDATE external_adapters SET ' + ','.join(fields) + ', updated_at=datetime("now","localtime") WHERE id=?',
            vals
        )
        get_db().commit()
    return jsonify({'ok': True})


@app.route('/api/admin/external-adapters/<eid>', methods=['DELETE'])
@admin_required
def admin_delete_external_adapter(eid):
    get_db().execute("DELETE FROM external_adapters WHERE id=?", (eid,))
    get_db().commit()
    return jsonify({'ok': True})


@app.route('/api/admin/external-adapters/<eid>/test', methods=['POST'])
@admin_required
def admin_test_external_adapter(eid):
    row = get_db().execute("SELECT platform, config FROM external_adapters WHERE id=?", (eid,)).fetchone()
    if not row:
        return jsonify({'error': '适配器不存在'}), 404
    config = json.loads(row['config'] or '{}')
    cls = _ADAPTER_REGISTRY.get(row['platform'])
    if not cls:
        return jsonify({'ok': False, 'message': '不支持的平台类型'}), 400
    adapter = cls(config)
    ok, msg = adapter.validate_config()
    if ok:
        return jsonify({'ok': True, 'message': '配置验证通过 ✅'})
    return jsonify({'ok': False, 'message': msg})


# ====== 缺陷操作（工程师端） ======

@app.route('/api/agent/tickets/<ticket_id>/external-links', methods=['GET'])
@logged_in_required
def agent_ticket_external_links(ticket_id):
    links = get_db().execute(
        "SELECT * FROM ticket_external_links WHERE ticket_id=? ORDER BY created_at DESC",
        (ticket_id,)
    ).fetchall()
    return jsonify([dict(l) for l in links])


@app.route('/api/agent/tickets/<ticket_id>/create-defect', methods=['POST'])
@logged_in_required
def agent_create_defect(ticket_id):
    data = request.get_json()
    adapter_id = data.get('adapter_id', '')
    title = data.get('title', '').strip()
    description = data.get('description', '').strip()
    priority = data.get('priority', 'normal')

    if not adapter_id or not title:
        return jsonify({'error': '请填写适配器和标题'}), 400

    # Verify ticket exists
    tk = get_db().execute("SELECT id, issue_description, ticket_number FROM service_tickets WHERE id=?", (ticket_id,)).fetchone()
    if not tk:
        return jsonify({'error': '工单不存在'}), 404

    adapter = get_external_adapter(adapter_id)
    if not adapter:
        return jsonify({'error': '适配器不存在或已停用'}), 400

    if not description:
        description = tk['issue_description'] or title

    result = adapter.create_ticket(title, description, priority, session.get('agent_name', ''))
    if not result.get('success'):
        return jsonify({'error': result.get('error', '创建失败')}), 500

    # Save external link
    link_id = 'extlink-' + uuid.uuid4().hex[:12]
    get_db().execute(
        "INSERT INTO ticket_external_links(id,ticket_id,external_system,external_id,external_url,link_type) "
        "VALUES(?,?,?,?,?,?)",
        (link_id, ticket_id, adapter_id, result['external_id'],
         result.get('external_url', ''), 'defect')
    )

    # Add system message to conversation
    conv = get_db().execute("SELECT conversation_id FROM service_tickets WHERE id=?", (ticket_id,)).fetchone()
    if conv:
        sys_msg = f'✅ 已提交缺陷 [{result["external_id"]}]({result.get("external_url", "")}) 到外部系统'
        get_db().execute(
            "INSERT INTO messages(id,conversation_id,role,content) VALUES(?,?,?,?)",
            (gen_id('msg-'), conv['conversation_id'], 'system', sys_msg)
        )
        get_db().execute(
            "UPDATE conversations SET updated_at=datetime('now','localtime') WHERE id=?",
            (conv['conversation_id'],)
        )

    get_db().commit()
    return jsonify({
        'ok': True,
        'external_id': result['external_id'],
        'external_url': result.get('external_url', ''),
        'link_id': link_id
    })


@app.route('/api/agent/tickets/<ticket_id>/query-defect', methods=['POST'])
@logged_in_required
def agent_query_defect(ticket_id):
    data = request.get_json()
    adapter_id = data.get('adapter_id', '')
    external_id = data.get('external_id', '')

    if not adapter_id or not external_id:
        # Try to find from existing links
        link = get_db().execute(
            "SELECT external_system, external_id FROM ticket_external_links "
            "WHERE ticket_id=? ORDER BY created_at DESC LIMIT 1",
            (ticket_id,)
        ).fetchone()
        if not link:
            return jsonify({'error': '未找到关联的缺陷，请指定 adapter_id 和 external_id'}), 400
        adapter_id = link['external_system']
        external_id = link['external_id']

    adapter = get_external_adapter(adapter_id)
    if not adapter:
        return jsonify({'error': '适配器不存在'}), 400

    result = adapter.get_status(external_id)
    return jsonify(result)


# ====== 外部链接管理（管理员） ======


# ====== SSL 证书管理 ======
@app.route('/api/admin/ssl/cert', methods=['GET'])
@admin_required
def admin_ssl_info():
    cert_path = '/etc/ssl/certs/smartcs.crt'
    key_path = '/etc/ssl/private/smartcs.key'
    info = {'cert_path': cert_path, 'key_path': key_path}
    if os.path.exists(cert_path):
        try:
            result = subprocess.run(['openssl', 'x509', '-in', cert_path, '-noout', '-dates', '-subject', '-issuer'], capture_output=True, text=True, timeout=5)
            for line in result.stdout.strip().split(chr(10)):
                if line.startswith('notBefore='): info['not_before'] = line[10:]
                elif line.startswith('notAfter='): info['not_after'] = line[9:]
                elif line.startswith('subject='): info['subject'] = line[8:]
                elif line.startswith('issuer='): info['issuer'] = line[7:]
        except: pass
        info['cert_exists'] = True
        info['cert_size'] = os.path.getsize(cert_path)
    else:
        info['cert_exists'] = False
    info['key_exists'] = os.path.exists(key_path)
    return jsonify(info)

@app.route('/api/admin/ssl/cert', methods=['POST'])
@admin_required
def admin_ssl_upload():
    data = request.get_json() or {}
    cert_content = data.get('certificate', '')
    key_content = data.get('private_key', '')
    if not cert_content or not key_content:
        return jsonify({'error': '证书和私钥都不能为空'}), 400
    try:
        with tempfile.NamedTemporaryFile(mode='w', suffix='.crt', delete=False) as f:
            f.write(cert_content); cert_tmp = f.name
        with tempfile.NamedTemporaryFile(mode='w', suffix='.key', delete=False) as f:
            f.write(key_content); key_tmp = f.name
        result = subprocess.run(['openssl', 'x509', '-in', cert_tmp, '-noout'], capture_output=True, text=True)
        if result.returncode != 0:
            os.unlink(cert_tmp); os.unlink(key_tmp)
            return jsonify({'error': '证书格式无效'}), 400
        os.chmod(cert_tmp, 0o644); os.chmod(key_tmp, 0o600)
        os.rename(cert_tmp, '/etc/ssl/certs/smartcs.crt')
        os.rename(key_tmp, '/etc/ssl/private/smartcs.key')
        subprocess.run(['systemctl', 'reload', 'nginx'], timeout=10)
        log_audit('', 'admin.ssl_update', session.get('agent_id',''), session.get('agent_name',''), {})
        return jsonify({'ok': True, 'message': 'SSL 证书已更新'})
    except Exception as e:
        return jsonify({'error': str(e)[:100]}), 500

# ====== 知识库 CRUD ======
@app.route('/api/admin/knowledge', methods=['GET'])
@admin_required
def admin_knowledge_list():
    q = request.args.get('q', '').strip().lower()
    system_id = request.args.get('system_id', '').strip()
    scenario = request.args.get('scenario', '').strip()
    category = request.args.get('category', '').strip()
    db = get_db()
    rows = db.execute('SELECT id, filename, word_count, uploaded_by, created_by, updated_by, tags, title, status, created_at, updated_at FROM knowledge_files ORDER BY created_at DESC').fetchall()
    results = []
    for r in rows:
        fp = os.path.join(KNOWLEDGE_DIR, r['filename'])
        snippet = ''
        if os.path.exists(fp):
            with open(fp, 'r', encoding='utf-8', errors='replace') as f:
                snippet = f.read()[:200]
        item = dict(r)
        item['snippet'] = snippet
        # Parse tags
        tags_str = item.get('tags', '{}') or '{}'
        try:
            tags = json.loads(tags_str) if isinstance(tags_str, str) else tags_str
        except:
            tags = {}
        if not isinstance(tags, dict):
            tags = {}
        item['tags'] = tags
        # Count history versions
        hcount = db.execute("SELECT COUNT(*) FROM knowledge_history WHERE knowledge_id=?", (item['id'],)).fetchone()[0]
        item['version_count'] = hcount
        # Filter by tags
        if system_id:
            sys_ids = tags.get('system_ids', [])
            if system_id not in sys_ids:
                continue
        if scenario and tags.get('scenario', '') != scenario:
            continue
        if category and tags.get('category', '') != category:
            continue
        if not q or q in r['filename'].lower() or q in snippet.lower():
            results.append(item)
    return jsonify(results)

@app.route('/api/admin/knowledge', methods=['POST'])
@admin_required
def admin_knowledge_create():
    data = request.get_json() or {}
    file_name = (data.get('filename') or data.get('file') or '').strip()
    content = (data.get('content') or data.get('snippet') or '').strip()
    if not file_name:
        return jsonify({'error': '文件名不能为空'}), 400
    if not file_name.endswith('.md'):
        file_name += '.md'
    if not content:
        return jsonify({'error': '内容不能为空'}), 400
    fp = os.path.join(KNOWLEDGE_DIR, file_name)
    db = get_db()
    existing = db.execute('SELECT id FROM knowledge_files WHERE filename=?', (file_name,)).fetchone()
    if existing:
        return jsonify({'error': '该文件名已存在'}), 409
    with open(fp, 'w', encoding='utf-8') as f:
        f.write(content)
    word_count = len(content)
    # Parse tags from request
    tags_dict = data.get('tags', {})
    if not isinstance(tags_dict, dict):
        tags_dict = {}
    if 'system_ids' not in tags_dict:
        tags_dict['system_ids'] = []
    if 'scenario' not in tags_dict:
        tags_dict['scenario'] = ''
    if 'category' not in tags_dict:
        tags_dict['category'] = ''
    if 'custom' not in tags_dict:
        tags_dict['custom'] = []
    tags_json = json.dumps(tags_dict, ensure_ascii=False)
    agent_name = session.get('agent_name','admin')
    kid = 'kf-' + uuid.uuid4().hex[:12]
    db.execute('INSERT INTO knowledge_files(id, filename, word_count, uploaded_by, created_by, updated_by, tags) VALUES(?,?,?,?,?,?,?)',
               (kid, file_name, word_count, agent_name, agent_name, agent_name, tags_json))
    # Add initial history record
    db.execute(
        "INSERT INTO knowledge_history(id,knowledge_id,updated_by,updated_by_name,change_summary,old_content,new_content) "
        "VALUES(?,?,?,?,?,?,?)",
        (gen_id('kh-'), kid, session.get('agent_id',''), agent_name, '创建知识条目', '', content[:500])
    )
    db.commit()
    log_audit('', 'admin.kb_create', session.get('agent_id',''), session.get('agent_name',''), {'file': file_name})
    return jsonify({'ok': True})

@app.route('/api/admin/knowledge/<kid>', methods=['PUT'])
@admin_required
def admin_knowledge_update(kid):
    data = request.get_json() or {}
    file_name = (data.get('filename') or data.get('file') or '').strip()
    content = data.get('content') or data.get('snippet') or ''
    db = get_db()
    row = db.execute('SELECT * FROM knowledge_files WHERE id=?', (kid,)).fetchone()
    if not row:
        return jsonify({'error': '条目不存在'}), 404
    old_fn = row['filename']
    target_fn = file_name if file_name else old_fn
    if not target_fn.endswith('.md'):
        target_fn += '.md'
    old_fp = os.path.join(KNOWLEDGE_DIR, old_fn)
    new_fp = os.path.join(KNOWLEDGE_DIR, target_fn)
    old_content = ''
    if os.path.exists(old_fp):
        with open(old_fp, 'r', encoding='utf-8', errors='replace') as f:
            old_content = f.read()
    if content:
        with open(new_fp, 'w', encoding='utf-8') as f:
            f.write(content)
    if target_fn != old_fn and os.path.exists(old_fp) and old_fp != new_fp:
        shutil.move(old_fp, new_fp)
    word_count = len(content) if content else len(old_content)
    agent_name = session.get('agent_name','admin')
    updates = ['filename=?', 'word_count=?', 'updated_by=?', "updated_at=datetime('now','localtime')"]
    params = [target_fn, word_count, agent_name]
    # Update tags if provided
    tags_dict = data.get('tags')
    if tags_dict is not None and isinstance(tags_dict, dict):
        tags_json = json.dumps(tags_dict, ensure_ascii=False)
        updates.append('tags=?')
        params.append(tags_json)
    params.append(kid)
    db.execute(f'UPDATE knowledge_files SET {chr(44).join(updates)} WHERE id=?', params)
    # Add history record if content changed
    if content and old_content != content:
        db.execute(
            "INSERT INTO knowledge_history(id,knowledge_id,updated_by,updated_by_name,change_summary,old_content,new_content) "
            "VALUES(?,?,?,?,?,?,?)",
            (gen_id('kh-'), kid, session.get('agent_id',''), agent_name, '编辑知识条目', old_content[:500], content[:500])
        )
    db.commit()
    log_audit('', 'admin.kb_update', session.get('agent_id',''), session.get('agent_name',''), {'kb_id': kid})
    return jsonify({'ok': True})

@app.route('/api/admin/knowledge/<kid>/history', methods=['GET'])
@admin_required
def admin_knowledge_history(kid):
    """Get version history for a knowledge item."""
    rows = get_db().execute(
        "SELECT * FROM knowledge_history WHERE knowledge_id=? ORDER BY updated_at DESC",
        (kid,)
    ).fetchall()
    return jsonify([dict(r) for r in rows])


@app.route('/api/admin/knowledge/<kid>', methods=['DELETE'])
@admin_required
def admin_knowledge_delete(kid):
    kb_id = kid
    if not kb_id:
        return jsonify({'error': '缺少条目ID'}), 400
    db = get_db()
    row = db.execute('SELECT filename FROM knowledge_files WHERE id=?', (kb_id,)).fetchone()
    if not row:
        return jsonify({'error': '条目不存在'}), 404
    fp = os.path.join(KNOWLEDGE_DIR, row['filename'])
    if os.path.exists(fp):
        os.remove(fp)
    db.execute('DELETE FROM knowledge_files WHERE id=?', (kb_id,))
    db.commit()
    log_audit('', 'admin.kb_delete', session.get('agent_id',''), session.get('agent_name',''), {'file': row['filename']})
    return jsonify({'ok': True})


# Keep old DELETE route for backward compatibility
@app.route('/api/admin/knowledge', methods=['DELETE'])
@admin_required
def admin_knowledge_delete_old():
    data = request.get_json() or {}
    kb_id = data.get('id', '')
    if not kb_id:
        return jsonify({'error': '缺少条目ID'}), 400
    db = get_db()
    row = db.execute('SELECT filename FROM knowledge_files WHERE id=?', (kb_id,)).fetchone()
    if not row:
        return jsonify({'error': '条目不存在'}), 404
    fp = os.path.join(KNOWLEDGE_DIR, row['filename'])
    if os.path.exists(fp):
        os.remove(fp)
    db.execute('DELETE FROM knowledge_files WHERE id=?', (kb_id,))
    db.commit()
    log_audit('', 'admin.kb_delete', session.get('agent_id',''), session.get('agent_name',''), {'file': row['filename']})
    return jsonify({'ok': True})

# ====== Knowledge Audit System (知识沉淀审核) ======

@app.route('/api/agent/tickets/knowledge', methods=['POST'])
@logged_in_required
def agent_knowledge_submit():
    """Engineer submits a knowledge document for review."""
    data = request.get_json() or {}
    title = (data.get('title') or '').strip()
    content = (data.get('content') or '').strip()
    if not title or not content:
        return jsonify({'error': '标题和内容不能为空'}), 400
    
    # Create filename from title
    safe_title = re.sub(r'[^\w\u4e00-\u9fff-]', '_', title)[:50]
    if not safe_title:
        safe_title = uuid.uuid4().hex[:8]
    filename = f"{safe_title}_{uuid.uuid4().hex[:6]}.md"
    
    fp = os.path.join(KNOWLEDGE_DIR, filename)
    with open(fp, 'w', encoding='utf-8') as f:
        f.write(content)
    
    word_count = len(content)
    kid = 'kf-' + uuid.uuid4().hex[:12]
    db = get_db()
    agent_name = session.get('agent_name','')
    db.execute(
        "INSERT INTO knowledge_files(id,filename,word_count,uploaded_by,title,status,submitted_by,created_by,updated_by,tags) "
        "VALUES(?,?,?,?,?,'pending',?,?,?,'{}')",
        (kid, filename, word_count, agent_name, title, agent_name, agent_name, agent_name)
    )
    # Add initial history record
    db.execute(
        "INSERT INTO knowledge_history(id,knowledge_id,updated_by,updated_by_name,change_summary,old_content,new_content) "
        "VALUES(?,?,?,?,?,?,?)",
        (gen_id('kh-'), kid, session.get('agent_id',''), agent_name, '提交知识条目（待审核）', '', content[:500])
    )
    db.commit()
    
    log_audit('', 'knowledge.submit', session.get('agent_id',''), session.get('agent_name',''),
              {'knowledge_id': kid, 'title': title})
    return jsonify({'ok': True, 'knowledge_id': kid})


@app.route('/api/agent/knowledge/mine', methods=['GET'])
@logged_in_required
def agent_knowledge_mine():
    """Engineer views their own knowledge submissions."""
    rows = get_db().execute(
        "SELECT id, filename, title, word_count, status, review_notes, created_at "
        "FROM knowledge_files WHERE submitted_by=? ORDER BY created_at DESC",
        (session.get('agent_name',''),)
    ).fetchall()
    result = []
    for r in rows:
        item = dict(r)
        fp = os.path.join(KNOWLEDGE_DIR, r['filename'])
        if os.path.exists(fp):
            with open(fp, 'r', encoding='utf-8', errors='replace') as f:
                item['snippet'] = f.read()[:200]
        else:
            item['snippet'] = ''
        result.append(item)
    return jsonify(result)


@app.route('/api/agent/knowledge/<kid>', methods=['GET'])
@logged_in_required
def agent_knowledge_detail(kid):
    """Get knowledge detail by id."""
    row = get_db().execute("SELECT * FROM knowledge_files WHERE id=?", (kid,)).fetchone()
    if not row:
        return jsonify({'error': '知识条目不存在'}), 404
    item = dict(row)
    fp = os.path.join(KNOWLEDGE_DIR, row['filename'])
    if os.path.exists(fp):
        with open(fp, 'r', encoding='utf-8', errors='replace') as f:
            item['content'] = f.read()
    else:
        item['content'] = ''
    return jsonify(item)


@app.route('/api/agent/knowledge/<kid>', methods=['PUT'])
@logged_in_required
def agent_knowledge_update(kid):
    """Update knowledge content (generates history entry)."""
    data = request.get_json() or {}
    db = get_db()
    row = db.execute("SELECT * FROM knowledge_files WHERE id=?", (kid,)).fetchone()
    if not row:
        return jsonify({'error': '知识条目不存在'}), 404
    
    new_content = data.get('content', '')
    agent_name = session.get('agent_name','')
    if new_content:
        fp = os.path.join(KNOWLEDGE_DIR, row['filename'])
        old_content = ''
        if os.path.exists(fp):
            with open(fp, 'r', encoding='utf-8', errors='replace') as f:
                old_content = f.read()
        with open(fp, 'w', encoding='utf-8') as f:
            f.write(new_content)
        # Save history
        if old_content != new_content:
            db.execute(
                "INSERT INTO knowledge_history(id,knowledge_id,updated_by,updated_by_name,change_summary,old_content,new_content) "
                "VALUES(?,?,?,?,?,?,?)",
                (gen_id('kh-'), kid, session.get('agent_id',''), agent_name, '更新内容', old_content[:500], new_content[:500])
            )
        db.execute(
            "UPDATE knowledge_files SET word_count=?,updated_by=?,updated_at=datetime('now','localtime') WHERE id=?",
            (len(new_content), agent_name, kid)
        )
    
    db.commit()
    return jsonify({'ok': True})


@app.route('/api/agent/knowledge/<kid>/history', methods=['GET'])
@logged_in_required
def agent_knowledge_history(kid):
    """Get knowledge update history."""
    rows = get_db().execute(
        "SELECT * FROM knowledge_history WHERE knowledge_id=? ORDER BY updated_at DESC",
        (kid,)
    ).fetchall()
    return jsonify([dict(r) for r in rows])


@app.route('/api/admin/knowledge/<kid>/approve', methods=['POST'])
@admin_required
def admin_knowledge_approve(kid):
    """Admin approves a pending knowledge document."""
    db = get_db()
    row = db.execute("SELECT * FROM knowledge_files WHERE id=?", (kid,)).fetchone()
    if not row:
        return jsonify({'error': '知识条目不存在'}), 404
    db.execute(
        "UPDATE knowledge_files SET status='approved',updated_by=?,updated_at=datetime('now','localtime') WHERE id=?",
        (session.get('agent_name','admin'), kid)
    )
    db.execute(
        "INSERT INTO knowledge_history(id,knowledge_id,updated_by,updated_by_name,change_summary,old_content,new_content) "
        "VALUES(?,?,?,?,?,?,?)",
        (gen_id('kh-'), kid, session.get('agent_id',''), session.get('agent_name','admin'), '审核通过', '', '')
    )
    db.commit()
    log_audit('', 'knowledge.approve', session.get('agent_id',''), session.get('agent_name',''),
              {'knowledge_id': kid, 'title': row['title']})
    return jsonify({'ok': True})


@app.route('/api/admin/knowledge/<kid>/reject', methods=['POST'])
@admin_required
def admin_knowledge_reject(kid):
    """Admin rejects a pending knowledge document with reason."""
    data = request.get_json() or {}
    reason = data.get('reason', '').strip()
    db = get_db()
    row = db.execute("SELECT * FROM knowledge_files WHERE id=?", (kid,)).fetchone()
    if not row:
        return jsonify({'error': '知识条目不存在'}), 404
    db.execute(
        "UPDATE knowledge_files SET status='rejected',review_notes=?,updated_by=?,updated_at=datetime('now','localtime') WHERE id=?",
        (reason, session.get('agent_name','admin'), kid)
    )
    db.execute(
        "INSERT INTO knowledge_history(id,knowledge_id,updated_by,updated_by_name,change_summary,old_content,new_content) "
        "VALUES(?,?,?,?,?,?,?)",
        (gen_id('kh-'), kid, session.get('agent_id',''), session.get('agent_name','admin'), f'驳回: {reason[:200]}', '', '')
    )
    db.commit()
    log_audit('', 'knowledge.reject', session.get('agent_id',''), session.get('agent_name',''),
              {'knowledge_id': kid, 'title': row['title'], 'reason': reason})
    return jsonify({'ok': True})


@app.route('/api/agent/tickets/escalate', methods=['POST'])
@logged_in_required
def agent_escalate_ticket():
    """Escalate a processing ticket to a higher level engineer."""
    data = request.get_json() or {}
    ticket_id = data.get('ticket_id', '')
    target_level = data.get('target_level', 0)
    reason = data.get('reason', '').strip()
    
    if not ticket_id or not target_level:
        return jsonify({'error': '请填写工单ID和目标级别'}), 400
    
    db = get_db()
    t = db.execute("SELECT * FROM service_tickets WHERE id=?", (ticket_id,)).fetchone()
    if not t:
        return jsonify({'error': '工单不存在'}), 404
    
    # Get current agent's level
    prof = db.execute("SELECT agent_level FROM agent_profiles WHERE agent_id=?", (session['agent_id'],)).fetchone()
    current_level = prof['agent_level'] if prof else 1
    
    if current_level >= target_level:
        return jsonify({'error': f'当前级别(L{current_level})必须低于目标级别(L{target_level})'}), 400
    
    if t['status'] != 'processing':
        return jsonify({'error': '仅处理中的工单可以升级'}), 400
    
    # Update ticket level
    db.execute(
        "UPDATE service_tickets SET level=?,updated_at=datetime('now','localtime') WHERE id=?",
        (target_level, ticket_id)
    )
    
    # Add system message to conversation
    conv = db.execute("SELECT conversation_id FROM service_tickets WHERE id=?", (ticket_id,)).fetchone()
    if conv:
        agent_name = session.get('agent_name', '工程师')
        msg_text = f'🚨 {agent_name} 已将该工单升级至L{target_level}级处理'
        if reason:
            msg_text += f'（原因：{reason}）'
        db.execute(
            "INSERT INTO messages(id,conversation_id,role,content) VALUES(?,?,?,?)",
            (gen_id('msg-'), conv['conversation_id'], 'system', msg_text)
        )
        db.execute(
            "UPDATE conversations SET updated_at=datetime('now','localtime') WHERE id=?",
            (conv['conversation_id'],)
        )
    
    db.commit()
    log_audit(ticket_id, 'ticket.escalate', session.get('agent_id',''), session.get('agent_name',''),
              {'from_level': current_level, 'to_level': target_level, 'reason': reason})
    return jsonify({'ok': True, 'message': f'工单已升级至L{target_level}级处理'})


@app.route('/api/admin/knowledge/pending', methods=['GET'])
@admin_required
def admin_knowledge_pending():
    """Admin gets list of pending knowledge docs with submitter display name."""
    rows = get_db().execute(
        "SELECT kf.*, COALESCE(ap.display_name, kf.submitted_by) as submitter_display_name "
        "FROM knowledge_files kf "
        "LEFT JOIN agents ag ON kf.submitted_by = ag.name "
        "LEFT JOIN agent_profiles ap ON ag.id = ap.agent_id "
        "WHERE kf.status='pending' ORDER BY kf.created_at DESC"
    ).fetchall()
    result = []
    for r in rows:
        item = dict(r)
        fp = os.path.join(KNOWLEDGE_DIR, r['filename'])
        if os.path.exists(fp):
            with open(fp, 'r', encoding='utf-8', errors='replace') as f:
                item['content'] = f.read()
        else:
            item['content'] = ''
        result.append(item)
    return jsonify(result)


@app.route('/api/knowledge/my-submissions', methods=['GET'])
def knowledge_my_submissions():
    """Engineer views their own submissions from non-agent context."""
    if not session.get('agent_id'):
        return jsonify({'error': '未登录'}), 401
    rows = get_db().execute(
        "SELECT kf.id, kf.filename, kf.title, kf.word_count, kf.status, kf.review_notes, "
        "kf.created_at, COALESCE(ap.display_name, kf.submitted_by) as submitter_display_name "
        "FROM knowledge_files kf "
        "LEFT JOIN agents ag ON kf.submitted_by = ag.name "
        "LEFT JOIN agent_profiles ap ON ag.id = ap.agent_id "
        "WHERE kf.submitted_by=? ORDER BY kf.created_at DESC",
        (session.get('agent_name',''),)
    ).fetchall()
    result = []
    for r in rows:
        item = dict(r)
        fp = os.path.join(KNOWLEDGE_DIR, r['filename'])
        if os.path.exists(fp):
            with open(fp, 'r', encoding='utf-8', errors='replace') as f:
                item['snippet'] = f.read()[:200]
        else:
            item['snippet'] = ''
        result.append(item)
    return jsonify(result)


# ====== End Knowledge Audit System ======


@app.route('/api/admin/external-links', methods=['POST'])
@admin_required
def admin_create_external_link():
    data = request.get_json()
    ticket_id = data.get('ticket_id', '')
    if not ticket_id:
        return jsonify({'error': '请填写工单 ID'}), 400
    link_id = 'extlink-' + uuid.uuid4().hex[:12]
    get_db().execute(
        "INSERT INTO ticket_external_links(id,ticket_id,external_system,external_id,external_url,link_type,sync_status) "
        "VALUES(?,?,?,?,?,?,?)",
        (link_id, ticket_id, data.get('external_system', ''),
         data.get('external_id', ''), data.get('external_url', ''),
         data.get('link_type', 'defect'), data.get('sync_status', 'active'))
    )
    get_db().commit()
    return jsonify({'ok': True, 'id': link_id})


@app.route('/api/admin/external-links/<link_id>', methods=['DELETE'])
@admin_required
def admin_delete_external_link(link_id):
    get_db().execute("DELETE FROM ticket_external_links WHERE id=?", (link_id,))
    get_db().commit()
    return jsonify({'ok': True})


@app.route('/api/admin/external-links', methods=['GET'])
@admin_required
def admin_list_external_links():
    ticket_id = request.args.get('ticket_id', '')
    where = []
    params = []
    if ticket_id:
        where.append("ticket_id=?")
        params.append(ticket_id)
    wsql = ' AND '.join(where) if where else '1=1'
    links = get_db().execute(
        f"SELECT * FROM ticket_external_links WHERE {wsql} ORDER BY created_at DESC",
        params
    ).fetchall()
    return jsonify([dict(l) for l in links])


@app.route('/icon-192.svg')
def icon_svg():
    return send_from_directory(os.path.join(BASE_DIR, 'static'), 'icon-192.svg')


# ====== Phase 5: 统一身份认证（LDAP/OIDC） ======

class AuthProvider:
    """认证提供者基类"""
    def __init__(self, config: dict):
        self.config = config

    def authenticate(self, username: str, password: str) -> dict:
        return {'success': False, 'error': 'not_implemented'}

    def get_login_html(self) -> str:
        return ''

    def get_config_fields(self) -> list:
        return []

    def test_connection(self) -> tuple:
        return True, ''


class LdapProvider(AuthProvider):
    """LDAP/AD 认证提供者"""
    def __init__(self, config):
        super().__init__(config)
        self.server = config.get('server', '')
        self.port = config.get('port', 389)
        self.base_dn = config.get('base_dn', '')
        self.bind_dn = config.get('bind_dn', '')
        self.bind_password = config.get('bind_password', '')
        self.user_filter = config.get('user_filter', '(sAMAccountName={username})')
        self.use_ssl = config.get('use_ssl', False)
        self.field_mapping = config.get('field_mapping', {
            'username': 'sAMAccountName', 'email': 'mail', 'display_name': 'displayName'
        })

    def authenticate(self, username: str, password: str) -> dict:
        try:
            import ldap3
            server = ldap3.Server(self.server, port=self.port, use_ssl=self.use_ssl, get_info=ldap3.ALL)
            conn = ldap3.Connection(server, self.bind_dn, self.bind_password, auto_bind=True)
            filter_str = self.user_filter.replace('{username}', username)
            conn.search(self.base_dn, filter_str, attributes=['*'])
            if len(conn.entries) == 0:
                return {'success': False, 'error': '用户不存在'}
            entry = conn.entries[0]
            try:
                user_dn = str(entry.entry_dn)
                user_conn = ldap3.Connection(server, user_dn, password, auto_bind=True)
                user_conn.unbind()
                return {
                    'success': True,
                    'user_id': str(getattr(entry, self.field_mapping.get('username', 'sAMAccountName'), '')),
                    'username': username,
                    'email': str(getattr(entry, self.field_mapping.get('email', 'mail'), '')),
                    'display_name': str(getattr(entry, self.field_mapping.get('display_name', 'displayName'), username)),
                }
            except:
                return {'success': False, 'error': '密码错误'}
        except ImportError:
            return {'success': False, 'error': 'LDAP 支持未安装（需要 ldap3 库）'}
        except Exception as e:
            return {'success': False, 'error': f'LDAP 连接失败: {str(e)}'}

    def test_connection(self) -> tuple:
        try:
            import ldap3
            server = ldap3.Server(self.server, port=self.port, use_ssl=self.use_ssl)
            conn = ldap3.Connection(server, self.bind_dn, self.bind_password, auto_bind=True)
            conn.unbind()
            return True, '连接成功'
        except ImportError:
            return False, 'LDAP 支持未安装（需要 ldap3 库）'
        except Exception as e:
            return False, f'连接失败: {str(e)}'

    def get_config_fields(self):
        return [
            {'key': 'server', 'label': 'LDAP 服务器地址', 'type': 'text', 'placeholder': 'ldap.company.com'},
            {'key': 'port', 'label': '端口', 'type': 'number', 'default': 389},
            {'key': 'base_dn', 'label': 'Base DN', 'type': 'text', 'placeholder': 'dc=company,dc=com'},
            {'key': 'bind_dn', 'label': 'Bind DN（管理员账号）', 'type': 'text'},
            {'key': 'bind_password', 'label': 'Bind 密码', 'type': 'password'},
            {'key': 'user_filter', 'label': '用户过滤条件', 'type': 'text', 'default': '(sAMAccountName={username})'},
            {'key': 'use_ssl', 'label': '使用 SSL', 'type': 'checkbox', 'default': False},
        ]


class OidcProvider(AuthProvider):
    """OIDC（OpenID Connect）认证提供者"""
    def __init__(self, config):
        super().__init__(config)
        self.client_id = config.get('client_id', '')
        self.client_secret = config.get('client_secret', '')
        self.discovery_url = config.get('discovery_url', '')
        self.redirect_uri = config.get('redirect_uri', '')
        self.authorization_endpoint = config.get('authorization_endpoint', '')
        self.token_endpoint = config.get('token_endpoint', '')
        self.userinfo_endpoint = config.get('userinfo_endpoint', '')
        self.scopes = config.get('scopes', 'openid profile email')

    def authenticate(self, username: str = '', password: str = '') -> dict:
        return {'success': False, 'error': 'OIDC 需要使用 SSO 登录页面'}

    def get_login_html(self) -> str:
        return f'''<div style="text-align:center;margin:16px 0">
            <a href="/api/auth/oidc/login/{self.config.get('_id','')}"
               style="display:inline-block;padding:12px 32px;background:var(--primary-color);color:#fff;
                      border-radius:8px;text-decoration:none;font-size:14px;">
               🔑 SSO 企业账号登录
            </a>
        </div>'''

    def get_authorization_url(self):
        state = uuid.uuid4().hex[:16]
        params = {
            'response_type': 'code',
            'client_id': self.client_id,
            'redirect_uri': self.redirect_uri,
            'scope': self.scopes,
            'state': state,
        }
        auth_url = self.authorization_endpoint or (self.discovery_url.rstrip('/') + '/auth')
        return f'{auth_url}?{urllib.parse.urlencode(params)}', state

    def exchange_code(self, code: str) -> dict:
        token_url = self.token_endpoint or (self.discovery_url.rstrip('/') + '/token')
        resp = requests.post(token_url, data={
            'grant_type': 'authorization_code',
            'code': code,
            'redirect_uri': self.redirect_uri,
            'client_id': self.client_id,
            'client_secret': self.client_secret,
        })
        return resp.json()

    def get_userinfo(self, access_token: str) -> dict:
        userinfo_url = self.userinfo_endpoint or (self.discovery_url.rstrip('/') + '/userinfo')
        resp = requests.get(userinfo_url, headers={'Authorization': f'Bearer {access_token}'})
        return resp.json()

    def get_config_fields(self):
        return [
            {'key': 'client_id', 'label': 'Client ID', 'type': 'text'},
            {'key': 'client_secret', 'label': 'Client Secret', 'type': 'password'},
            {'key': 'discovery_url', 'label': 'Discovery URL', 'type': 'text',
             'placeholder': 'https://auth.company.com/.well-known/openid-configuration'},
            {'key': 'redirect_uri', 'label': '回调地址 (Redirect URI)', 'type': 'text',
             'placeholder': 'http://localhost:5000/agent/oidc/callback'},
            {'key': 'scopes', 'label': 'Scopes', 'type': 'text', 'default': 'openid profile email'},
        ]


_AUTH_PROVIDER_REGISTRY = {
    'ldap': LdapProvider,
    'oidc': OidcProvider,
}


def get_auth_provider(provider_id: str) -> AuthProvider:
    """根据 provider_id 获取认证提供者实例"""
    row = get_db().execute('SELECT * FROM auth_providers WHERE id=?', (provider_id,)).fetchone()
    if not row:
        return None
    config = json.loads(row['config']) if isinstance(row['config'], str) else row['config']
    config['_id'] = row['id']
    cls = _AUTH_PROVIDER_REGISTRY.get(row['provider_type'])
    if cls:
        return cls(config)
    return None


def get_enabled_providers() -> list:
    """获取所有启用的非本地认证提供者"""
    rows = get_db().execute('SELECT * FROM auth_providers WHERE enabled=1 ORDER BY provider_type').fetchall()
    return [dict(r) for r in rows]


# ====== 认证提供者管理 API（管理员） ======

@app.route('/api/admin/auth-providers', methods=['GET'])
@secadmin_required
def admin_list_auth_providers():
    providers = get_db().execute('SELECT * FROM auth_providers ORDER BY provider_type, created_at').fetchall()
    return jsonify([dict(p) for p in providers])


@app.route('/api/admin/auth-providers', methods=['POST'])
@secadmin_required
def admin_create_auth_provider():
    data = request.get_json()
    aid = 'auth-' + uuid.uuid4().hex[:12]
    ptype = data.get('provider_type', '')
    if ptype not in ('ldap', 'oidc'):
        return jsonify({'error': '不支持的认证提供者类型'}), 400
    get_db().execute(
        'INSERT INTO auth_providers(id,name,provider_type,enabled,config) VALUES(?,?,?,?,?)',
        (aid, data.get('name', ptype), ptype,
         data.get('enabled', 1), json.dumps(data.get('config', {}), ensure_ascii=False))
    )
    get_db().commit()
    log_audit('', 'admin.auth_provider.create', session.get('agent_id',''), session.get('agent_name',''),
              {'provider_id': aid, 'type': ptype, 'name': data.get('name', ptype)})
    return jsonify({'ok': True, 'id': aid})


@app.route('/api/admin/auth-providers/<pid>', methods=['PUT'])
@secadmin_required
def admin_update_auth_provider(pid):
    data = request.get_json()
    existing = get_db().execute('SELECT id FROM auth_providers WHERE id=?', (pid,)).fetchone()
    if not existing:
        return jsonify({'error': '认证提供者不存在'}), 404
    get_db().execute(
        'UPDATE auth_providers SET name=?, enabled=?, config=?, updated_at=datetime(\'now\',\'localtime\') WHERE id=?',
        (data.get('name', ''), data.get('enabled', 1),
         json.dumps(data.get('config', {}), ensure_ascii=False), pid)
    )
    get_db().commit()
    return jsonify({'ok': True})


@app.route('/api/admin/auth-providers/<pid>', methods=['DELETE'])
@secadmin_required
def admin_delete_auth_provider(pid):
    get_db().execute('DELETE FROM auth_identity_mappings WHERE provider_id=?', (pid,))
    get_db().execute('DELETE FROM auth_providers WHERE id=?', (pid,))
    get_db().commit()
    log_audit('', 'admin.auth_provider.delete', session.get('agent_id',''), session.get('agent_name',''),
              {'provider_id': pid})
    return jsonify({'ok': True})


@app.route('/api/admin/auth-providers/<pid>/test', methods=['POST'])
@secadmin_required
def admin_test_auth_provider(pid):
    provider = get_auth_provider(pid)
    if not provider:
        return jsonify({'error': '认证提供者不存在'}), 404
    ok, msg = provider.test_connection()
    return jsonify({'ok': ok, 'message': msg})


@app.route('/api/auth/providers-enabled', methods=['GET'])
def auth_providers_enabled():
    providers = get_enabled_providers()
    return jsonify(providers)


# ====== LDAP 登录 ======

@app.route('/api/auth/ldap/login', methods=['POST'])
def auth_ldap_login():
    """LDAP 统一登录接口"""
    data = request.get_json() or request.form
    username = data.get('username', '')
    password = data.get('password', '')
    provider_id = data.get('provider_id', '')
    
    if not username or not password or not provider_id:
        return jsonify({'error': '请填写用户名、密码和认证提供者'}), 400
    
    provider = get_auth_provider(provider_id)
    if not provider:
        return jsonify({'error': '认证提供者不存在'}), 404
    
    result = provider.authenticate(username, password)
    if not result['success']:
        return jsonify({'error': result.get('error', '认证失败')}), 401
    
    # 查找或创建本地映射
    mapping = get_db().execute(
        'SELECT * FROM auth_identity_mappings WHERE provider_id=? AND external_user_id=?',
        (provider_id, result['user_id'])
    ).fetchone()
    
    if mapping:
        smartcs_agent_id = mapping['smartcs_agent_id']
        get_db().execute('UPDATE auth_identity_mappings SET last_login_at=datetime(\'now\',\'localtime\') WHERE id=?',
                         (mapping['id'],))
    else:
        # 尝试按邮箱匹配已有IT工程师
        email = result.get('email', '')
        agent = None
        if email:
            agent = get_db().execute('SELECT id FROM agents WHERE email=?', (email,)).fetchone()
        if not agent:
            # 创建临时IT工程师账号
            aid = 'auth-agent-' + uuid.uuid4().hex[:8]
            display = result.get('display_name', result['username'])
            pwd_hash = hashlib.sha256(f'admin:{uuid.uuid4().hex[:12]}'.encode()).hexdigest()
            get_db().execute(
                'INSERT INTO agents(id,name,email,password_hash,role) VALUES(?,?,?,?,?)',
                (aid, display, email or f'{result["username"]}@ldap.local', pwd_hash, 'agent')
            )
            get_db().execute(
                'INSERT INTO agent_profiles(id,agent_id,display_name,department) VALUES(?,?,?,?)',
                ('ap-' + uuid.uuid4().hex[:8], aid, display, 'LDAP导入')
            )
            smartcs_agent_id = aid
        else:
            smartcs_agent_id = agent['id']
        
        mid = 'authmap-' + uuid.uuid4().hex[:12]
        get_db().execute(
            'INSERT INTO auth_identity_mappings(id,provider_id,external_user_id,smartcs_agent_id,external_username,external_email,last_login_at) '
            'VALUES(?,?,?,?,?,?,datetime(\'now\',\'localtime\'))',
            (mid, provider_id, result['user_id'], smartcs_agent_id,
             result.get('username', ''), result.get('email', ''))
        )
    
    get_db().commit()
    
    # 建立会话
    agent_info = get_db().execute('SELECT id,name,role FROM agents WHERE id=?', (smartcs_agent_id,)).fetchone()
    if not agent_info:
        return jsonify({'error': '登录失败：无法建立会话'}), 500
    
    session['agent_id'] = agent_info['id']
    session['agent_name'] = agent_info['name']
    session['agent_role'] = agent_info['role']
    get_db().execute('UPDATE agents SET status=\'online\' WHERE id=?', (agent_info['id'],))
    get_db().commit()
    
    redirect_url = '/admin/dashboard' if agent_info['role'] == 'admin' else '/agent/dashboard'
    return jsonify({'ok': True, 'redirect': redirect_url, 'role': agent_info['role']})


# ====== OIDC 登录流程 ======

@app.route('/api/auth/oidc/login/<provider_id>')
def auth_oidc_login(provider_id):
    provider = get_auth_provider(provider_id)
    if not provider or not isinstance(provider, OidcProvider):
        return jsonify({'error': 'OIDC 提供者不存在或类型错误'}), 404
    auth_url, state = provider.get_authorization_url()
    session['oidc_state'] = state
    session['oidc_provider_id'] = provider_id
    return redirect(auth_url)


@app.route('/agent/oidc/callback')
def auth_oidc_callback():
    """OIDC 回调处理"""
    error = request.args.get('error', '')
    if error:
        return jsonify({'error': f'OIDC 授权失败: {error}'})
    
    code = request.args.get('code', '')
    state = request.args.get('state', '')
    saved_state = session.pop('oidc_state', '')
    provider_id = session.pop('oidc_provider_id', '')
    
    if state != saved_state:
        return jsonify({'error': 'state 不匹配，可能存在 CSRF 攻击'}), 400
    if not code or not provider_id:
        return jsonify({'error': '缺少授权码或提供者 ID'}), 400
    
    provider = get_auth_provider(provider_id)
    if not provider:
        return jsonify({'error': '认证提供者不存在'}), 404
    
    # 用授权码换 token
    token_data = provider.exchange_code(code)
    access_token = token_data.get('access_token')
    if not access_token:
        return jsonify({'error': '换取 Token 失败', 'detail': token_data}), 400
    
    # 获取用户信息
    userinfo = provider.get_userinfo(access_token)
    if not userinfo or userinfo.get('sub'):
        pass
    
    external_id = userinfo.get('sub') or userinfo.get('id') or ''
    external_email = userinfo.get('email') or ''
    external_username = userinfo.get('preferred_username') or userinfo.get('name') or external_email.split('@')[0]
    display_name = userinfo.get('name') or external_username
    
    if not external_id:
        return jsonify({'error': '无法获取用户标识'}), 400
    
    # 查找或创建映射
    mapping = get_db().execute(
        'SELECT * FROM auth_identity_mappings WHERE provider_id=? AND external_user_id=?',
        (provider_id, external_id)
    ).fetchone()
    
    if mapping:
        smartcs_agent_id = mapping['smartcs_agent_id']
        get_db().execute('UPDATE auth_identity_mappings SET last_login_at=datetime(\'now\',\'localtime\') WHERE id=?',
                         (mapping['id'],))
    else:
        agent = None
        if external_email:
            agent = get_db().execute('SELECT id FROM agents WHERE email=?', (external_email,)).fetchone()
        if not agent:
            aid = 'auth-agent-' + uuid.uuid4().hex[:8]
            pwd_hash = hashlib.sha256(f'admin:{uuid.uuid4().hex[:12]}'.encode()).hexdigest()
            get_db().execute(
                'INSERT INTO agents(id,name,email,password_hash,role) VALUES(?,?,?,?,?)',
                (aid, display_name, external_email or f'{external_username}@oidc.local', pwd_hash, 'agent')
            )
            get_db().execute(
                'INSERT INTO agent_profiles(id,agent_id,display_name,department) VALUES(?,?,?,?)',
                ('ap-' + uuid.uuid4().hex[:8], aid, display_name, 'OIDC导入')
            )
            smartcs_agent_id = aid
        else:
            smartcs_agent_id = agent['id']
        
        mid = 'authmap-' + uuid.uuid4().hex[:12]
        get_db().execute(
            'INSERT INTO auth_identity_mappings(id,provider_id,external_user_id,smartcs_agent_id,external_username,external_email,last_login_at) '
            'VALUES(?,?,?,?,?,?,datetime(\'now\',\'localtime\'))',
            (mid, provider_id, external_id, smartcs_agent_id, external_username, external_email)
        )
    
    get_db().commit()
    
    # 建立会话
    agent_info = get_db().execute('SELECT id,name,role FROM agents WHERE id=?', (smartcs_agent_id,)).fetchone()
    if not agent_info:
        return jsonify({'error': '登录失败'}), 500
    
    session['agent_id'] = agent_info['id']
    session['agent_name'] = agent_info['name']
    session['agent_role'] = agent_info['role']
    get_db().execute('UPDATE agents SET status=\'online\' WHERE id=?', (agent_info['id'],))
    get_db().commit()
    
    redirect_url = '/admin/dashboard' if agent_info['role'] == 'admin' else '/agent/dashboard'
    return redirect(redirect_url)


@app.route('/api/admin/auth-mappings', methods=['GET'])
@secadmin_required
def admin_list_auth_mappings():
    provider_id = request.args.get('provider_id', '')
    where = ['1=1']
    params = []
    if provider_id:
        where.append('provider_id=?')
        params.append(provider_id)
    mappings = get_db().execute(
        f'SELECT m.*, a.name as agent_name FROM auth_identity_mappings m '
        f'LEFT JOIN agents a ON m.smartcs_agent_id = a.id '
        f'WHERE {" AND ".join(where)} ORDER BY m.created_at DESC', params
    ).fetchall()
    return jsonify([dict(m) for m in mappings])


# ====== End Phase 5 ======

# ====== P0: Auto-close / Auto-rate Timeout Background Task ======
AUTO_CLOSE_MINUTES = int(os.environ.get('AUTO_CLOSE_MINUTES', '20'))
AUTO_RATE_HOURS = int(os.environ.get('AUTO_RATE_HOURS', '24'))

def _auto_timeout_check(app_ref):
    """Background task: auto-close created tickets idle > AUTO_CLOSE_MINUTES,
    auto-rate resolved tickets idle > AUTO_RATE_HOURS"""
    with app_ref.app_context():
        try:
            db = get_db()
            now = datetime.now()
            # Auto-close: created tickets with no new messages in AUTO_CLOSE_MINUTES
            cutoff_close = (now - timedelta(minutes=AUTO_CLOSE_MINUTES)).strftime('%Y-%m-%d %H:%M:%S')
            created_idle = db.execute(
                "SELECT t.id, t.conversation_id FROM service_tickets t "
                "LEFT JOIN (SELECT conversation_id, MAX(created_at) as last_msg "
                "FROM messages GROUP BY conversation_id) m ON t.conversation_id=m.conversation_id "
                "WHERE t.status='created' AND (m.last_msg IS NULL OR m.last_msg<?)",
                (cutoff_close,)
            ).fetchall()
            for t in created_idle:
                transition_ticket(t['id'], 'closed', 'system', '系统',
                                {'close_reason': f'AI工单超时自动关闭（{AUTO_CLOSE_MINUTES}分钟无交互）'})
            # Auto-rate: resolved tickets idle > AUTO_RATE_HOURS
            cutoff_rate = (now - timedelta(hours=AUTO_RATE_HOURS)).strftime('%Y-%m-%d %H:%M:%S')
            resolved_idle = db.execute(
                "SELECT t.id, t.conversation_id FROM service_tickets t "
                "LEFT JOIN (SELECT conversation_id, MAX(created_at) as last_msg "
                "FROM messages GROUP BY conversation_id) m ON t.conversation_id=m.conversation_id "
                "WHERE t.status='resolved' AND (m.last_msg IS NULL OR m.last_msg<?)",
                (cutoff_rate,)
            ).fetchall()
            for t in resolved_idle:
                transition_ticket(t['id'], 'rated', 'system', '系统',
                                {'customer_rating': 5, 'customer_feedback': '超时自动评价',
                                 'close_reason': f'超时自动评价（{AUTO_RATE_HOURS}小时无交互）'})
        except Exception as e:
            print(f'Auto-timeout check error: {e}')

def start_timeout_checker(app_ref):
    """Start background thread for periodic timeout checks"""
    def checker():
        while True:
            cfg = get_system_config()
            interval = int(cfg.get('auto_check_interval', '300'))
            time.sleep(interval)
            _auto_timeout_check(app_ref)
    t = threading.Thread(target=checker, daemon=True)
    t.start()


if __name__ == '__main__':
    os.makedirs(os.path.join(BASE_DIR,'data'),exist_ok=True)
    os.makedirs(os.path.join(BASE_DIR,"uploads"),exist_ok=True)
    print(f"🚀 SmartCS v3.0 启动")
    print(f"   💬 工程师端: http://localhost:5000/")
    print(f"   👤 登录页: http://localhost:5000/agent/login")
    print(f"   📊 工作台: http://localhost:5000/agent/dashboard")
    print(f"   ⚙️  后台: http://localhost:5000/admin/dashboard")
    print(f"   📚 知识库: http://localhost:5000/upload")
    print(f"   ⏱️  超时: AI工单{AUTO_CLOSE_MINUTES}分自动关闭 / 已解决{AUTO_RATE_HOURS}小时自动评价")
    start_timeout_checker(app)
    app.run(host="0.0.0.0",port=5000,debug=False)
